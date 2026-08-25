"""
教学设计修复脚本v2 - 函数先定义，后调用
"""
import os, re, json
from docx import Document
from docx.shared import Pt
from docx.oxml.ns import qn
from pptx import Presentation
import store

fp = r"生成结果\精修版\2023-2024-2《商务数据分析》教学设计 杜媛.docx"
ppt_base = r"原始资料\教材\商务数据分析\大数据分析方法项目实战\03 课程PPT"
code_base = r"原始资料\教材\商务数据分析\大数据分析方法项目实战\04 实训源代码"

# ============================================================
# 辅助函数（全部在主逻辑之前定义）
# ============================================================

def write_cell_safe(cell, text):
    """安全写入单元格"""
    while len(cell.paragraphs) > 1:
        cell.paragraphs[-1]._element.getparent().remove(cell.paragraphs[-1]._element)
    p0 = cell.paragraphs[0]
    for r in p0.runs:
        r.text = ""
    lines = text.split("\n") if text else [""]
    for li, line in enumerate(lines):
        if li == 0:
            if p0.runs:
                p0.runs[0].text = line
            else:
                p0.add_run(line)
        else:
            np = cell.add_paragraph()
            np.add_run(line)

def get_import_code(ch):
    imports = {
        "初识数据分析": "import numpy as np\nimport pandas as pd",
        "Excel数据分析工具": "import pandas as pd",
        "Numpy数学运算库": "import numpy as np",
        "Pandas数据分析库": "import pandas as pd\nimport numpy as np",
        "SciPy科学计算库": "from scipy import stats, linalg\nimport numpy as np",
        "Sklearn数据统计基础": "from sklearn import datasets, preprocessing\nfrom sklearn.model_selection import train_test_split",
        "Sklearn数据统计进阶": "from sklearn.linear_model import LinearRegression\nfrom sklearn.cluster import KMeans",
        "Seaborn可视化分析库": "import seaborn as sns\nimport matplotlib.pyplot as plt",
        "综合评价与课程总结": "import numpy as np\nimport pandas as pd\nimport seaborn as sns",
    }
    return imports.get(ch, "import numpy as np")

def gen_intro(ch, ti_in_ch, all_kps, case, title):
    if ti_in_ch == 0:
        return (f"【案例展示】{case}\n\n"
                f"【提出问题】\n  · 这个案例中用到了什么数据分析方法？\n  · 需要掌握哪些知识？\n\n"
                f"【学生讨论】分组讨论3分钟，每组提出1-2个观点\n\n"
                f"【引出课题】今天学习{title}，这是{ch}模块的第一节课。")
    elif ti_in_ch == 1:
        return (f"【复习导入】上节课学习了{all_kps[0]['标题'] if all_kps else title}。\n"
                f"  回顾要点：{all_kps[0].get('概念','')[:60] if all_kps else ''}...\n\n"
                f"【案例展示】{case}\n\n"
                f"【引出课题】本节课学习{title}，在上节课基础上深入。")
    elif ti_in_ch == 2:
        return (f"【复习导入】回顾前两节内容：\n"
                f"  · {all_kps[0]['标题'] if len(all_kps)>0 else ''}\n"
                f"  · {all_kps[1]['标题'] if len(all_kps)>1 else ''}\n\n"
                f"【案例展示】{case}\n\n"
                f"【引出课题】本节课学习{title}，是本章核心内容。")
    else:
        return (f"【复习导入】回顾本章所有知识点：\n"
                f"  · {all_kps[0]['标题'] if len(all_kps)>0 else ''}\n"
                f"  · {all_kps[1]['标题'] if len(all_kps)>1 else ''}\n"
                f"  · {all_kps[2]['标题'] if len(all_kps)>2 else ''}\n\n"
                f"【任务引入】本章知识点已学完，进入任务实施环节。\n"
                f"  需要综合运用所学知识完成{title}。\n\n"
                f"【明确要求】提交代码+报告+截图，下节课前提交。")

def gen_goals(kp, title, task):
    kg = task.get("knowledge_goal","") or f"理解{kp.get('标题',title)}的基本概念"
    ag = task.get("ability_goal","") or f"能够运用{kp.get('标题',title)}完成对应功能"
    return (f"教师归纳讨论结果，明确学习目标：\n"
            f"1. 知识目标：{kg}\n2. 能力目标：{ag}\n"
            f"3. 素质目标：培养规范操作意识\n4. 思政目标：培养数据驱动思维")

def gen_knowledge_content(kp, ppt_slides, title):
    content = f"任务1：{kp.get('标题',title)}（知识讲解）\n\n一、基本概念\n"
    concept = kp.get("概念","")
    if concept:
        content += f"  {concept}\n\n"
    content += "二、详细知识点\n"
    for dk in [k for k in kp if k not in ("标题","概念","要点")]:
        content += f"  【{dk}】\n    {kp[dk]}\n\n"
    # PPT补充
    ppt_pts = []
    for sl in ppt_slides[:12]:
        for t in sl["texts"]:
            if 5 < len(t) < 200 and t not in ppt_pts:
                ppt_pts.append(t)
            if len(ppt_pts) >= 8: break
        if len(ppt_pts) >= 8: break
    if ppt_pts:
        content += "三、PPT补充要点\n"
        for pt in ppt_pts[3:8]:
            content += f"  · {pt[:100]}\n"
    content += f"\n四、本节要点\n  {kp.get('要点','掌握基本概念和操作方法。')}"
    return content

def gen_demo_content(kp, title, code_files, ch):
    content = f"任务2：{kp.get('标题',title)}（技术演示）\n\n"
    content += "一、环境准备\n  1. 启动Jupyter Notebook\n  2. 新建Python 3笔记本\n  3. 导入库\n\n"
    content += f"二、导入库\n```python\n{get_import_code(ch)}\n```\n\n"
    content += "三、核心代码演示\n"
    if code_files:
        for ci,(cname,ccode) in enumerate(code_files[:2]):
            content += f"  示例{ci+1}：{cname}\n"
            lines = ccode.split("\n")[:20]
            content += "```python\n" + "\n".join(lines) + "\n```\n\n"
    else:
        content += "```python\n# 演示代码\nprint('演示" + title + "')\n```\n"
    content += f"四、运行结果分析\n  · 观察{kp.get('标题',title)}的输出\n  · 对比不同参数结果\n  · 记录常见错误"
    return content

def gen_practice(kp, title):
    return (f"任务3：{kp.get('标题',title)}（实操练习）\n\n"
            f"一、实训目标\n  1. 独立完成{title}代码实现\n  2. 调试并解决错误\n  3. 记录结果并分析\n\n"
            f"二、实训步骤\n  1. 新建Notebook\n  2. 导入库\n  3. 按演示逐步实现\n  4. 修改参数观察变化\n  5. 扩展功能（选做）\n\n"
            f"三、要求\n  · 代码规范、注释完整\n  · 截图保存结果\n  · 遇错先自查（报错→文档→同学→老师）\n\n"
            f"四、成果展示\n  完成3组展示，其他同学互评，教师点评。")

def gen_practice_task(title, ch):
    return (f"任务实施：{title}\n\n"
            f"一、任务描述\n  综合运用{ch}章节知识完成实训项目。\n\n"
            f"二、任务要求\n  1. 数据准备：使用课堂数据集或自选\n  2. 代码实现：独立编写完整代码\n"
            f"  3. 结果分析：对输出结果分析\n  4. 报告撰写：过程+截图+结论\n\n"
            f"三、评分标准\n  · 代码正确性40%  · 规范性20%\n  · 报告质量25%  · 创新拓展15%\n\n"
            f"四、提交\n  源码+报告+截图，下节课前提交学习通。")

def gen_check(title):
    return (f"检查与优化：{title}\n\n"
            f"一、代码检查\n  1. 能否正确运行\n  2. 输出结果是否合理\n  3. 代码是否规范\n\n"
            f"二、常见问题\n  · 数据类型错误：检查astype()\n  · 索引越界：检查shape\n"
            f"  · 空值：检查isnull()/fillna()\n  · 性能：向量化替代循环\n\n"
            f"三、优化建议\n  · 添加异常处理\n  · 增加日志输出\n  · 封装函数\n  · 添加图表展示")

def gen_optimize(title):
    return (f"成果制作与优化：{title}\n\n"
            f"一、成果整合\n  1. 整理代码确保可运行\n  2. 补充报告内容\n  3. 确保截图清晰\n\n"
            f"二、成果展示\n  · 每人/组展示3分钟\n  · 说明：目标→方法→代码→结果\n  · 其他同学提问\n\n"
            f"三、互评与自评\n  · 互评：给出评价和建议\n  · 自评：总结收获和不足\n"
            f"  · 教师点评：共性问题、优秀做法\n\n"
            f"四、课后提交\n  最终版提交学习通，截止本周日24:00")

def gen_scene(ch, idx, room):
    opts = [
        f"{room}，配备多媒体教学设备、Python开发环境（Anaconda+Jupyter Notebook）、投影仪、白板。学生每人一台计算机。",
        f"{room}，配备多媒体教学设备、Python开发环境。本节课以实操为主，学生需同步操作。",
        f"{room}，配备多媒体教学设备、Python开发环境。课中安排分组讨论环节，每组3-4人。",
        f"{room}，配备多媒体教学设备、Python开发环境。本节课为任务实施课，教师巡回指导。",
    ]
    return opts[idx % len(opts)]

def gen_resources(ch, idx):
    opts = [
        f"1. 多媒体课件：{ch}PPT课件（约20页）\n2. 实训资源：配套数据集（CSV）和源代码\n3. 在线平台：学习通微课和练习题\n4. 开发工具：Anaconda+Jupyter\n5. 参考资料：《大数据分析方法项目实战》教材",
        f"1. 多媒体课件：{ch}PPT课件\n2. 实训资源：新增函数和代码示例\n3. 开发工具：Anaconda+Jupyter\n4. 学习通课前测验（5题）\n5. 参考资料：Python官方文档",
        f"1. 多媒体课件：{ch}PPT课件含综合案例\n2. 实训资源：综合实训数据集和项目模板\n3. 开发工具：Anaconda+Jupyter\n4. 分组工具：随机分组\n5. 评价工具：实训评分表",
        f"1. 多媒体课件：{ch}任务实施指导PPT\n2. 实训资源：完整项目源码参考和数据集\n3. 开发工具：Anaconda+Jupyter\n4. 提交平台：学习通作业入口\n5. 评价标准：评分表和自评表",
    ]
    return opts[idx % len(opts)]

def gen_method(idx):
    k = ["讲授法、案例教学法、启发式提问；\n学生听讲、思考、做笔记、回答问题",
         "讲授法+演示法+对比教学法；\n学生听讲、对比分析、提问讨论",
         "任务驱动法+探究式学习；\n学生自主探究、小组讨论、成果分享",
         "项目教学法+翻转课堂；\n学生课前预习、课中实操、课后巩固"]
    d = ["操作演示法、逐步讲解法、巡回指导法；\n学生跟随操作、提问讨论、记录笔记",
         "演示法+对比实验法；\n学生同步操作、对比参数、记录差异",
         "任务驱动法+巡回指导法；\n学生独立操作、小组互助、成果展示",
         "项目教学法+个性化指导；\n学生自主实践、教师针对性指导"]
    p = ["实操练习法、任务驱动法、巡回指导法、成果展示法；\n学生独立操作、小组互助、展示分享",
         "实操练习法+同伴互查法；\n学生独立完成、互相检查、讨论改进",
         "项目实践法+成果汇报法；\n学生完成项目、汇报展示、互评学习",
         "综合实践法+反思总结法；\n学生完成项目、总结反思、提交成果"]
    return k[idx%4], d[idx%4], p[idx%4]

def gen_goal_text(kp, title):
    return (f"理解{kp.get('标题',title)}的基本概念、原理和分类；建立知识框架",
            f"掌握{kp.get('标题',title)}的操作步骤和代码实现；培养调试能力",
            f"能独立完成{kp.get('标题',title)}实训；培养问题解决和协作能力")


# ============================================================
# 知识点数据（每章节的核心知识点）
# ============================================================

chapter_knowledge = {
    "初识数据分析": {
        "case": "某农产品电商平台通过用户购买行为数据分析，发现70%用户浏览3个以上商品后才下单，平台据此优化推荐策略，转化率提升25%。这背后就是数据分析的基本方法。",
        "kps": [
            {"标题":"认识数据分析","概念":"数据分析是用适当的统计分析方法对收集来的大量数据进行分析，提取有用信息和形成结论的过程。数据分析与数据挖掘密切相关，数据分析的目的在于把隐没在大量数据中的信息提炼出来，为决策提供依据。",
             "分类":"数据分析分为描述性分析（发生了什么）、诊断性分析（为什么发生）、预测性分析（将会发生什么）、规范性分析（应该做什么）四类。",
             "流程":"数据分析流程：明确分析目标→数据采集→数据清洗→数据分析→数据可视化→撰写报告。",
             "要点":"数据分析是数学与计算机科学结合的产物，为企业决策提供数据支撑。"},
            {"标题":"常用数据分析方法","概念":"常用数据分析方法包括对比分析法、分组分析法、结构分析法、平均分析法、交叉分析法等。",
             "对比分析":"对比分析法是将两个或多个数据进行比较，分析差异。包括同比（与去年同期比）、环比（与上月比）、定基比（与固定期比）。",
             "分组分析":"分组分析法将总体数据按某一标志分成若干组，分析各组特征。如按年龄段分组分析购买力差异。",
             "要点":"实际应用中往往多种方法结合使用，才能全面分析数据。"},
            {"标题":"数据分析指标","概念":"常见指标包括转化率、留存率、活跃用户数（DAU/MAU）、用户获取成本（CAC）、客户终身价值（LTV）等。",
             "转化率":"转化率=完成目标行为的用户数/总访问用户数×100%。如1000访客中50人下单，转化率5%。",
             "留存率":"留存率=某日后仍活跃用户数/当日新增用户数×100%。次日留存率和7日留存率是核心指标。",
             "要点":"不同业务场景关注不同指标，选择合适指标是关键。"},
            {"标题":"常用数据分析工具","概念":"常用工具：Excel（基础处理）、Python（NumPy/Pandas/sklearn/Seaborn）、R语言、SQL、Tableau、Power BI等。",
             "Excel":"Excel适合小规模数据处理，提供数据透视表、图表、函数等功能，操作简单但不适合大数据。",
             "Python":"Python是数据分析主流语言。NumPy数值计算、Pandas数据处理、Matplotlib/Seaborn可视化、sklearn机器学习。",
             "要点":"工具选择取决于数据规模和分析深度，Python+SQL是电商数据分析主流组合。"}
        ]
    },
    "Excel数据分析工具": {
        "case": "某农产品电商运营团队需分析月度销售数据，包含5000条订单记录。用Excel数据透视表5分钟完成按品类、地区、月份的多维汇总，发现某些品类的季节性销售规律。",
        "kps": [
            {"标题":"Excel概述","概念":"Excel是Microsoft Office套件中的电子表格软件，广泛用于数据记录、计算、分析和可视化。",
             "特点":"优势：操作简单、界面直观、适合小规模分析。局限：处理大数据（>10万行）性能差，不支持复杂统计。",
             "要点":"Excel常用于快速数据汇总、趋势图制作和简单报表。"},
            {"标题":"Excel数据分析技巧","概念":"核心技巧：数据透视表（PivotTable）、VLOOKUP函数、条件格式、图表制作等。",
             "透视表":"数据透视表是Excel最强大分析工具，可快速汇总、分类、交叉分析。创建：选中数据→插入→数据透视表→拖拽字段。",
             "VLOOKUP":"VLOOKUP查找指定值并返回对应数据。语法：=VLOOKUP(查找值,数据范围,列序号,匹配方式)。",
             "要点":"数据透视表和VLOOKUP是Excel数据分析基础，解决80%日常需求。"},
            {"标题":"函数相关概念","概念":"Excel函数分数学函数(SUM/AVERAGE)、统计函数(STDEV/MEDIAN)、逻辑函数(IF/AND)、文本函数(LEFT/RIGHT)、日期函数(YEAR/MONTH)等。",
             "SUMIF":"SUMIF按条件求和：=SUMIF(条件范围,条件,求和范围)。如=SUMIF(A:A,\">100\",B:B)。",
             "COUNTIF":"COUNTIF按条件计数：=COUNTIF(范围,条件)。如=COUNTIF(B:B,\">90\")统计大于90的个数。",
             "要点":"函数组合可实现复杂分析逻辑，如SUMIFS多条件求和。"},
            {"标题":"Excel图表可视化","概念":"Excel支持柱状图、折线图、饼图、散点图等。选择图表取决于数据特征和分析目的。",
             "柱状图":"柱状图适合比较不同类别数值大小。如各品类销售额对比。",
             "折线图":"折线图适合展示数据随时间的变化趋势。如月度销售额趋势。",
             "要点":"图表设计原则：标题清晰、标注完整、颜色协调。避免3D图表。"}
        ]
    },
    "Numpy数学运算库": {
        "case": "某农产品电商平台需分析10万条用户行为数据，计算用户活跃度的平均值、标准差等统计指标。用Excel会很慢，而NumPy几行代码毫秒级完成。",
        "kps": [
            {"标题":"位运算函数","概念":"NumPy位运算：bitwise_and(位与)、bitwise_or(位或)、bitwise_xor(异或)、invert(取反)、left_shift(左移)、right_shift(右移)。",
             "bitwise_and":"位与运算：二进制位都为1则为1。np.bitwise_and([1,2,3],[2,3,4])返回[0,2,0]。",
             "left_shift":"左移运算：二进制位左移，相当于乘2的n次方。np.left_shift(5,1)返回10。",
             "要点":"位运算在图像处理、加密算法、权限管理中有广泛应用。"},
            {"标题":"统计函数","概念":"NumPy统计函数：sum、mean、median、std、var、max、min、ptp、average等。",
             "mean_std":"np.mean()计算均值，np.std()计算标准差。如ages=np.array([18,20,22])，np.mean=20.0，np.std=1.63。",
             "cumsum":"np.cumsum()计算累计和，np.cumprod()计算累计积。用于计算累计销售额等。",
             "要点":"统计函数支持按轴计算：axis=0按列，axis=1按行。"},
            {"标题":"数学函数","概念":"NumPy数学函数：三角函数(sin/cos/tan)、舍入函数(around/ceil/floor)、幂和对数(power/sqrt/exp/log)等。",
             "三角函数":"np.sin(np.pi/2)=1.0。注意使用弧度制，np.deg2rad()转换角度。",
             "舍入函数":"np.around(3.14159,2)=3.14，np.ceil(3.1)=4.0，np.floor(3.9)=3.0。",
             "要点":"NumPy数学函数都是元素级的，比Python循环快100倍以上。"},
            {"标题":"取反与任务实施","概念":"np.invert()对整数按位取反。任务实施要求学生用NumPy分析学生成绩数据。",
             "invert":"np.invert(np.array([1,2,3]))返回[-2,-3,-4]。补码表示下取反等于-(x+1)。",
             "任务实施":"学生信息统计：加载成绩数据，计算总分、平均分、标准差、最高/最低分。",
             "要点":"任务实施要求独立完成数据加载、统计分析、结果输出的完整流程。"}
        ]
    },
    "Pandas数据分析库": {
        "case": "某航空公司拥有10万条航班数据，需分析不同航线的准点率、客流量和收益。用Pandas的pivot_table和groupby快速完成多维分析。",
        "kps": [
            {"标题":"统计函数与transform","概念":"Pandas统计函数：sum/mean/std/count。transform()将聚合结果广播回原始DataFrame形状。",
             "transform":"df.groupby('品类')['销售额'].transform('mean')返回每品类平均销售额，与原DataFrame等长。",
             "要点":"transform返回等长结果，apply可返回任意形状。"},
            {"标题":"apply函数与透视表","概念":"apply()对行或列应用自定义函数。pivot_table()是数据透视表功能。",
             "apply":"df.apply(np.sum, axis=0)对每列求和。df['列'].apply(lambda x: x*2)对列每个元素乘2。",
             "pivot_table":"pd.pivot_table(df, values='销售额', index='地区', columns='月份', aggfunc='sum')。",
             "要点":"pivot_table核心参数：values、index、columns、aggfunc、fill_value。"},
            {"标题":"agg函数与交叉表","概念":"agg()支持对不同列应用不同聚合函数。crosstab()计算两个分类变量的频数表。",
             "agg":"df.agg({'销售额':'sum', '利润':'mean', '数量':'count'})一次完成多种统计。",
             "crosstab":"pd.crosstab(df['性别'], df['品类'])创建性别与品类的交叉频数表。",
             "要点":"agg支持多函数：df.agg(['sum','mean','std'])。"},
            {"标题":"标准化数据","概念":"数据标准化将不同量级数据转换到同一标准。常用：小数定标、Z-score、Min-Max标准化。",
             "小数定标":"x' = x / 10^j，j是最大绝对值的整数位数。如最大值987，则j=3，x'=x/1000。",
             "标准差标准化":"Z-score：x' = (x - mean) / std，转换后均值0标准差1。",
             "要点":"标准化是机器学习数据预处理的关键步骤。"}
        ]
    },
    "SciPy科学计算库": {
        "case": "某投资公司需分析股票收益率风险特征，计算相关性矩阵、正态性检验、拟合概率分布。这些高级统计分析用SciPy可轻松完成。",
        "kps": [
            {"标题":"SciPy简介与模块","概念":"SciPy是依赖NumPy的科学计算工具包，含cluster、fft、integrate、interpolate、linalg、ndimage、stats等模块。",
             "安装":"pip install scipy。导入：from scipy import stats, linalg, fftpack。",
             "要点":"SciPy在NumPy基础上构建高级科学计算功能。"},
            {"标题":"cluster模块与统计函数","概念":"cluster模块提供K-Means聚类。stats模块提供概率分布、统计检验等。",
             "kmeans":"from scipy.cluster.vq import kmeans, vq。centroids = kmeans(data, 3)聚为3类。",
             "stats":"from scipy import stats。stats.norm.rvs(size=10)生成正态分布随机数。stats.ttest_ind(g1,g2)做t检验。",
             "要点":"stats模块支持100多种概率分布和各种假设检验。"},
            {"标题":"线性代数与傅里叶变换","概念":"linalg提供矩阵运算：det、eig、svd、inv。fftpack提供快速傅里叶变换。",
             "det_eig":"linalg.det(A)计算行列式。linalg.eig(A)计算特征值和特征向量。",
             "fft":"fftpack.fft(signal)将时域信号转换为频域，用于信号分析和滤波。",
             "要点":"SciPy线性代数比NumPy更全面，适合复杂数学和工程问题。"}
        ]
    },
    "Sklearn数据统计基础": {
        "case": "泰坦尼克号乘客数据集包含891名乘客的姓名、性别、年龄、票价、是否生还等信息。用sklearn可以预测哪些乘客更可能生还，是经典机器学习入门案例。",
        "kps": [
            {"标题":"sklearn简介与数据集","概念":"sklearn是基于NumPy/SciPy的机器学习工具包，提供分类、回归、聚类、降维、预处理等功能。",
             "数据集":"自带数据集：load_iris(鸢尾花)、load_boston(波士顿房价)、load_digits(手写数字)。也可用make_blobs生成模拟数据。",
             "要点":"API统一：fit()训练、predict()预测、transform()转换、score()评估。"},
            {"标题":"特征提取与数据处理","概念":"特征提取：DictVectorizer、CountVectorizer、TfidfVectorizer、OneHotEncoder、LabelEncoder等。",
             "OneHotEncoder":"将分类变量转为0/1向量。如['男','女','男']变为[[1,0],[0,1],[1,0]]。",
             "StandardScaler":"将数据标准化为均值0标准差1。fit_transform()先拟合再转换。",
             "要点":"好的特征工程比算法选择更重要。"},
            {"标题":"数据集分割","概念":"train_test_split将数据分割为训练集和测试集，一般7:3或8:2。",
             "split":"from sklearn.model_selection import train_test_split。X_train,X_test,y_train,y_test = train_test_split(X,y,test_size=0.3,random_state=42)。",
             "要点":"random_state保证可复现。stratify=y保证类别比例一致。"}
        ]
    },
    "Sklearn数据统计进阶": {
        "case": "某银行需根据客户历史交易数据和信用记录，预测客户是否违约。用sklearn分类模型可构建违约预测模型，准确率可达85%以上。",
        "kps": [
            {"标题":"分类模型","概念":"分类模型：LogisticRegression(逻辑回归)、DecisionTreeClassifier(决策树)、GaussianNB(朴素贝叶斯)、SVC(支持向量机)。",
             "LogisticRegression":"逻辑回归用于二分类。model=LogisticRegression()，model.fit(X_train,y_train)，model.predict(X_test)。",
             "DecisionTree":"决策树自动学习分类规则。优点可解释性强，缺点容易过拟合。",
             "要点":"线性可分用逻辑回归，非线性用决策树，文本分类用朴素贝叶斯。"},
            {"标题":"回归模型与聚类","概念":"回归：LinearRegression、Ridge、PolynomialFeatures。聚类：KMeans、DBSCAN。",
             "LinearRegression":"线性回归预测连续值。model.coef_获取系数，model.intercept_获取截距。",
             "KMeans":"K-Means聚类分K组。model=KMeans(n_clusters=3)，model.fit(X)，model.labels_获取标签。",
             "要点":"回归评估用R²和MSE，聚类评估用轮廓系数。"},
            {"标题":"模型评估与保存","概念":"评估指标：accuracy、precision、recall、f1-score、ROC-AUC。交叉验证cross_val_score。",
             "metrics":"from sklearn.metrics import classification_report。输出精确率、召回率、F1值。",
             "save_model":"import joblib。joblib.dump(model,'model.pkl')保存，joblib.load加载。",
             "要点":"准确率不是唯一指标，不平衡数据要关注精确率和召回率。"}
        ]
    },
    "Seaborn可视化分析库": {
        "case": "某房产公司需分析波士顿房价数据集，探索房价与房间数、犯罪率等特征的关系。用Seaborn可视化可快速绘制清晰美观的统计图表。",
        "kps": [
            {"标题":"seaborn概述与分类图","概念":"seaborn是基于matplotlib的高级可视化库，提供更美观样式和更简洁API。支持分类图、关系图、分布图、回归图。",
             "分类图":"sns.barplot()条形图、sns.countplot()计数图、sns.boxplot()箱线图、sns.violinplot()小提琴图。",
             "样式":"sns.set_style('whitegrid')设置白色网格背景。比matplotlib默认更美观。",
             "要点":"seaborn是matplotlib的补充而非替代，两者配合效果最佳。"},
            {"标题":"关系图与分布图","概念":"关系图：scatterplot、lineplot、relplot。分布图：histplot、kdeplot、jointplot。",
             "scatterplot":"sns.scatterplot(x='面积',y='价格',hue='区域',data=df)绘制散点图，用颜色区分区域。",
             "distplot":"sns.histplot(df['价格'], kde=True)绘制分布直方图并叠加核密度曲线。",
             "要点":"看关系用散点图，看分布用直方图，看比较用条形图，看趋势用折线图。"},
            {"标题":"回归图与样式管理","概念":"回归图：regplot、lmplot。样式管理：set_style、set_context、set_palette。",
             "regplot":"sns.regplot(x='面积',y='价格',data=df)绘制散点图并拟合回归线。",
             "样式管理":"sns.set_style()设5种主题：white/dark/whitegrid/darkgrid/ticks。",
             "要点":"好的可视化不仅要数据准确，还要美观易读。"}
        ]
    },
    "综合评价与课程总结": {
        "case": "回顾本学期学习：从数据分析基础概念到Excel工具、Python库（NumPy、Pandas、SciPy、sklearn、Seaborn），同学们已具备完整数据分析能力。现在整合所有知识完成综合项目。",
        "kps": [
            {"标题":"课程成果汇报与评价","概念":"课程综合作品要求完成完整电商数据分析项目：数据采集、清洗、分析、可视化和报告。",
             "汇报要求":"展示项目背景、数据来源、分析方法、可视化图表、结论建议。时间8-10分钟。",
             "评价标准":"完整性30%、正确性20%、可视化20%、规范性15%、创新性15%。",
             "要点":"综合作品占总评60%，体现数据驱动的分析思维和规范操作流程。"},
            {"标题":"课程总结与复习","概念":"课程涵盖8个学习情境：数据分析基础、Excel、NumPy、Pandas、SciPy、sklearn基础、sklearn进阶、Seaborn。",
             "知识图谱":"流程：采集→清洗→分析→建模→可视化→报告。工具链：Excel→NumPy→Pandas→SciPy→sklearn→Seaborn。",
             "复习重点":"重点：透视表、分组聚合、分类回归模型、可视化图表。难点：特征工程、模型选择。",
             "要点":"课程目标达成自评：知识、能力、素质、思政四个维度。"}
        ]
    }
}

chapter_ppt_map = {
    "初识数据分析":"CORE-01","Excel数据分析工具":"CORE-02",
    "Numpy数学运算库":"CORE-03","Pandas数据分析库":"CORE-04",
    "SciPy科学计算库":"CORE-05","Sklearn数据统计基础":"CORE-06",
    "Sklearn数据统计进阶":"CORE-07","Seaborn可视化分析库":"CORE-08",
}
chapter_code_map = {
    "初识数据分析":"SOURCE-CORE-01","Numpy数学运算库":"SOURCE-CORE-03",
    "Pandas数据分析库":"SOURCE-CORE-04","SciPy科学计算库":"SOURCE-CORE-05",
    "Sklearn数据统计基础":"SOURCE-CORE-06","Sklearn数据统计进阶":"SOURCE-CORE-07",
    "Seaborn可视化分析库":"SOURCE-CORE-08",
}

# ============================================================
# 主执行
# ============================================================

# 加载PPT
ppt_data = {}
if os.path.exists(ppt_base):
    for f in sorted(os.listdir(ppt_base)):
        if f.endswith(".pptx"):
            key = f.replace(".pptx","")
            try:
                prs = Presentation(os.path.join(ppt_base, f))
                slides = []
                for si, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t and len(t) < 300:
                                    texts.append(t)
                    if texts:
                        slides.append({"index":si+1,"texts":texts})
                ppt_data[key] = slides
            except: pass

# 加载源码
code_data = {}
if os.path.exists(code_base):
    for root, dirs, files in os.walk(code_base):
        for f in files:
            if f.endswith(".py"):
                rel = os.path.relpath(os.path.join(root, f), code_base)
                try:
                    with open(os.path.join(root, f), "r", encoding="utf-8", errors="ignore") as fh:
                        code_data[rel] = fh.read(1500)
                except: pass

# 加载任务
tasks = store.rows("SELECT * FROM tasks WHERE offering_id=20 ORDER BY seq")

# 按章节分组
chapter_tasks = {}
chapter_order = []
for t in tasks:
    ch = t["chapter"]
    if ch not in chapter_tasks:
        chapter_tasks[ch] = []
        chapter_order.append(ch)
    chapter_tasks[ch].append(t)

# 打开文档
doc = Document(fp)
print(f"表格数: {len(doc.tables)}")

# ============================================================
# 修复1：基本信息表模板残留
# ============================================================
fix_count = 0
for ti in range(5, len(doc.tables), 2):
    t = doc.tables[ti]
    if len(t.rows) > 2:
        for ci in range(len(t.rows[2].cells)):
            ct = t.rows[2].cells[ci].text.strip()
            if "××" in ct:
                for p in t.rows[2].cells[ci].paragraphs:
                    for r in p.runs:
                        r.text = ""
                fix_count += 1
print(f"修复1 - 清除模板残留: {fix_count}处")

# ============================================================
# 修复2：重写教学过程内容
# ============================================================
ch_idx = 0
for ch in chapter_order:
    ch_tasks = chapter_tasks[ch]
    knowledge = chapter_knowledge.get(ch, {"case":"", "kps":[{"标题":ch,"概念":"","要点":""}]})
    kps = knowledge.get("kps", [])
    case = knowledge.get("case", "")
    ch_idx += 1

    for ti_in_ch, task in enumerate(ch_tasks):
        task_seq_idx = task["seq"] - 1  # 0-based
        org_idx = 6 + task_seq_idx * 2
        if org_idx >= len(doc.tables):
            break
        t = doc.tables[org_idx]

        kp_idx = min(ti_in_ch, len(kps)-1) if kps else 0
        kp = kps[kp_idx] if kps else {"标题":ch,"概念":"","要点":""}
        raw_title = task["title"].split("：",1)[1] if "：" in task["title"] else task["title"]
        # 只有标题以"任务实施"或"成果制作"开头才算任务实施课
        # "数据分析方法理论、任务实施"是知识课不是实施课
        is_practice = raw_title.startswith("任务实施") or raw_title.startswith("成果制作")
        title = raw_title.replace("任务实施：", "", 1) if raw_title.startswith("任务实施：") else raw_title

        ppt_key = chapter_ppt_map.get(ch)
        ppt_slides = ppt_data.get(ppt_key, []) if ppt_key else []
        code_dir = chapter_code_map.get(ch)
        code_files = [(k,v) for k,v in code_data.items() if code_dir and k.startswith(code_dir)] if code_dir else []

        room = "801教室" if task_seq_idx % 2 == 0 else "802教室"
        m_k, m_d, m_p = gen_method(task_seq_idx)
        g_k, g_d, g_p = gen_goal_text(kp, title)

        # R0: 教学场景
        if len(t.rows) > 0 and len(t.rows[0].cells) > 2:
            write_cell_safe(t.rows[0].cells[2], gen_scene(ch, task_seq_idx, room))
        # R1: 教学资源
        if len(t.rows) > 1 and len(t.rows[1].cells) > 2:
            write_cell_safe(t.rows[1].cells[2], gen_resources(ch, task_seq_idx))
        # R3: 导入
        if len(t.rows) > 3 and len(t.rows[3].cells) > 2:
            write_cell_safe(t.rows[3].cells[2], gen_intro(ch, ti_in_ch, kps, case, title))
        # R4: 目标
        if len(t.rows) > 4 and len(t.rows[4].cells) > 2:
            write_cell_safe(t.rows[4].cells[2], gen_goals(kp, title, task))
        # R5: 知识讲解
        if len(t.rows) > 5 and len(t.rows[5].cells) > 2:
            if is_practice:
                write_cell_safe(t.rows[5].cells[2], gen_practice_task(title, ch))
            else:
                write_cell_safe(t.rows[5].cells[2], gen_knowledge_content(kp, ppt_slides, title))
            if len(t.rows[5].cells) > 3: write_cell_safe(t.rows[5].cells[3], m_k)
            if len(t.rows[5].cells) > 4: write_cell_safe(t.rows[5].cells[4], g_k)
        # R6: 技术演示
        if len(t.rows) > 6 and len(t.rows[6].cells) > 2:
            if is_practice:
                write_cell_safe(t.rows[6].cells[2], gen_check(title))
            else:
                write_cell_safe(t.rows[6].cells[2], gen_demo_content(kp, title, code_files, ch))
            if len(t.rows[6].cells) > 3: write_cell_safe(t.rows[6].cells[3], m_d)
            if len(t.rows[6].cells) > 4: write_cell_safe(t.rows[6].cells[4], g_d)
        # R7: 实操练习
        if len(t.rows) > 7 and len(t.rows[7].cells) > 2:
            if is_practice:
                write_cell_safe(t.rows[7].cells[2], gen_optimize(title))
            else:
                write_cell_safe(t.rows[7].cells[2], gen_practice(kp, title))
            if len(t.rows[7].cells) > 3: write_cell_safe(t.rows[7].cells[3], m_p)
            if len(t.rows[7].cells) > 4: write_cell_safe(t.rows[7].cells[4], g_p)

        label = "[任务实施]" if is_practice else ""
        print(f"  任务{task_seq_idx+1} ({ch}): {label}")

doc.save(fp)
print(f"\n保存完成: {fp}")
