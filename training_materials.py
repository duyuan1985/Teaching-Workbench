import re
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path

import pdfplumber
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Pt
from openpyxl import load_workbook
from pptx import Presentation

import store


ROOT = Path(__file__).parent
TEMPLATE = ROOT / "原始资料" / "模板" / "模板8：《XXX》实训资料.docx"
LEGACY_EXTRACTOR = ROOT / "extract_legacy_training_doc.ps1"


def _docx_text(path):
    document = Document(path)
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.extend(cell.text.strip() for cell in row.cells if cell.text.strip())
    return "\n".join(parts)


def _legacy_doc_text(path):
    with tempfile.TemporaryDirectory(prefix="training-doc-") as temp_dir:
        output = Path(temp_dir) / "content.txt"
        result = subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass", "-File", str(LEGACY_EXTRACTOR),
             "-InputPath", str(path), "-OutputPath", str(output)],
            capture_output=True, text=True, encoding="utf-8",
        )
        if result.returncode != 0:
            return ""
        return output.read_text(encoding="utf-8-sig", errors="ignore") if output.exists() else ""


def _source_text(path):
    suffix = path.suffix.lower()
    try:
        if suffix == ".docx":
            return _docx_text(path)
        if suffix == ".doc":
            return _legacy_doc_text(path)
        if suffix == ".pdf":
            with pdfplumber.open(path) as pdf:
                return "\n".join(page.extract_text() or "" for page in pdf.pages)
        if suffix == ".pptx":
            presentation = Presentation(path)
            return "\n".join(shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text"))
        if suffix in {".xlsx", ".xlsm"}:
            workbook = load_workbook(path, read_only=True, data_only=True)
            return "\n".join("\t".join(str(value) for value in row if value is not None) for sheet in workbook for row in sheet.iter_rows(values_only=True))
        if suffix in {".txt", ".md", ".csv", ".html", ".htm", ".py", ".js", ".css", ".json", ".xml"}:
            return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""
    return ""


def collect_sources(source_dir):
    root = Path(source_dir)
    if not root.is_dir():
        raise ValueError("实训资料目录不存在。")
    files = [path for path in root.rglob("*") if path.is_file() and not path.name.startswith("~$")]
    supported = {".doc", ".docx", ".pdf", ".pptx", ".xlsx", ".xlsm", ".txt", ".md", ".csv", ".html", ".htm", ".py", ".js", ".css", ".json", ".xml"}
    selected = [path for path in files if path.suffix.lower() in supported]
    if not selected:
        raise ValueError("目录中没有可读取的实训资料。")
    text = "\n".join(_source_text(path) for path in selected)
    return selected, text


def _first_match(patterns, text, default=""):
    for pattern in patterns:
        match = re.search(pattern, text, re.I)
        if match:
            return match.group(1).strip()
    return default


def _set_paragraph(paragraph, text, font_name=None, size=None):
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
        run = paragraph.runs[0]
    else:
        run = paragraph.add_run(text)
    if font_name:
        run.font.name = font_name
        run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
    if size:
        run.font.size = Pt(size)


def _set_cell(cell, text):
    paragraph = cell.paragraphs[0]
    _set_paragraph(paragraph, str(text), "仿宋_GB2312", 10.5)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER


def _replace_xxx(document, course_name):
    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            if "XXX" in run.text:
                run.text = run.text.replace("XXX", course_name)
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if "XXX" in run.text:
                            run.text = run.text.replace("XXX", course_name)


def _prevent_row_split(row):
    properties = row._tr.get_or_add_trPr()
    if properties.find(qn("w:cantSplit")) is None:
        properties.append(OxmlElement("w:cantSplit"))


def _repeat_header(row):
    properties = row._tr.get_or_add_trPr()
    if properties.find(qn("w:tblHeader")) is None:
        properties.append(OxmlElement("w:tblHeader"))


def _course_units(offering_id):
    units = store.rows(
        "SELECT project_title,source_skills,suggested_hours FROM curriculum_units WHERE offering_id=? AND review_action<>'删除' ORDER BY seq",
        (offering_id,),
    )
    if units:
        return units
    return store.rows(
        """SELECT chapter AS project_title,
        CASE WHEN ability_goal<>'' THEN ability_goal ELSE title END AS source_skills,
        hours AS suggested_hours FROM tasks WHERE offering_id=? ORDER BY seq""",
        (offering_id,),
    )


def generate_training_materials(offering_id, source_dir, class_name=""):
    rows = store.rows("SELECT * FROM offerings WHERE id=?", (offering_id,))
    if not rows:
        raise ValueError("课程实例不存在。")
    offering = rows[0]
    files, source_text = collect_sources(source_dir)
    units = _course_units(offering_id)
    unit_names = [str(item["project_title"] or "").strip() for item in units if str(item["project_title"] or "").strip()]
    skill_names = []
    for item in units:
        skill_names.extend(part.strip() for part in str(item["source_skills"] or "").replace("；", ";").split(";") if part.strip())
    class_hint = _first_match((r"(?:农商|营销|全媒体)?(\d{6})",), " ".join(path.name for path in files))
    classes = [part.strip() for part in str(offering.get("teaching_class") or "").replace("；", ";").split(";") if part.strip()]
    if not class_name:
        class_name = next((item for item in classes if class_hint and class_hint in item), "、".join(classes))
    enrollment = store.rows("SELECT enrollment_count FROM offering_classes WHERE offering_id=? AND class_name=?", (offering_id, class_name))
    student_count = str(enrollment[0]["enrollment_count"]) if enrollment and enrollment[0]["enrollment_count"] else _first_match((r"(?:共计?|实训人数[:：]?)\s*(\d{1,3})\s*名?\s*学生", r"实训人数[:：]?\s*(\d{1,3})"), source_text, "按实际参训名单")
    location = _first_match((r"(?:实训地点|地点)[:：\s]*([^\n，。；]{2,20})", r"(\d{3,4}(?:机房|教室))"), source_text, "校内实训室")
    date_match = re.search(r"(20\d{2})\s*年\s*(\d{1,2})\s*月\s*(\d{1,2})\s*日\s*[-—至]+\s*(?:(20\d{2})\s*年\s*)?(\d{1,2})\s*月\s*(\d{1,2})\s*日", source_text)
    if date_match:
        year1, month1, day1, year2, month2, day2 = date_match.groups()
        year2 = year2 or year1
        date_text = f"{year1}年{int(month1)}月{int(day1)}日至{year2}年{int(month2)}月{int(day2)}日"
    else:
        date_text = "以学校实训教学安排为准"
    hours = str(offering["total_hours"]) if offering.get("offering_kind") == "实训课程" else _first_match((r"实训学时[:：\s]*(\d+)",), source_text, "30")
    teacher = store.get_setting("teacher_name", "杜媛")
    core = "、".join(unit_names[:6]) or "课程核心项目"
    skills = "、".join(dict.fromkeys(skill_names[:10])) or "数据采集、处理、分析、呈现与报告撰写"
    assessment_rows = store.rows(
        "SELECT component_name,weight FROM grade_components WHERE offering_id=? ORDER BY sort_order,id",
        (offering_id,),
    )
    assessment_text = "、".join(f"{item['component_name']}{float(item['weight']):g}分" for item in assessment_rows)
    if not assessment_text:
        assessment_text = "过程表现与成果质量综合评价"
    is_platform_training = offering.get("offering_kind") == "实训课程" and "商务数据分析" in offering["course_name"]

    document = Document(TEMPLATE)
    _replace_xxx(document, offering["course_name"])
    paragraphs = document.paragraphs
    body_font = "仿宋_GB2312"
    content = {
        25: (f"本实训使用中联商务大数据分析实训平台V1.0独立组织教学，不继承《商务数据分析》课程教材任务。学生通过平台视频课程学习、章节测验、实践练习和实训报告完成24学时训练，形成平台操作、商务数据分析实践、结果检查和规范表达能力；同时强化诚信学习、考勤纪律、数据安全和独立完成任务意识。" if is_platform_training else f"本实训依据《{offering['course_name']}》课程标准、课程教学安排和《职业学校学生实习管理规定》组织实施。围绕{core}开展集中训练，使学生掌握{skills}等知识与技能，能够按任务要求完成分析、检查和成果交付；同时强化诚实守信、劳动纪律、数据安全、设备安全、知识产权和团队协作意识。"),
        28: "实训按“安全与任务说明—教师示范—分组实践—过程检查—成果完善—展示评价”的流程实施，模块和日程见下表。",
        31: "采用过程评价与成果评价相结合的方式。过程评价重点考查出勤、操作规范、任务进度、协作表现、安全纪律和劳动表现；成果评价重点考查成果完整性、技术正确性、分析逻辑、表达规范及改进质量。对弄虚作假、违规使用数据或严重违反安全纪律的行为按规定处理。",
        33: "系部负责实训统筹与质量检查，教研室负责方案审查、资源协调和过程督导，指导教师负责安全教育、任务讲解、技术指导、考勤记录、过程评价和成果验收；学生组长协助完成资料收集、设备检查、进度反馈和现场整理。涉及企业指导时，由企业教师承担岗位规范、真实案例和成果评价指导。",
        36: "实训开始前开展设备用电、账号密码、数据保密、个人信息保护、网络资源使用和应急处置教育。学生须按操作规程使用软硬件，不得私接设备、传播敏感数据或使用来源不明的软件与素材；发现设备、网络或数据异常应立即停止操作并报告指导教师。",
        38: "原则上使用现有实训室、教学软件和数字资源，不另行安排经费；确需打印、耗材或校外实践支出的，按学校财务制度履行审批手续并据实核算。",
        40: "实训前准备方案、指导书、任务书、分组表、安全教育记录和评价标准；实训中留存考勤、过程记录、阶段成果和指导记录；实训后提交源文件、成果报告、评价表、实训总结及代表性成果，统一命名并分类归档。",
        42: "学生应按时到岗并完成签到，服从分组和任务安排；遵守课堂、操作、劳动和安全纪律，保持工位整洁，规范保存与提交文件。迟到、早退、请假和缺勤按学校规定记录；抄袭、伪造数据、擅自外传资料或违规操作的，视情节给予整改、扣分或取消相应任务成绩处理。",
        45: "附件包括实训任务书、安全教育及责任确认记录、学生分组与考勤表、过程检查表、成果评价表和优秀成果清单，可根据课程与班级实际情况增减。",
        55: f"通过《{offering['course_name']}》项目实训，使学生在连续任务中综合运用课程知识，完成从任务理解、方案设计、工具操作、结果检查到成果表达的完整工作过程，提升岗位实践、问题解决、团队协作和规范交付能力。",
        59: f"知识目标：理解{skills}的基本原理、操作规范和质量要求。\n能力目标：能够围绕{core}独立或协作完成任务实施、结果检查、问题修正和成果提交。\n素质目标：形成严谨求实、诚实守信、精益求精和持续改进意识，落实劳动教育、安全教育、数据保护与知识产权要求。",
        63: f"以{core}为主要载体，完成需求分析、任务分解、实践操作、过程记录、成果检查、展示汇报和复盘改进，形成可核验、可展示、可归档的实训成果。",
        65: "具体内容包括：①接受安全教育并明确任务、分组和评价要求；②根据任务书分析需求并制定实施计划；③运用课程工具和方法完成各模块实践；④检查数据、操作步骤和成果质量，记录并解决问题；⑤完成成果报告、展示汇报、互评与修改；⑥整理源文件、过程证据和最终成果并规范归档。",
        67: f"综合成绩按百分制评定：{assessment_text}，合计100分。平台视频学习、章节测验和实践练习成绩以中联商务大数据分析实训平台V1.0记录为依据；考勤由指导教师记录；实训报告根据完整性、分析过程、结果准确性和表达规范评分。",
        70: f"由{teacher}担任校内指导教师，负责教学组织、技术指导和评价。实训地点为{location}，使用课程配套教材、课件、任务书、案例素材、实训软件及计算机设备。实训前检查设备、网络、账号和安全设施，确保资料来源合法、环境稳定、应急通道畅通。",
        73: "实行指导教师安全负责制和学生岗位责任制。首次实训须完成安全培训和责任确认；操作中遵守设备用电、网络账号、数据访问、文件备份及场地管理规范；发生断电、设备故障、网络安全、数据泄露或人员不适等情况时，立即停止操作、保护现场并报告，按学校应急预案处置。",
        91: f"本次实训围绕{core}组织实施，完成安全教育、任务讲解、分组实践、过程指导、成果验收和总结评价。学生总体能够遵守考勤、劳动和安全要求，按计划完成主要任务并提交成果；具体人数、完成率和优秀成果情况以实际考勤及评价记录为准。",
        94: f"知识方面，学生进一步理解了{skills}的原理与规范；能力方面，能够完成任务分析、工具操作、结果检查、问题修正和成果表达；素质方面，规范意识、数据安全意识、劳动意识、团队协作和责任意识得到强化。个别学生仍需在综合应用、细节检查和成果说明方面继续提高。",
        97: "实训组织总体有序，但不同学生的知识基础、操作速度和自主排错能力存在差异，集中指导与个别辅导的时间分配仍需优化；部分过程记录和成果命名不够统一，设备与网络的课前检查、数据备份以及劳动教育的过程性记录还需进一步细化。",
        100: "后续将提前发布预习资料和工具检查清单，增加分层任务、关键步骤示范和易错点微课；完善过程检查、阶段验收和成果评价量规，加强对基础薄弱学生的跟踪指导；落实设备、账号、数据和场地安全检查，细化劳动教育记录，并将优秀成果与典型问题纳入下一轮实训资源。",
    }
    for index, text in content.items():
        _set_paragraph(paragraphs[index], text, body_font, 10.5)
        paragraphs[index].paragraph_format.keep_together = True
    _set_paragraph(paragraphs[74], f"《{offering['course_name']}》实训作业、成果、报告等", "黑体", 16)
    paragraphs[74].alignment = WD_ALIGN_PARAGRAPH.CENTER
    paragraphs[74].paragraph_format.page_break_before = True
    _set_paragraph(paragraphs[75], "本部分归档学生实训任务书、过程记录、源文件、成果截图或数据结果、实训报告、展示材料、评价表及优秀成果。电子资料按“班级-学号-姓名-任务名称”统一命名；纸质材料按任务顺序整理。实际学生成果由指导教师在实训结束后补充装订或作为电子附件归档。", body_font, 10.5)
    paragraphs[75].paragraph_format.keep_together = True
    paragraphs[89].paragraph_format.page_break_before = True

    overview = document.tables[0]
    values = (offering["major"], class_name, student_count, hours, date_text, location, teacher, "校内实训，无需外出；如临时安排校外活动，另行履行审批和安全告知手续。")
    for row, value in zip(overview.rows, values):
        _set_cell(row.cells[1], value)
    schedule = document.tables[1]
    modules = unit_names[:5] or ["任务准备与安全教育", "核心任务实践", "成果完善与评价"]
    while len(schedule.rows) - 2 < len(modules):
        schedule.add_row()
    module_hours = max(1, round(float(hours) / len(modules))) if str(hours).isdigit() else ""
    project_cells = []
    project_number = 0
    for index, module in enumerate(modules, 2):
        row = schedule.rows[index]
        project_cell = row.cells[0]
        if not any(project_cell._tc is existing for existing in project_cells):
            project_number += 1
            project_cells.append(project_cell._tc)
            project_label = f"项目{'一' if project_number == 1 else '二' if project_number == 2 else project_number}"
            _set_cell(project_cell, project_label)
        for cell, value in zip(row.cells[1:], (module, module_hours, location, teacher, "")):
            _set_cell(cell, value)
    guide = document.tables[2]
    for row_index, pair in enumerate(((offering["major"], class_name), (student_count, location), (hours, date_text), (teacher, ""))):
        _set_cell(guide.rows[row_index].cells[1], pair[0])
        _set_cell(guide.rows[row_index].cells[3], pair[1])
    detail = document.tables[3]
    while len(detail.rows) - 1 < len(modules):
        detail.add_row()
    for index, module in enumerate(modules, 1):
        row = detail.rows[index]
        date_value = date_text if index == 1 else "按实训日程"
        expected = f"完成“{module}”任务成果及过程记录"
        criterion = "技术正确、过程规范、成果完整；首项含安全与劳动教育记录" if index == 1 else "按任务书与评价量规验收"
        for cell, value in zip(row.cells, (index, date_value, "按课表", location, module, expected, criterion)):
            _set_cell(cell, value)

    for table in document.tables:
        for row in table.rows:
            _prevent_row_split(row)
    _repeat_header(schedule.rows[0])
    _repeat_header(schedule.rows[1])
    _repeat_header(detail.rows[0])
    for index in (90, 93, 96, 99):
        paragraphs[index].paragraph_format.keep_with_next = True

    output_root = Path(store.get_setting("output_root", ROOT / "生成结果"))
    output_dir = output_root / offering["term"] / offering["course_name"] / "实训资料"
    output_dir.mkdir(parents=True, exist_ok=True)
    class_suffix = f" {class_name}" if class_name else ""
    output_path = output_dir / f"{offering['term']}《{offering['course_name']}》{class_suffix}实训资料.docx"
    document.core_properties.title = f"《{offering['course_name']}》实训资料"
    document.core_properties.author = teacher
    document.core_properties.modified = datetime.now()
    document.save(output_path)
    store.execute(
        "INSERT INTO training_documents(offering_id,class_name,source_dir,output_path,generated_at) VALUES (?,?,?,?,?)",
        (offering_id, class_name, str(Path(source_dir).resolve()), str(output_path), datetime.now().strftime("%Y-%m-%d %H:%M:%S")),
    )
    return output_path
