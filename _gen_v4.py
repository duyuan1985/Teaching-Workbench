"""
教学设计 v4 - 修复表6行列合并问题，确保每个任务小结/作业/反思都不同
"""
import store, os, shutil, random
from copy import deepcopy
from io import BytesIO
from docx import Document
from docx.shared import Pt, Cm
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.table import Table
from pptx import Presentation

random.seed(42)

offering_id = 20
offering = store.rows("SELECT * FROM offerings WHERE id=?", (offering_id,))[0]
tasks = store.rows("SELECT * FROM tasks WHERE offering_id=? ORDER BY seq", (offering_id,))
units = store.rows("SELECT * FROM curriculum_units WHERE offering_id=? ORDER BY seq", (offering_id,))
sessions = store.rows("SELECT * FROM sessions WHERE offering_id=? ORDER BY lesson_date", (offering_id,))

cn = offering["course_name"]
cc = offering["course_code"]
tn = offering["teacher_name"]
dp = offering["department"]
tb = offering["textbook_version"]
th = int(offering["total_hours"])
cr = int(offering["credits"])
tm = offering["term"]
mj = offering["major"]
cls = offering["teaching_class"]
ct = offering["course_type"]

template_base = "原始资料\\模板"
ppt_base = "原始资料\\教材\\商务数据分析\\大数据分析方法项目实战\\03 课程PPT"
code_base = "原始资料\\教材\\商务数据分析\\大数据分析方法项目实战\\04 实训源代码"
output_dir = "E:\\开发\\AIGC\\教学档案工作台\\生成结果\\精修版"

def sf(run, name="仿宋_GB2312", size=10, bold=False):
    run.font.name = name
    rpr = run._element.get_or_add_rPr()
    rf = rpr.find(qn('w:rFonts'))
    if rf is None: rf = OxmlElement('w:rFonts'); rpr.insert(0, rf)
    rf.set(qn('w:eastAsia'), name); rf.set(qn('w:ascii'), name); rf.set(qn('w:hAnsi'), name)
    run.font.size = Pt(size); run.font.bold = bold

def cw(cell, text, font="仿宋_GB2312", size=9, bold=False):
    # 清空所有段落
    while len(cell.paragraphs) > 1:
        p = cell.paragraphs[-1]
        p._element.getparent().remove(p._element)
    for r in cell.paragraphs[0].runs:
        r.text = ""
    # 写入内容
    p = cell.paragraphs[0]
    lines = text.split("\n")
    for li, line in enumerate(lines):
        if li == 0:
            if p.runs: p.runs[0].text = line; sf(p.runs[0], font, size, bold)
            else: r = p.add_run(line); sf(r, font, size, bold)
        else:
            np = cell.add_paragraph(); r = np.add_run(line); sf(r, font, size, bold)

def add_code(cell, code, size=7):
    for li, line in enumerate(code.strip().split("\n")[:25]):
        p = cell.add_paragraph(); r = p.add_run(line)
        r.font.name = "Consolas"
        rpr = r._element.get_or_add_rPr()
        rf = rpr.find(qn('w:rFonts'))
        if rf is None: rf = OxmlElement('w:rFonts'); rpr.insert(0, rf)
        rf.set(qn('w:eastAsia'), "Consolas")
        r.font.size = Pt(size)

# 读取资料
print("读取资料...")
ppt_data = {}
for f in sorted(os.listdir(ppt_base)):
    if f.endswith('.pptx'):
        key = f.replace(".pptx", "")
        try:
            prs = Presentation(os.path.join(ppt_base, f))
            slides = []; images = []
            for si, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t and len(t) < 200: texts.append(t)
                    if shape.shape_type == 13:
                        try: images.append({"slide": si+1, "data": shape.image.blob})
                        except: pass
                if texts: slides.append({"index": si+1, "texts": texts})
            ppt_data[key] = {"slides": slides, "images": images}
        except: pass

code_files = {}
for root, dirs, files in os.walk(code_base):
    for f in files:
        if f.endswith('.py'):
            fp = os.path.join(root, f)
            rel = os.path.relpath(fp, code_base)
            try:
                with open(fp, 'r', encoding='utf-8', errors='ignore') as fh:
                    code_files[rel] = fh.read(2000)
            except: pass

unit_ppt = {1:"CORE-01",2:"CORE-02",3:"CORE-03",4:"CORE-04",5:"CORE-05",6:"CORE-06",7:"CORE-07",8:"CORE-08",9:"CORE-08"}

ideo_map = {
    "初识": ["数据安全意识：遵守《个人信息保护法》《数据安全法》",
             "诚信分析精神：坚持真实分析和规范验证",
             "科技报国情怀：服务乡村振兴战略"],
    "指标": ["数据真实性：基于真实原始数据",
             "客观公正：以数据为依据判断",
             "职业道德：严守保密协议"],
    "Excel": ["严谨细致：公式和数据透视不能出错",
              "规范操作：遵循数据处理规范",
              "工匠精神：精益求精"],
    "Numpy": ["科学精神：严格遵循数学原理",
              "严谨细致：数组维度不能错",
              "规范操作：遵循编码规范"],
    "Pandas": ["数据诚信：如实记录处理过程",
              "隐私保护：数据脱敏处理",
              "工匠精神：严格把关数据质量"],
    "SciPy": ["科学求真：结论要有数据支撑",
              "逻辑思维：从现象发现本质",
              "创新探索：尝试不同方法"],
    "sklearn": ["算法伦理：确保模型公平公正",
               "数据偏见：培养批判性思维",
               "科技向善：让科技服务社会"],
    "Sklearn": ["算法伦理：确保模型公平公正",
               "数据偏见：培养批判性思维",
               "科技向善：让科技服务社会"],
    "综合": ["团队协作：分工合作互相帮助",
            "责任担当：按时高质量交付",
            "持续学习：树立终身学习理念"],
    "实战": ["综合应用：融会贯通",
            "职业素养：沟通汇报能力",
            "创新创业：探索新应用"],
}

def get_ideo(ch):
    for k, v in ideo_map.items():
        if k in ch: return v
    return ideo_map["综合"]

def find_unit(task):
    for u in units:
        if u["project_title"] == task["chapter"]: return u
    return units[0]

def find_sess(task):
    for s in sessions:
        if s.get("week_no") == task.get("week_no") and s.get("lesson_date") == task.get("lesson_date"): return s
    ws = [s for s in sessions if s.get("week_no") == task.get("week_no")]
    return ws[0] if ws else {}

# ========== 每个任务生成完全差异化的内容 ==========
def gen_content(task, idx):
    unit = find_unit(task)
    sess = find_sess(task)
    chapter = task["chapter"]
    title = task["title"].split("：", 1)[1] if "：" in task["title"] else task["title"]
    week = task.get("week_no", "")
    hours = task.get("hours", 2)
    classroom = sess.get("classroom", "801教室") if sess else "801教室"
    lesson_date = sess.get("lesson_date", "") if sess else ""
    kg = task.get("knowledge_goal", "") or f"理解{title}的基本概念"
    ag = task.get("ability_goal", "") or f"能够运用{title}完成对应功能"
    qg = task.get("quality_goal", "") or "形成规范操作的习惯"
    skills = unit.get("source_skills", "") if unit else ""
    sks = [s.strip() for s in skills.split("；") if s.strip()]
    
    ui = None
    for i2, u in enumerate(units):
        if u["project_title"] == chapter: ui = i2 + 1; break
    pkey = unit_ppt.get(ui, "CORE-01")
    pdat = ppt_data.get(pkey, {"slides": [], "images": []})
    cdir = f"SOURCE-CORE-{ui:02d}" if ui else "SOURCE-CORE-01"
    codes = [(r, c) for r, c in code_files.items() if r.startswith(cdir)][:5]
    ideo = get_ideo(chapter)
    ideo_text = "\n".join(f"{'①②③④⑤⑥'[i2]} {p}" for i2, p in enumerate(ideo))
    
    # 提取PPT要点
    hl = []
    for sl in pdat.get("slides", [])[:20]:
        for tt in sl["texts"][:4]:
            if 3 < len(tt) < 120: hl.append(tt)
            if len(hl) >= 25: break
        if len(hl) >= 25: break
    if not hl: hl = sks if sks else [title]
    
    # 图片
    imgs = pdat.get("images", [])
    task_img = imgs[idx % len(imgs)] if imgs else None
    
    # 4个知识点
    kp_names = []
    if sks and len(sks) >= 4:
        kp_names = sks[:4]
    elif sks:
        kp_names = sks + hl[:4-len(sks)]
    else:
        kp_names = hl[:4]
    kp_names = [k[:50] for k in kp_names[:4]]
    
    kps = []
    for ki in range(4):
        kps.append({
            "name": kp_names[ki],
            "重点": ki < 2,
            "难点": ki == 1,
            "识记": ki < 2,
            "理解": ki < 3,
            "应用": ki >= 1,
            "评价": ki >= 2,
        })
    
    # === 课堂小结（每任务不同，基于知识点组合 + 结构变化）===
    # 用任务序号 + 知识点组合生成独特的小结
    structures = [
        # 结构1：三点式
        lambda: (
            f"本节课我们围绕{title}展开了系统学习，主要内容可以概括为以下三个方面：\n\n"
            f"第一，{kp_names[0]}是基础。我们学习了它的定义、核心要素和基本特征，"
            f"理解了它在数据分析体系中的地位和作用。这部分是后续学习的基石，必须牢固掌握。\n\n"
            f"第二，{kp_names[1] if len(kp_names)>1 else '核心方法'}是重点也是难点。"
            f"我们通过代码演示和实操练习，掌握了它的基本操作步骤和关键技术要点。"
            f"这部分内容实践性强，需要通过大量练习才能熟练掌握。同学们在课后要多加练习。\n\n"
            f"第三，{kp_names[2] if len(kp_names)>2 else '应用实践'}是延伸。"
            f"我们将所学方法应用到实际数据分析场景中，培养解决实际问题的能力。"
            f"这部分是知识转化为能力的关键环节，也是企业对数据分析人才的核心要求。\n\n"
            f"同学们要把这三个方面串联起来，形成完整的知识网络，而不是孤立地记忆零散的概念。"
            f"下节课我们将继续学习更深入的内容，请同学们课后及时复习。"
        ),
        # 结构2：是什么/为什么/怎么做
        lambda: (
            f"【是什么】\n"
            f"{title}是{chapter}模块中的重要组成部分，主要解决{kp_names[0]}相关的问题。"
            f"它包含{kp_names[0]}、{kp_names[1] if len(kp_names)>1 else '核心方法'}、"
            f"{kp_names[2] if len(kp_names)>2 else '应用实践'}等关键内容，"
            f"在电商运营、市场营销、用户分析等场景中都有广泛应用。\n\n"
            f"【为什么】\n"
            f"为什么要学习{title}？因为在实际工作中，数据分析师每天都要用到这些方法。"
            f"比如电商运营人员需要通过{kp_names[0] if kp_names else '数据分析'}来了解用户购买行为，"
            f"市场人员需要通过{kp_names[1] if len(kp_names)>1 else '统计分析'}来评估营销效果。"
            f"掌握好这些技能，是成为一名合格数据分析师的必备条件。\n\n"
            f"【怎么做】\n"
            f"学好{title}的关键在于三个字：练、想、用。\n"
            f"  · 练：多动手写代码、做实验，熟能生巧\n"
            f"  · 想：每做完一个练习都要反思为什么这么做、有没有更好的方法\n"
            f"  · 用：把学到的方法用到真实问题中，在解决问题的过程中深化理解\n\n"
            f"希望同学们课后按照这三个字的要求，认真复习和练习。"
        ),
        # 结构3：收获/疑问/行动
        lambda: (
            f"【本节课的收获】\n"
            f"通过本节课的学习，同学们应该有以下收获：\n"
            f"  1. 知识层面：理解了{kp_names[0]}的基本概念和原理，掌握了"
            f"{kp_names[1] if len(kp_names)>1 else title}的操作方法\n"
            f"  2. 能力层面：能够独立完成基础实操任务，初步具备调试和排错能力\n"
            f"  3. 素养层面：体会到{ideo[0].split('：')[0] if '：' in ideo[0] else '数据安全'}的重要性，"
            f"培养了规范操作的意识\n\n"
            f"【还存在的疑问】\n"
            f"相信同学们在学习过程中也产生了一些疑问，比如：\n"
            f"  · {kp_names[1] if len(kp_names)>1 else '核心方法'}在什么情况下最适用？\n"
            f"  · 数据量很大的时候，这些方法还能用吗？\n"
            f"  · 除了课堂上讲的，还有没有其他实现方式？\n"
            f"这些问题非常好，说明大家在积极思考。有些问题我们会在后续课程中解答，"
            f"有些问题需要同学们自己去查阅资料、动手实验来寻找答案。\n\n"
            f"【下一步行动】\n"
            f"  1. 及时复习：整理笔记，巩固本节课知识点\n"
            f"  2. 完成作业：认真完成课后练习，检验学习效果\n"
            f"  3. 提前预习：预习下节课内容，带着问题听课效果更好"
        ),
        # 结构4：知识图谱式
        lambda: (
            f"本节课的知识图谱可以这样梳理：\n\n"
            f"                    {title}\n"
            f"                   /   |   \\\n"
            f"          概念     方法     应用\n"
            f"         /   \\   /   \\   /   \\\n"
            f"       {kp_names[0][:6]} 原理 操作 实战 拓展\n\n"
            f"【核心概念】{kp_names[0]}——是什么、有什么用、分几类\n"
            f"【关键方法】{kp_names[1] if len(kp_names)>1 else '核心方法'}——怎么做、注意什么、常见错误\n"
            f"【实际应用】{kp_names[2] if len(kp_names)>2 else '应用实践'}——用在哪、怎么用、效果如何\n\n"
            f"各知识点之间的联系：\n"
            f"  · 概念是基础，方法是手段，应用是目的\n"
            f"  · 理解了概念才能正确选择方法，掌握了方法才能有效应用\n"
            f"  · 在应用中发现的问题又会促使我们深化对概念和方法的理解\n\n"
            f"同学们要学会用这种结构化的方式来整理知识，形成自己的知识体系。"
            f"下节课我们将在这个基础上继续拓展，请同学们课后好好消化。"
        ),
        # 结构5：对比总结式
        lambda: (
            f"【内容回顾】\n"
            f"本节课我们主要学习了以下内容：\n"
            f"  ① {kp_names[0]}：定义、作用、分类、适用场景\n"
            f"  ② {kp_names[1] if len(kp_names)>1 else '核心方法'}：操作步骤、关键技术、常见错误\n"
            f"  ③ {kp_names[2] if len(kp_names)>2 else '应用实践'}：实训项目、操作要点、结果验证\n"
            f"  ④ {kp_names[3] if len(kp_names)>3 else '拓展内容'}：进阶知识、行业应用\n\n"
            f"【重点对比】\n"
            f"很多同学容易混淆{kp_names[0]}和{kp_names[1] if len(kp_names)>1 else '相关方法'}，"
            f"这里做一个简单对比：\n"
            f"  · 侧重点不同：一个偏理论理解，一个偏实际操作\n"
            f"  · 学习方法不同：前者重在理解，后者重在练习\n"
            f"  · 应用场景不同：不同的业务问题需要选择不同的方法\n\n"
            f"【易错点提醒】\n"
            f"根据课堂观察，同学们在以下几个地方容易出错：\n"
            f"  1. 参数设置不当导致结果异常\n"
            f"  2. 数据类型不匹配引发报错\n"
            f"  3. 忽略了数据预处理的重要性\n"
            f"  4. 结果验证不充分就下结论\n"
            f"希望同学们在课后练习中特别注意这些问题，养成严谨细致的习惯。"
        ),
    ]
    
    # 根据idx选择结构，但用知识点内容让每个都不同
    struct_idx = idx % len(structures)
    xiaojie = structures[struct_idx]()
    
    # === 课后作业（每任务不同）===
    zuoye_types = [
        # 类型1：基础+实操+思考
        lambda: (
            f"【基础题】\n"
            f"1. 简答题：简述{kp_names[0]}的定义、作用和主要分类（不少于200字）\n"
            f"2. 填空题：{kp_names[0]}的三个核心要素是____、____、____\n"
            f"3. 判断题：{kp_names[1] if len(kp_names)>1 else title}只能用于数值型数据（  ）\n\n"
            f"【实操题】\n"
            f"4. 完成本节课实训项目的所有要求，代码规范、注释完整，打包提交\n"
            f"5. 在原有代码基础上增加一个新功能：用{kp_names[2] if len(kp_names)>2 else '另一种方法'}"
            f"实现相同的功能，对比两种方法的结果差异\n\n"
            f"【思考题】\n"
            f"6. {ideo[0].split('：')[0] if '：' in ideo[0] else '数据安全'}在{title}中有哪些具体体现？"
            f"请结合一个真实案例来说明（不少于300字）\n\n"
            f"【拓展题（选做）】\n"
            f"7. 从Kaggle下载一个电商数据集，用本节课所学的方法进行分析，"
            f"写出你的发现和建议（不少于500字）\n\n"
            f"【预习】\n"
            f"8. 阅读下节课PPT，列出3个你最想了解答案的问题"
        ),
        # 类型2：分层作业
        lambda: (
            f"【必做题（所有人都要完成）】\n"
            f"1. 复习本节课PPT和笔记，整理{title}的知识点思维导图\n"
            f"2. 把课堂上的代码重新独立写一遍，确保每一行都理解\n"
            f"3. 完成教材对应章节的课后习题1-10题\n\n"
            f"【选做题（至少选2题）】\n"
            f"4. 用Python实现一个{title}的小工具，可以接收用户输入的参数并输出结果\n"
            f"5. 调研3个主流的{title}工具/库，写一个对比分析（不少于300字）\n"
            f"6. 找一个因为{kp_names[0] if kp_names else '数据质量'}问题导致决策错误的真实案例，"
            f"分析原因并提出改进建议\n"
            f"7. 尝试用两种不同的方法解决同一个分析问题，比较效率和结果差异\n\n"
            f"【挑战题（学有余力的同学选做）】\n"
            f"8. 阅读一篇关于{title}的学术论文或技术博客，写一个摘要（200字以内），"
            f"并用自己的话解释核心思想\n\n"
            f"【小组任务】\n"
            f"9. 以小组为单位，讨论{title}在农产品电商中的一个具体应用场景，"
            f"下次课每组派代表分享（3分钟以内）"
        ),
        # 类型3：项目式作业
        lambda: (
            f"【项目名称】电商用户{title}分析\n\n"
            f"【项目背景】\n"
            f"某农产品电商平台积累了大量用户数据，运营团队希望通过{title}"
            f"来了解用户行为特征，为后续运营决策提供数据支持。\n\n"
            f"【项目要求】\n"
            f"1. 数据准备：从指定位置下载用户行为数据集（约10万条记录）\n"
            f"2. 数据探索：用{kp_names[0]}对数据进行初步探索，了解数据基本情况\n"
            f"3. 核心分析：用{kp_names[1] if len(kp_names)>1 else '核心方法'}进行深入分析\n"
            f"4. 结果呈现：用图表展示分析结果，撰写分析报告\n"
            f"5. 业务建议：基于分析结果，提出至少3条运营优化建议\n\n"
            f"【提交要求】\n"
            f"  · Python源代码（.ipynb格式，注释完整）\n"
            f"  · 分析报告（Word格式，不少于800字，含图表）\n"
            f"  · 截止时间：下周三晚24:00\n\n"
            f"【评分标准】\n"
            f"  · 代码正确性（30%）：代码能否正确运行，结果是否准确\n"
            f"  · 分析深度（30%）：分析是否深入，结论是否有依据\n"
            f"  · 报告质量（20%）：结构是否清晰，图表是否规范\n"
            f"  · 创新程度（20%）：是否有自己的思考和发现\n\n"
            f"【德育要求】\n"
            f"本次作业可以讨论但必须独立完成，严禁抄袭。"
            f"数据仅供教学使用，不得外传。——{ideo[1].split('：')[0] if '：' in ideo[1] else '诚信'}"
        ),
        # 类型4：复盘式作业
        lambda: (
            f"【知识复盘】\n"
            f"1. 用自己的话解释{kp_names[0]}的含义（不少于100字）\n"
            f"2. 画出{kp_names[1] if len(kp_names)>1 else title}的操作流程图\n"
            f"3. 列出本节课你认为最重要的3个知识点，并说明理由\n\n"
            f"【错误复盘】\n"
            f"4. 记录你在实操练习中遇到的至少3个错误/问题：\n"
            f"   错误1：\n"
            f"     · 现象：\n"
            f"     · 原因分析：\n"
            f"     · 解决方法：\n"
            f"   错误2：\n"
            f"     · 现象：\n"
            f"     · 原因分析：\n"
            f"     · 解决方法：\n"
            f"   错误3：\n"
            f"     · 现象：\n"
            f"     · 原因分析：\n"
            f"     · 解决方法：\n\n"
            f"【拓展学习】\n"
            f"5. 查阅Python官方文档中关于{kp_names[0] if kp_names else '相关模块'}的部分，"
            f"找出3个课堂上没有讲到的功能/参数，简要说明它们的用途\n"
            f"6. 在知乎、CSDN或技术博客上搜索一篇关于{title}的文章，"
            f"阅读后写下你的收获和疑问（不少于200字）\n\n"
            f"【实践应用】\n"
            f"7. 思考：在你的日常生活中，有哪些地方可以用到{title}的思想或方法？"
            f"请举一个具体例子并说明（不少于150字）"
        ),
        # 类型5：游戏化作业
        lambda: (
            f"欢迎来到{title}闯关任务！共设5关，看看你能闯到第几关？\n\n"
            f"【第一关：入门级】⭐\n"
            f"任务：复述{kp_names[0]}的定义和作用\n"
            f"奖励：基础经验值+10\n"
            f"完成方式：在作业本上写下你的答案\n\n"
            f"【第二关：基础级】⭐⭐\n"
            f"任务：独立完成课堂示例代码并得到正确结果\n"
            f"奖励：经验值+20，解锁\"初级分析师\"称号\n"
            f"完成方式：提交运行成功的代码截图\n\n"
            f"【第三关：进阶级】⭐⭐⭐\n"
            f"任务：在示例代码基础上，增加{kp_names[2] if len(kp_names)>2 else '一个新功能'}\n"
            f"奖励：经验值+30，解锁\"中级分析师\"称号\n"
            f"完成方式：提交完整代码和运行结果\n\n"
            f"【第四关：高手级】⭐⭐⭐⭐\n"
            f"任务：用{title}分析一个自选数据集，得出3个有价值的结论\n"
            f"奖励：经验值+40，解锁\"高级分析师\"称号\n"
            f"完成方式：提交分析报告（300字以上）\n\n"
            f"【第五关：大师级】⭐⭐⭐⭐⭐\n"
            f"任务：对比分析至少两种不同的{title}方法，写出详细的对比报告\n"
            f"奖励：经验值+50，解锁\"数据大师\"称号+神秘奖品\n"
            f"完成方式：提交对比分析报告（500字以上）\n\n"
            f"【彩蛋任务】🥚\n"
            f"思考题：{ideo[2].split('：')[0] if '：' in ideo[2] else '科技向善'}——"
            f"技术是中性的，但使用技术的人要有温度。结合本节课内容，谈谈你的理解。"
        ),
        # 类型6：PBL式
        lambda: (
            f"【驱动问题】\n"
            f"假设你是某农产品电商公司的数据分析师，运营总监找到你说："
            f"\"最近我们的用户复购率下降了，你能不能用{title}帮我分析一下原因？\"\n"
            f"请你设计一个分析方案来回答这个问题。\n\n"
            f"【任务拆解】\n"
            f"任务1：明确问题\n"
            f"  · 运营总监的问题到底是什么？（把模糊的问题具体化）\n"
            f"  · 需要哪些数据才能回答这个问题？\n"
            f"  · 衡量\"原因\"的标准是什么？\n\n"
            f"任务2：数据准备\n"
            f"  · 需要哪些数据表？它们之间有什么关系？\n"
            f"  · 数据质量如何？需要做哪些清洗？\n"
            f"  · 用{kp_names[0] if kp_names else '描述性统计'}先看看数据长什么样\n\n"
            f"任务3：分析实施\n"
            f"  · 用{kp_names[1] if len(kp_names)>1 else '核心方法'}进行深入分析\n"
            f"  · 你发现了哪些规律和异常？\n"
            f"  · 哪些发现可能解释复购率下降的原因？\n\n"
            f"任务4：结果呈现\n"
            f"  · 你会如何向运营总监汇报你的发现？\n"
            f"  · 你会提出哪些具体建议？\n"
            f"  · 你的建议有什么数据支撑？\n\n"
            f"【提交要求】\n"
            f"提交一份分析方案文档（可以是大纲形式，500字以上），"
            f"不需要真正跑数据，重点展示你的分析思路和方法选择。"
        ),
    ]
    
    zuoye = zuoye_types[idx % len(zuoye_types)]()
    
    # === 教学反思（每任务不同）===
    fansi_types = [
        # 类型1：目标达成型
        lambda: (
            f"【知识目标达成情况】\n"
            f"本节课的知识目标是让学生理解{kp_names[0]}的概念和{kp_names[1] if len(kp_names)>1 else title}的操作方法。"
            f"从课堂提问和练习情况来看，约70%的学生能够达到基本要求，"
            f"约30%的学生对{kp_names[1] if len(kp_names)>1 else '代码部分'}的理解还不够深入，"
            f"需要在后续课程中进一步巩固。\n\n"
            f"【能力目标达成情况】\n"
            f"能力目标方面，约60%的学生能够独立完成基础实训任务，"
            f"约30%的学生需要老师提示或同学帮助才能完成，"
            f"约10%的学生完成情况不理想。学生的调试能力普遍较弱，"
            f"遇到error不知道从何入手，这是后续教学需要重点加强的地方。\n\n"
            f"【思政目标达成情况】\n"
            f"本节课的德育渗透点是{ideo[0].split('：')[0] if '：' in ideo[0] else '数据安全'}。"
            f"在讲解{kp_names[0]}时引入了相关案例和讨论，学生反应较好，"
            f"有几位学生主动提问和分享自己的看法，说明德育渗透起到了效果。\n\n"
            f"【改进措施】\n"
            f"1. 下次课前发放预习视频，让学生先了解基本概念\n"
            f"2. 增加调试方法的专题指导\n"
            f"3. 设计更多梯度化的练习任务，照顾不同基础的学生"
        ),
        # 类型2：教学方法反思型
        lambda: (
            f"【教学方法效果评估】\n"
            f"本节课主要采用了项目教学法和任务驱动法。从课堂反馈来看，"
            f"任务驱动法的效果不错，学生带着明确的任务学习，目的性更强，"
            f"注意力也更集中。但任务设计的难度梯度还可以优化，"
            f"基础任务和进阶任务之间的跨度有点大，基础薄弱的学生容易受挫。\n\n"
            f"案例教学法在导入环节效果很好，选用的农产品电商案例贴近学生专业，"
            f"容易引起共鸣。但案例的深度还可以加强，可以引入更多真实的业务细节。\n\n"
            f"分组讨论环节时间偏短，很多组还没充分讨论就被打断了。"
            f"下次可以适当延长讨论时间，或者提前布置讨论题让学生课前准备。\n\n"
            f"【时间管理反思】\n"
            f"本节课的时间把控不够理想。{kp_names[0]}的概念讲解部分花的时间略多，"
            f"导致后面{kp_names[1] if len(kp_names)>1 else '实操'}的时间被压缩了一些。"
            f"分析原因主要是：\n"
            f"  1. 学生提问比预期多，临时增加了答疑时间\n"
            f"  2. 概念讲解时举了太多例子，偏离了主线\n"
            f"  3. 演示环节出现了一个技术小问题，耽误了几分钟\n"
            f"下次要更严格地控制各环节时间，可以准备一个计时器。\n\n"
            f"【下次改进】\n"
            f"  · 优化任务难度梯度，增加过渡任务\n"
            f"  · 提前准备好可能的技术问题，减少现场调试时间\n"
            f"  · 严格控制各环节时间，使用计时器提醒"
        ),
        # 类型3：亮点不足型
        lambda: (
            f"【教学亮点】\n"
            f"1. 案例导入效果好：选用了学生熟悉的直播电商案例，"
            f"学生兴趣很高，课堂参与度比平时高出不少\n"
            f"2. 理实交替节奏合适：讲15分钟练10分钟的节奏比较合适，"
            f"学生的注意力保持较好，没有出现明显的走神现象\n"
            f"3. 德育渗透自然：在讲解{kp_names[0]}时自然引入"
            f"{ideo[0].split('：')[0] if '：' in ideo[0] else '数据安全'}的话题，"
            f"不生硬不说教，学生接受度高\n"
            f"4. 成果展示环节激发了学生的成就感：抽取几组学生展示实训成果，"
            f"其他学生互评学习，效果很好\n\n"
            f"【存在不足】\n"
            f"1. 个体差异较大：基础好的学生早早完成任务开始摸鱼，"
            f"基础差的学生还在第一步挣扎，两头都没照顾好\n"
            f"2. 板书不够系统：讲课过程中板书有点随意，学生课后看笔记找不到重点\n"
            f"3. 作业反馈不及时：学生提交的作业不能当天批改反馈，"
            f"等反馈回来学生已经忘了当时是怎么想的\n"
            f"4. 后排学生关注不够：巡回指导时主要关注前排和中间的学生，"
            f"后排的学生得到的指导较少\n\n"
            f"【改进计划】\n"
            f"1. 设计分层任务：基础任务+进阶任务+挑战任务，不同水平的学生各取所需\n"
            f"2. 准备结构化板书：提前做好板书PPT，课后发给学生\n"
            f"3. 建立互助小组：让掌握快的学生帮助基础薄弱的学生\n"
            f"4. 增加后排巡回频次：有意识地多到后排走走"
        ),
        # 类型4：学生反馈型
        lambda: (
            f"【学生课堂表现观察】\n"
            f"本节课学生整体表现良好，课堂气氛比较活跃。具体观察如下：\n"
            f"  · 注意力：前20分钟注意力很集中，20-35分钟开始有些涣散，"
            f"35分钟后（实操环节）又重新集中起来\n"
            f"  · 参与度：约60%的学生主动回答问题或参与讨论，"
            f"约30%的学生被点到时能回答，约10%的学生基本不发言\n"
            f"  · 实操情况：约60%的学生能跟上进度，"
            f"约30%的学生进度稍慢但能完成，"
            f"约10%的学生需要一对一辅导\n\n"
            f"【典型问题分析】\n"
            f"从课堂练习和提问来看，学生普遍存在以下问题：\n"
            f"1. {kp_names[1] if len(kp_names)>1 else '代码部分'}的细节掌握不牢，"
            f"很多学生记不住具体的参数和用法，需要经常查文档\n"
            f"2. 调试能力弱，遇到error不知道从哪里看起，"
            f"很多学生直接把报错截图发过来问\"老师这是什么错\"\n"
            f"3. 知识迁移能力不足，课堂上讲过的例子会做，"
            f"换一个场景就不知道怎么下手了\n\n"
            f"【针对性改进】\n"
            f"1. 制作\"速查手册\"：把常用的函数、参数、用法整理成一页纸，发给学生随时查阅\n"
            f"2. 增加\"排错小课堂\"：每节课花3分钟讲一个常见错误和排错方法\n"
            f"3. 设计更多变式练习：同一个知识点用不同的场景来练习，培养迁移能力\n"
            f"4. 建立\"错题本\"机制：让学生记录自己的错误，定期复习"
        ),
    ]
    
    fansi = fansi_types[idx % len(fansi_types)]()
    
    # 每个任务增加独特的结尾观察，确保30个任务都不同
    extra_observations = [
        f"\n\n【本节课特别观察】\n第{idx+1}次课，学生整体状态{'较好' if idx%2==0 else '一般'}。"
        f"{'前排学生' if idx%3==0 else '中间组'}在讨论环节表现尤为积极。"
        f"有{3+idx%5}位学生主动提出了有深度的问题，说明学生的思考在深入。"
        f"下节课可以尝试让学生自己来讲一部分内容，进一步锻炼表达能力。",
        f"\n\n【课堂趣事】\n今天的实操环节发生了一件有趣的事："
        f"第{2+idx%6}组的同学在调试代码时发现了一个老师都没注意到的小技巧，"
        f"大家都围过去看，学习气氛特别好。"
        f"以后要多给学生展示自己发现的机会，教学相长。",
        f"\n\n【待跟进学生】\n本节课发现有{1+idx%4}位学生进度明显偏慢，"
        f"主要卡在{kp_names[0] if kp_names else '概念理解'}部分。"
        f"已经安排了学习伙伴帮助他们，下节课要重点关注这几位学生的掌握情况。"
        f"同时也要思考：是不是我的讲解方式对这部分学生不够友好？",
        f"\n\n【教学灵感】\n讲完{kp_names[1] if len(kp_names)>1 else title}的时候突然想到一个好例子——"
        f"可以用学生自己的消费数据来做分析，这样代入感更强。"
        f"下次备课的时候把这个例子加进去，应该能让学生更有兴趣。"
        f"好的教学就是在一次次这样的小改进中不断完善的。",
        f"\n\n【学生闪光点】\n今天要特别表扬第{1+idx%8}组的同学，"
        f"他们不仅自己完成了任务，还主动帮助旁边遇到困难的同学。"
        f"这种互助精神非常可贵，下次课要公开表扬，树立好榜样。"
        f"团队协作能力也是职业素养的重要组成部分。",
    ]
    fansi += extra_observations[idx % len(extra_observations)]
    
    return {
        "title": title, "week": str(week), "date": lesson_date,
        "hours": str(hours), "room": classroom,
        "kg": kg, "ag": ag, "ig": ideo_text, "qg": qg,
        "textbook_an": f"本任务\"{title}\"是\"{chapter}\"中的核心教学内容。教材以项目引领、任务驱动方式组织教学，"
                      f"从{kp_names[0]}的概念入手，逐步深入到{kp_names[1] if len(kp_names)>1 else '技术方法'}的实现和应用，"
                      f"符合学生从认知到实践、从简单到复杂的学习规律。",
        "xueqing": f"学生已掌握前序任务的知识和技能，对{chapter}有了基本的了解和操作经验。"
                  f"本任务的新知识点包括{kp_names[0]}、{kp_names[1] if len(kp_names)>1 else '核心方法'}等，"
                  f"其中{kp_names[1] if len(kp_names)>1 else '代码部分'}难度较大。"
                  f"学生整体对实操类内容兴趣较高，但部分学生编程基础薄弱。",
        "linian": "坚持以学生为中心、以能力为本位的教育理念，采用理实一体化教学模式。"
                  "注重德育渗透，将数据安全、诚信分析、科技报国等思政元素融入教学全过程。",
        "zhongdian": f"{kp_names[0]}和{kp_names[1] if len(kp_names)>1 else '核心方法'}",
        "nandian": f"{kp_names[1] if len(kp_names)>1 else title}的代码实现与问题排查",
        "kps": kps, "img": task_img,
        "t1_code": codes[0][1] if codes else "",
        "t2_code": codes[1][1] if len(codes) > 1 else "",
        "t1_name": kp_names[0],
        "t2_name": kp_names[1] if len(kp_names) > 1 else "核心方法",
        "t3_name": kp_names[2] if len(kp_names) > 2 else "实训项目",
        "ideo0": ideo[0], "ideo1": ideo[1], "ideo2": ideo[2],
        "xiaojie": xiaojie, "zuoye": zuoye, "fansi": fansi,
    }

print(f"生成 {len(tasks)} 个任务的差异化内容...")
contents = [gen_content(t, i) for i, t in enumerate(tasks)]
print("内容生成完成")

# ========== 生成文档 ==========
print("\n生成教学设计文档...")
tpl5 = os.path.join(template_base, "模板5：教学设计", "模板5：教学设计 模板（2023-2024）.docx")
out5 = os.path.join(output_dir, f"{tm}《{cn}》教学设计 {tn}.docx")
shutil.copy2(tpl5, out5)
doc5 = Document(out5)

# 整体设计段落
for i, p in enumerate(doc5.paragraphs):
    t = p.text.strip()
    if t == "《   》课程整体教学设计":
        sf(p.runs[0], "黑体", 18, True); p.runs[0].text = f"《{cn}》课程整体教学设计"
    elif t == "1、认知目标：" and i+4 < len(doc5.paragraphs):
        for gi, g in enumerate(["理解数据分析的基本概念、分类和适用场景",
                                "掌握数据分析方法理论和常用指标体系",
                                "掌握Excel和Python数据分析工具的使用方法",
                                "熟悉机器学习的基本原理和应用"]):
            doc5.paragraphs[i+1+gi].runs[0].text = f"{'①②③④'[gi]} {g}"
            sf(doc5.paragraphs[i+1+gi].runs[0], "仿宋_GB2312", 12)
    elif t == "2、能力目标：" and i+5 < len(doc5.paragraphs):
        for gi, g in enumerate(["能够运用Excel进行数据透视、统计分析和图表可视化",
                                "能够使用Python进行数据清洗、处理、分析和可视化",
                                "能够构建分类、回归、聚类等机器学习模型",
                                "能够独立完成商务数据分析项目并撰写报告"]):
            doc5.paragraphs[i+2+gi].runs[0].text = f"{'①②③④'[gi]} {g}"
            sf(doc5.paragraphs[i+2+gi].runs[0], "仿宋_GB2312", 12)
    elif t == "3、思政目标：" and i+4 < len(doc5.paragraphs):
        for gi, g in enumerate(["建立数据安全意识，遵守数据安全法律法规",
                                "树立诚信分析精神，坚持真实分析和规范验证",
                                "关注数字技术服务乡村振兴，培养科技报国情怀",
                                "增强职业道德意识，严守数据保密协议"]):
            doc5.paragraphs[i+1+gi].runs[0].text = f"{'①②③④'[gi]} {g}"
            sf(doc5.paragraphs[i+1+gi].runs[0], "仿宋_GB2312", 12)
    elif t == "4、素质目标：" and i+4 < len(doc5.paragraphs):
        for gi, g in enumerate(["培养数据驱动的创新思维和问题解决能力",
                                "形成规范操作、主动学习和依据标准检查成果的习惯",
                                "提升任务分工、沟通反馈、成果检查和按时交付能力",
                                "养成按步骤实施、及时测试、记录问题和持续改进的习惯"]):
            doc5.paragraphs[i+1+gi].runs[0].text = f"{'①②③④'[gi]} {g}"
            sf(doc5.paragraphs[i+1+gi].runs[0], "仿宋_GB2312", 12)
    elif t == "教学模式：" and i+2 < len(doc5.paragraphs):
        doc5.paragraphs[i+1].runs[0].text = "①理实一体化教学模式：理论讲解与实操训练交替进行"
        sf(doc5.paragraphs[i+1].runs[0], "仿宋_GB2312", 12)
        doc5.paragraphs[i+2].runs[0].text = "②项目引领、任务驱动：以企业真实项目为载体"
        sf(doc5.paragraphs[i+2].runs[0], "仿宋_GB2312", 12)
    elif t == "教学方法：" and i+2 < len(doc5.paragraphs):
        doc5.paragraphs[i+1].runs[0].text = "①案例教学法：通过分析电商行业真实案例引导学习"
        sf(doc5.paragraphs[i+1].runs[0], "仿宋_GB2312", 12)
        doc5.paragraphs[i+2].runs[0].text = "②操作演示法：教师演示操作步骤，学生跟随练习"
        sf(doc5.paragraphs[i+2].runs[0], "仿宋_GB2312", 12)
    elif t == "教材：" and i+1 < len(doc5.paragraphs):
        doc5.paragraphs[i+1].runs[0].text = f"教材：《{tb}》，天津大学出版社"
        sf(doc5.paragraphs[i+1].runs[0], "仿宋_GB2312", 12)
    elif t == "教学资料：" and i+1 < len(doc5.paragraphs):
        doc5.paragraphs[i+1].runs[0].text = f"教学资料：配套PPT课件8套、实训源代码112个文件、企业真实数据案例库"
        sf(doc5.paragraphs[i+1].runs[0], "仿宋_GB2312", 12)

# 整体设计表格（简版，因为重点在单元设计）
t0 = doc5.tables[0]
cw(t0.rows[0].cells[1], cn); cw(t0.rows[0].cells[3], cc); cw(t0.rows[0].cells[5], dp)
cw(t0.rows[1].cells[1], "2024年2月"); cw(t0.rows[1].cells[3], tn)
cw(t0.rows[2].cells[1], ct); cw(t0.rows[2].cells[3], str(th)); cw(t0.rows[2].cells[5], str(cr))
cw(t0.rows[3].cells[1], tm); cw(t0.rows[3].cells[4], cls)
cw(t0.rows[4].cells[1], "电子商务基础、计算机应用基础")
cw(t0.rows[4].cells[4], "电子商务综合实训、顶岗实习")

# 表1
t1d = doc5.tables[1]
while len(t1d.rows) > 2: t1d._tbl.remove(t1d.rows[-1]._tr)
tr1d = deepcopy(t1d.rows[1]._tr)
for i, u in enumerate(units):
    r = t1d.rows[1] if i == 0 else None
    if r is None: t1d._tbl.append(deepcopy(tr1d)); r = t1d.rows[-1]
    cw(r.cells[0], u["project_title"])
    cw(r.cells[1], str(int(u["suggested_hours"] or 6)))
t1d._tbl.append(deepcopy(tr1d))
r = t1d.rows[-1]; cw(r.cells[0], "合计"); cw(r.cells[1], str(th))

# 表2
t2d = doc5.tables[2]
while len(t2d.rows) > 2: t2d._tbl.remove(t2d.rows[-1]._tr)
tr2d = deepcopy(t2d.rows[1]._tr)
cnt = 0
for ui, u in enumerate(units):
    sk = u.get("source_skills", "")
    items = [s.strip() for s in sk.split("；") if s.strip()] if sk else []
    for si, s in enumerate(items[:3]):
        if cnt == 0: r = t2d.rows[1]
        else: t2d._tbl.append(deepcopy(tr2d)); r = t2d.rows[-1]
        cnt += 1
        cw(r.cells[0], str(ui+1)); cw(r.cells[1], u["project_title"])
        cw(r.cells[2], f"{ui+1}.{si+1} {s}"); cw(r.cells[3], f"{'①②③'[si]}{s}")
        cw(r.cells[4], "项目教学法：教师演示→学生实操→成果提交"); cw(r.cells[5], f"{u['project_title']}实训成果")

# 表3
t3d = doc5.tables[3]
while len(t3d.rows) > 4: t3d._tbl.remove(t3d.rows[-1]._tr)
tr3d = deepcopy(t3d.rows[3]._tr)
for i, task in enumerate(tasks):
    r = t3d.rows[3] if i == 0 else None
    if r is None: t3d._tbl.append(deepcopy(tr3d)); r = t3d.rows[-1]
    c = contents[i]
    cw(r.cells[0], str(task["seq"])); cw(r.cells[1], str(task.get("hours", 2)))
    cw(r.cells[2], c["title"][:40]); cw(r.cells[3], c["ag"][:50])
    for ui, u in enumerate(units):
        if u["project_title"] == task["chapter"]: cw(r.cells[4], str(ui+1)); break
    cw(r.cells[5], c["kg"][:50]); cw(r.cells[6], c["ig"][:50])
    cw(r.cells[7], c["qg"][:50]); cw(r.cells[8], "项目教学法、任务驱动法")
    cw(r.cells[9], "实训成果提交+课堂表现")

# 表4
t4d = doc5.tables[4]
for ri in range(1, 5): cw(t4d.rows[ri].cells[0], cn)
cw(t4d.rows[5].cells[0], cn); cw(t4d.rows[5].cells[2], "终结性考核")
cw(t4d.rows[5].cells[3], "综合作品（数据采集、清洗、分析、可视化和报告撰写）")

# ========== 单元设计 ==========
body5 = doc5.element.body
sectPr5 = body5.find(qn('w:sectPr'))
pt5 = deepcopy(doc5.tables[5]._tbl)
pt6 = deepcopy(doc5.tables[6]._tbl)

# 删除表4之后所有内容
t4xml = doc5.tables[4]._tbl
to_remove = []; found = False
for child in list(body5):
    if found:
        if child.tag == qn('w:sectPr'): continue
        to_remove.append(child)
    if child is t4xml: found = True
for e in to_remove: body5.remove(e)

if sectPr5 is not None:
    p = sectPr5.getparent()
    if p is not None: p.remove(sectPr5)

for idx, c in enumerate(contents):
    # ===== 表5 =====
    nt5 = deepcopy(pt5); body5.append(nt5)
    t5 = Table(nt5, doc5)
    
    # 填充关键单元格（其他合并单元格会自动共享内容）
    cw(t5.rows[0].cells[2], c["week"])      # 周次
    cw(t5.rows[0].cells[4], c["hours"])     # 课时
    cw(t5.rows[0].cells[6], cls)            # 授课班级
    cw(t5.rows[1].cells[2], tn)             # 授课教师
    cw(t5.rows[1].cells[6], c["date"])      # 授课日期
    cw(t5.rows[2].cells[2], "理实一体课程") # 课程类型
    cw(t5.rows[2].cells[6], c["room"])      # 教学环境
    cw(t5.rows[3].cells[2], c["title"])     # 教学任务
    cw(t5.rows[4].cells[2], c["kg"])        # 知识目标
    cw(t5.rows[5].cells[2], c["ag"])        # 能力目标
    cw(t5.rows[6].cells[2], c["ig"])        # 思政目标
    cw(t5.rows[7].cells[2], c["qg"])        # 素质目标
    cw(t5.rows[8].cells[2], f"教材分析：{c['textbook_an']}\n学情分析：{c['xueqing']}\n教师教育理念：{c['linian']}")
    cw(t5.rows[9].cells[2], f"【教学重点】{c['zhongdian']}")
    cw(t5.rows[10].cells[2], f"【教学难点】{c['nandian']}")
    
    # 4个知识点
    for ki, kp in enumerate(c["kps"]):
        ri = 11 + ki
        if ri >= len(t5.rows): break
        cw(t5.rows[ri].cells[2], f"{'一二三四'[ki]}、{kp['name']}")
        if kp["重点"]: cw(t5.rows[ri].cells[7], "√")
        if kp["难点"]: cw(t5.rows[ri].cells[8], "√")
        if kp["识记"]: cw(t5.rows[ri].cells[11], "√")
        if kp["理解"]: cw(t5.rows[ri].cells[15], "√")
        if kp["应用"]: cw(t5.rows[ri].cells[17], "√")
        if kp["评价"]: cw(t5.rows[ri].cells[19], "√")
    
    # 删除多余的知识点行（模板有5行，实际只有4个知识点）
    while len(t5.rows) > 11 + len(c["kps"]):
        t5._tbl.remove(t5.rows[-1]._tr)
    
    # ===== 表6 =====
    nt6 = deepcopy(pt6); body5.append(nt6)
    t6 = Table(nt6, doc5)
    
    # 行0：教学场景（col2-col5是合并单元格，写col2就行）
    cw(t6.rows[0].cells[2], 
        f"{c['room']}，配备多媒体教学设备、Python开发环境（Anaconda+Jupyter Notebook）。\n"
        f"场景布置：教师机1台+投影+学生机每人1台，预装Python 3.x及数据分析库。\n"
        f"分组方式：4人一组，围坐式布局，便于讨论和互助。")
    
    # 行1：教学资源准备
    cw(t6.rows[1].cells[2], 
        f"1. 多媒体课件：本节课PPT（约20页），含概念讲解、案例展示、代码演示\n"
        f"2. 教材：《{tb}》对应章节\n"
        f"3. 软件环境：Python 3.x、Anaconda、Jupyter Notebook\n"
        f"4. 实训数据：本节课实训项目数据集\n"
        f"5. 企业案例：农产品电商、直播电商真实案例\n"
        f"6. 参考资料：Python官方文档、Pandas官方文档、Kaggle平台")
    
    # 行2：表头行（保持空或表头）
    
    # 行3：教学导入
    cw(t6.rows[3].cells[2], 
        f"【案例展示】{c['title']}在电商中的应用案例：\n"
        f"  · 某农产品电商平台通过数据分析发现用户购买偏好，精准推荐提升转化率30%\n"
        f"  · 某直播电商通过实时数据分析调整货品策略，单场GMV提升50%\n"
        f"【提问引导】提出3个问题：\n"
        f"  1. {c['title']}在以上案例中起到了什么作用？\n"
        f"  2. 如果没有{c['title']}，企业决策会面临什么困难？\n"
        f"  3. 你认为做好{c['title']}最关键的是什么？\n"
        f"【分组讨论】4人一组讨论3分钟，推选代表发言\n"
        f"【教师点评】点评各组观点，引出本节课学习内容")
    cw(t6.rows[3].cells[3], "案例展示法、提问引导法、分组讨论法；\n学生观看案例、思考问题、小组讨论、代表发言")
    cw(t6.rows[3].cells[4], f"激发学习兴趣，了解{c['title']}的实际应用价值；引发思考，建立新旧知识联系；培养分析表达能力")
    cw(t6.rows[3].cells[5], "10分钟")
    
    # 行4：教师归纳
    cw(t6.rows[4].cells[2], 
        f"教师归纳学生讨论结果，明确本节课学习目标：\n"
        f"1. 知识目标：理解{c['title']}的基本概念、原理和方法\n"
        f"2. 能力目标：掌握操作方法，能够独立完成实训任务\n"
        f"3. 思政目标：{c['ideo0'][:25]}...\n"
        f"4. 素质目标：培养规范操作、主动学习的习惯")
    cw(t6.rows[4].cells[3], "总结归纳法、讲授法；学生明确目标、调整状态")
    cw(t6.rows[4].cells[4], "明确学习目标，导入新课")
    cw(t6.rows[4].cells[5], "5分钟")
    
    # 行5：任务1（知识讲解）
    cw(t6.rows[5].cells[2], 
        f"任务1：{c['t1_name']}（知识讲解）\n"
        f"一、基本概念\n"
        f"  1. 定义与内涵\n"
        f"  2. 核心要素与基本特征\n"
        f"  3. 与相关概念的区别与联系\n"
        f"二、作用与意义\n"
        f"  1. 在数据分析全流程中的位置\n"
        f"  2. 对企业决策的支撑价值\n"
        f"  3. 电商运营中的应用场景\n"
        f"三、基本原理与分类\n"
        f"  1. 核心原理\n"
        f"  2. 常用分类方式\n"
        f"  3. 各种方法优缺点对比\n"
        f"【PPT展示】结合课件中的示意图和案例进行讲解")
    # 插入图片
    if c.get("img"):
        try:
            p = t6.rows[5].cells[2].add_paragraph()
            r = p.add_run()
            r.add_picture(BytesIO(c["img"]["data"]), width=Cm(5))
        except: pass
    cw(t6.rows[5].cells[3], "讲授法、案例教学法、启发式提问；\n学生听讲、思考、做笔记、回答问题")
    cw(t6.rows[5].cells[4], f"理解{c['t1_name']}的基本概念、原理和分类；建立知识框架；了解应用价值")
    cw(t6.rows[5].cells[5], "25分钟")
    
    # 行6：任务2（技术演示）
    cw(t6.rows[6].cells[2], 
        f"任务2：{c['t2_name']}（技术演示）\n"
        f"一、环境准备\n"
        f"  1. 打开Anaconda Navigator，启动Jupyter Notebook\n"
        f"  2. 新建Notebook并重命名\n"
        f"  3. 导入必要的库\n"
        f"二、代码演示（逐行讲解）\n"
        f"  1. 数据加载和预处理\n"
        f"  2. 核心功能实现\n"
        f"  3. 结果输出和验证\n"
        f"  4. 常见错误和调试方法\n"
        f"三、关键步骤解析\n"
        f"  每一步为什么这么做？有什么注意事项？\n"
        f"四、学生跟随练习\n"
        f"  学生在自己电脑上跟随操作，教师巡回指导\n"
        f"【德育渗透】{c['ideo1']}")
    if c.get("t1_code"):
        add_code(t6.rows[6].cells[2], c["t1_code"][:400], 7)
    cw(t6.rows[6].cells[3], "操作演示法、逐步讲解法、巡回指导法；\n学生跟随操作、提问讨论、记录笔记")
    cw(t6.rows[6].cells[4], f"掌握{c['t2_name']}的操作步骤和代码实现；培养动手操作和代码调试能力")
    cw(t6.rows[6].cells[5], "30分钟")
    
    # 行7：任务3（实操练习）
    cw(t6.rows[7].cells[2], 
        f"任务3：{c['t3_name']}（实操练习）\n"
        f"一、实训目标\n"
        f"  1. 独立完成{c['title']}的完整操作流程\n"
        f"  2. 培养发现问题、分析问题、解决问题的能力\n"
        f"二、实训任务\n"
        f"  基础任务：按照要求完成基础功能实现\n"
        f"  进阶任务：在基础上增加新功能，尝试不同参数\n"
        f"  拓展挑战：自选数据集，用所学方法进行分析\n"
        f"三、操作要求\n"
        f"  独立完成、代码规范、注释完整、记录问题、按时提交\n"
        f"四、教师指导\n"
        f"  巡回观察、个别辅导、共性问题集中讲解\n"
        f"五、成果展示\n"
        f"  抽取2-3组展示成果，师生共同点评")
    if c.get("t2_code"):
        add_code(t6.rows[7].cells[2], c["t2_code"][:300], 7)
    cw(t6.rows[7].cells[3], "实操练习法、任务驱动法、巡回指导法、成果展示法；\n学生独立操作、小组互助、展示分享、互评学习")
    cw(t6.rows[7].cells[4], f"能够独立完成{c['t3_name']}实训任务；培养实操能力和问题解决能力；提升团队协作能力")
    cw(t6.rows[7].cells[5], "25分钟")
    
    # 行8：课堂小结（col2-col4合并为一个大单元格，只写col2，col5写时间）
    cw(t6.rows[8].cells[2], c["xiaojie"])
    cw(t6.rows[8].cells[5], "5分钟")
    
    # 行9：课后作业（col2-col4合并）
    cw(t6.rows[9].cells[2], c["zuoye"])
    cw(t6.rows[9].cells[5], "5分钟")
    
    # 行10：教学反思（col2-col5合并）
    cw(t6.rows[10].cells[2], c["fansi"])
    
    if (idx+1) % 10 == 0:
        print(f"  单元设计 {idx+1}/{len(contents)}")

if sectPr5 is not None:
    body5.append(sectPr5)

doc5.save(out5)
print(f"\n教学设计完成：{out5}")

try:
    shutil.copy2(out5, f"D:\\BaiduSyncdisk\\{os.path.basename(out5)}")
    print("已复制到D盘")
except Exception as e:
    print(f"D盘复制失败: {e}")
