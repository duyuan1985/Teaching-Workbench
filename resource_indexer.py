import hashlib
import re
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from docx import Document

import store
from legacy_ppt import extract_text as extract_legacy_ppt_text


TEXT_EXTENSIONS = {".py", ".html", ".htm", ".css", ".js", ".json", ".xml", ".csv", ".txt", ".md"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
PROJECT_EXTENSIONS = {".psd", ".psb", ".ai", ".xd", ".fig"}


def digest(path):
    hasher = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            hasher.update(chunk)
    return hasher.hexdigest()


def compact(text, limit=8000):
    text = re.sub(r"\s+", " ", text or "").strip()
    return text[:limit]


def pptx_excerpt(path):
    from resource_analyzer import _detect_ppt_noise_texts, _is_noise_text
    presentation = Presentation(path)
    noise_texts = _detect_ppt_noise_texts(presentation)
    slides = []
    for index, slide in enumerate(presentation.slides, 1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if not _is_noise_text(text, noise_texts):
                    parts.append(text)
        if parts:
            slides.append(f"第{index}页：{'；'.join(parts)}")
        if len(" ".join(slides)) >= 8000:
            break
    return compact(" ".join(slides))


def legacy_ppt_text(path):
    return extract_legacy_ppt_text(path)


def legacy_ppt_excerpt(path):
    return compact(legacy_ppt_text(path))


def docx_text(path):
    """提取 .docx 文件全文（段落 + 表格），跳过目录，标注章节"""
    from resource_analyzer import _is_toc_paragraph, _is_heading_paragraph, _heading_level, classify_document_type
    doc = Document(path)
    parts = []
    current_section = None
    in_toc = False
    toc_ended = False
    para_index = 0

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        para_index += 1

        # 目录检测与跳过
        if not toc_ended:
            if _is_toc_paragraph(text):
                in_toc = True
                continue
            if in_toc:
                if _heading_level(paragraph) == 1 or para_index > 50:
                    toc_ended = True
                    in_toc = False
                else:
                    continue

        # 标题识别
        if _is_heading_paragraph(paragraph):
            level = _heading_level(paragraph)
            if level == 1:
                current_section = text
            parts.append(f"【{text}】")
        else:
            if current_section:
                parts.append(f"[{current_section}] {text}")
            else:
                parts.append(text)

    # 表格
    for table in doc.tables:
        for row in table.rows:
            row_texts = []
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    row_texts.append(t)
            if row_texts:
                parts.append(" | ".join(row_texts))

    return "\n".join(parts)


def docx_excerpt(path):
    return compact(docx_text(path))


def legacy_doc_text(path):
    """用 Word COM 提取 .doc 文件全文（通过 PowerShell 脚本）"""
    script = Path(__file__).parent / "extract_legacy_training_doc.ps1"
    if not script.exists():
        return ""
    with tempfile.TemporaryDirectory(prefix="doc-extract-") as temp_dir:
        output = Path(temp_dir) / "content.txt"
        try:
            result = subprocess.run(
                ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass",
                 "-File", str(script), "-InputPath", str(path), "-OutputPath", str(output)],
                capture_output=True, text=True, encoding="utf-8", timeout=60,
            )
            if result.returncode == 0 and output.exists():
                return output.read_text(encoding="utf-8-sig", errors="ignore")
        except (subprocess.TimeoutExpired, OSError):
            pass
    return ""


def legacy_doc_excerpt(path):
    return compact(legacy_doc_text(path))


def text_excerpt(path):
    for encoding in ("utf-8", "gb18030", "utf-8-sig"):
        try:
            return compact(path.read_text(encoding=encoding, errors="strict"))
        except (UnicodeDecodeError, OSError):
            continue
    return ""


def classify(path):
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return "PPT课件"
    if suffix == ".ppt":
        return "旧版PPT课件"
    if suffix in TEXT_EXTENSIONS:
        return "实训源码/数据"
    if suffix in IMAGE_EXTENSIONS:
        return "图片素材"
    if suffix in PROJECT_EXTENSIONS:
        return "项目源文件"
    if suffix in {".doc", ".docx", ".pdf"}:
        # 更细粒度的文档分类
        from resource_analyzer import classify_document_type
        doc_type, _ = classify_document_type(path)
        if doc_type in ("教学大纲", "课程标准"):
            return "教学大纲"
        elif doc_type in ("教案", "教学设计"):
            return "教案"
        elif doc_type in ("实训指导书", "实验指导"):
            return "实训指导书"
        elif doc_type in ("实训大纲", "实训方案"):
            return "实训方案"
        elif doc_type in ("习题答案", "参考答案"):
            return "习题答案"
        elif doc_type == "考核方案":
            return "考核方案"
        else:
            return "教材/实训文档"
    if suffix in {".mp4", ".avi", ".mov", ".wmv", ".mp3", ".wav"}:
        return "音视频资源"
    return "其他资源"


def build_resource_index(offering):
    sources = store.rows(
        "SELECT * FROM source_files WHERE offering_id=? AND source_type IN ('教材目录','PPT目录','实训资源目录')",
        (offering["id"],),
    )
    roots = [Path(source["source_path"]) for source in sources]
    if not roots:
        raise ValueError("尚未登记教材目录、PPT目录或实训资源目录。")
    items = []
    for root in roots:
        if not root.exists():
            continue
        paths = [root] if root.is_file() else root.rglob("*")
        for path in paths:
            if not path.is_file() or path.name.startswith("~$"):
                continue
            resource_type = classify(path)
            excerpt = ""
            try:
                if path.suffix.lower() == ".pptx":
                    excerpt = pptx_excerpt(path)
                elif path.suffix.lower() == ".ppt":
                    excerpt = legacy_ppt_excerpt(path)
                elif path.suffix.lower() == ".docx":
                    excerpt = docx_excerpt(path)
                elif path.suffix.lower() == ".doc":
                    excerpt = legacy_doc_excerpt(path)
                elif path.suffix.lower() in TEXT_EXTENSIONS and path.stat().st_size <= 2_000_000:
                    excerpt = text_excerpt(path)
            except Exception:
                excerpt = ""
            items.append((resource_type, str(path), path.stem, excerpt, digest(path)))
    with store.connect() as db:
        db.execute("DELETE FROM resource_items WHERE offering_id=?", (offering["id"],))
        db.executemany(
            """INSERT INTO resource_items
            (offering_id,resource_type,file_path,title,content_excerpt,source_hash)
            VALUES (?,?,?,?,?,?)""",
            [(offering["id"], *item) for item in items],
        )
        db.commit()
    counts = {}
    for resource_type, *_ in items:
        counts[resource_type] = counts.get(resource_type, 0) + 1
    return len(items), counts
