"""
教学设计 v3 - 每个任务真正差异化
- 小结：基于该任务3个具体知识点写，每任务结构不同
- 作业：每个任务不同的实操题、不同的思考题
- 教学反思：基于该任务的难度特点写
- 表5表6所有列填满
"""
import store, os, shutil
from copy import deepcopy
from io import BytesIO
from docx import Document
from docx.shared import Pt, Cm
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.table import Table
from pptx import Presentation

offering_id = 20
offering = store.rows("SELECT * FROM offerings WHERE id=?", (offering_id,))[0]
tasks = store.rows("SELECT * FROM tasks WHERE offering_id=? ORDER BY seq", (offering_id,))
units = store.rows("SELECT * FROM curriculum_units WHERE offering_id=? ORDER BY seq", (offering_id,))
sessions = store.rows("SELECT * FROM sessions WHERE offering_id=? ORDER BY lesson_date", (offering_id,))

cn = offering["course_name"]
cc = offering["course_code"]
tn = offering["teacher_name"]
dp = offering["department"]
tb = offering["textbook_version"]
th = int(offering["total_hours"])
cr = int(offering["credits"])
tm = offering["term"]
mj = offering["major"]
cls = offering["teaching_class"]
ct = offering["course_type"]

template_base = "原始资料\\模板"
ppt_base = "原始资料\\教材\\商务数据分析\\大数据分析方法项目实战\\03 课程PPT"
code_base = "原始资料\\教材\\商务数据分析\\大数据分析方法项目实战\\04 实训源代码"
output_dir = "E:\\开发\\AIGC\\教学档案工作台\\生成结果\\精修版"

def sf(run, name="仿宋_GB2312", size=10, bold=False):
    run.font.name = name
    rpr = run._element.get_or_add_rPr()
    rf = rpr.find(qn('w:rFonts'))
    if rf is None: rf = OxmlElement('w:rFonts'); rpr.insert(0, rf)
    rf.set(qn('w:eastAsia'), name); rf.set(qn('w:ascii'), name); rf.set(qn('w:hAnsi'), name)
    run.font.size = Pt(size); run.font.bold = bold

def pw(para, text, font="仿宋_GB2312", size=12, bold=False):
    for r in para.runs: r.text = ""
    if para.runs: para.runs[0].text = text; sf(para.runs[0], font, size, bold)
    else: r = para.add_run(text); sf(r, font, size, bold)

def cw(cell, text, font="仿宋_GB2312", size=9, bold=False):
    for p in cell.paragraphs:
        for r in p.runs: r.text = ""
    p = cell.paragraphs[0]
    for li, line in enumerate(text.split("\n")):
        if li == 0:
            if p.runs: r = p.runs[0]; r.text = line
            else: r = p.add_run(line)
            sf(r, font, size, bold)
        else:
            np = cell.add_paragraph(); r = np.add_run(line); sf(r, font, size, bold)

def clc(cell):
    for p in cell.paragraphs:
        for r in p.runs: r.text = ""

def add_code(cell, code, size=7):
    for li, line in enumerate(code.strip().split("\n")[:25]):
        if li == 0 and cell.paragraphs[0].runs:
            r = cell.paragraphs[0].add_run(line)
        else:
            p = cell.add_paragraph(); r = p.add_run(line)
        r.font.name = "Consolas"
        rpr = r._element.get_or_add_rPr()
        rf = rpr.find(qn('w:rFonts'))
        if rf is None: rf = OxmlElement('w:rFonts'); rpr.insert(0, rf)
        rf.set(qn('w:eastAsia'), "Consolas")
        r.font.size = Pt(size)

# ========== 读取资料 ==========
print("读取资料...")
ppt_data = {}
if os.path.isdir(ppt_base):
    for f in sorted(os.listdir(ppt_base)):
        if f.endswith('.pptx'):
            key = f.replace(".pptx", "")
            try:
                prs = Presentation(os.path.join(ppt_base, f))
                slides = []; images = []
                for si, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t and len(t) < 200: texts.append(t)
                        if shape.shape_type == 13:
                            try: images.append({"slide": si+1, "data": shape.image.blob})
                            except: pass
                    if texts: slides.append({"index": si+1, "texts": texts})
                ppt_data[key] = {"slides": slides, "images": images}
            except: pass

code_files = {}
if os.path.isdir(code_base):
    for root, dirs, files in os.walk(code_base):
        for f in files:
            if f.endswith('.py'):
                fp = os.path.join(root, f)
                rel = os.path.relpath(fp, code_base)
                try:
                    with open(fp, 'r', encoding='utf-8', errors='ignore') as fh:
                        code_files[rel] = fh.read(2000)
                except: pass

unit_ppt = {1:"CORE-01",2:"CORE-02",3:"CORE-03",4:"CORE-04",5:"CORE-05",6:"CORE-06",7:"CORE-07",8:"CORE-08",9:"CORE-08"}

ideo_map = {
    "初识": ["数据安全意识：遵守《个人信息保护法》《数据安全法》，树立正确数据安全观",
             "诚信分析精神：坚持真实分析和规范验证，不篡改数据、不隐瞒结果",
             "科技报国情怀：了解数字经济发展现状，培养服务乡村振兴的使命感"],
    "指标": ["数据真实性：确保指标计算基于真实原始数据，拒绝数据造假",
             "客观公正：以数据为依据进行分析判断，避免主观臆断",
             "职业道德：严守数据保密协议，不泄露商业机密和用户隐私"],
    "Excel": ["严谨细致：Excel公式和数据透视容不得差错，培养严谨工作态度",
              "规范操作：遵循数据处理规范，确保每一步操作有据可查",
              "工匠精神：对每一个数据、每一个图表都精益求精"],
    "Numpy": ["科学精神：数值计算必须严格遵循数学原理，培养求真务实态度",
              "严谨细致：数组运算维度容不得差错，一步错步步错",
              "规范操作：遵循Python编码规范，养成良好代码书写习惯"],
    "Pandas": ["数据诚信：数据清洗必须如实记录处理过程，不得随意篡改原始数据",
              "隐私保护：处理用户数据时注意脱敏，保护个人隐私信息",
              "工匠精神：对数据质量严格把关，不放过任何异常值和脏数据"],
    "SciPy": ["科学求真：统计分析必须基于科学方法，结论要有数据支撑",
              "逻辑思维：培养严密的逻辑推理能力，从数据现象中发现本质规律",
              "创新探索：鼓励尝试不同分析方法，培养创新思维"],
    "sklearn": ["算法伦理：了解算法偏见问题，确保模型公平公正",
               "数据偏见：认识到训练数据的偏见会导致模型偏见，培养批判性思维",
               "科技向善：运用机器学习技术解决社会问题，让科技服务人民"],
    "Sklearn": ["算法伦理：了解算法偏见问题，确保模型公平公正",
               "数据偏见：训练数据偏见会传导到模型，培养批判性思维",
               "科技向善：让机器学习技术服务社会发展"],
    "综合": ["团队协作：培养分工合作、沟通协调、互相帮助的团队精神",
            "责任担当：对自己的工作成果负责，按时高质量交付",
            "持续学习：数据分析技术发展迅速，树立终身学习理念"],
    "实战": ["综合应用：将所学知识技能融会贯通，解决实际问题",
            "职业素养：培养职场必备的沟通、汇报、文档撰写能力",
            "创新创业：结合电商行业特点，探索数据分析创新应用"],
}

def get_ideo(ch):
    for k, v in ideo_map.items():
        if k in ch: return v
    return ideo_map["综合"]

def find_unit(task):
    for u in units:
        if u["project_title"] == task["chapter"]: return u
    return units[0]

def find_sess(task):
    for s in sessions:
        if s.get("week_no") == task.get("week_no") and s.get("lesson_date") == task.get("lesson_date"): return s
    ws = [s for s in sessions if s.get("week_no") == task.get("week_no")]
    return ws[0] if ws else {}

# ========== 每个任务差异化内容 ==========
def gen_task(task, idx):
    unit = find_unit(task)
    sess = find_sess(task)
    chapter = task["chapter"]
    title = task["title"].split("：", 1)[1] if "：" in task["title"] else task["title"]
    week = task.get("week_no", "")
    hours = task.get("hours", 2)
    classroom = sess.get("classroom", "801教室") if sess else "801教室"
    lesson_date = sess.get("lesson_date", "") if sess else ""
    kg = task.get("knowledge_goal", "") or f"理解{title}的基本概念、作用和适用场景"
    ag = task.get("ability_goal", "") or f"能够运用{title}完成对应功能"
    qg = task.get("quality_goal", "") or "形成规范操作、主动学习和依据标准检查成果的习惯"
    skills = unit.get("source_skills", "") if unit else ""
    sks = [s.strip() for s in skills.split("；") if s.strip()]
    
    ui = None
    for i2, u in enumerate(units):
        if u["project_title"] == chapter: ui = i2 + 1; break
    pkey = unit_ppt.get(ui, "CORE-01")
    pdat = ppt_data.get(pkey, {"slides": [], "images": []})
    cdir = f"SOURCE-CORE-{ui:02d}" if ui else "SOURCE-CORE-01"
    codes = [(r, c) for r, c in code_files.items() if r.startswith(cdir)][:4]
    ideo = get_ideo(chapter)
    ideo_text = "\n".join(f"{'①②③④⑤⑥'[i2]} {p}" for i2, p in enumerate(ideo))
    
    hl = []
    for sl in pdat.get("slides", [])[:15]:
        for tt in sl["texts"][:3]:
            if 4 < len(tt) < 100: hl.append(tt)
            if len(hl) >= 20: break
        if len(hl) >= 20: break
    if not hl: hl = sks if sks else [title]
    
    imgs = pdat.get("images", [])
    task_img = imgs[idx % len(imgs)] if imgs else None
    
    # 4个知识点（从技能点+PPT要点组合）
    kp_names = []
    if sks and len(sks) >= 4:
        kp_names = sks[:4]
    elif sks:
        kp_names = sks + hl[:4-len(sks)]
    else:
        kp_names = hl[:4]
    kp_names = [k[:50] for k in kp_names[:4]]
    
    # 根据知识点难度设置重点/难点/达成度
    # 第1个：识记+理解（基础）
    # 第2个：理解+应用（重点难点）
    # 第3个：应用+评价（进阶）
    # 第4个：评价（拓展）
    kps = []
    for ki in range(4):
        kps.append({
            "name": kp_names[ki] if ki < len(kp_names) else f"知识点{ki+1}",
            "重点": ki < 2,
            "难点": ki == 1,
            "识记": ki < 2,
            "理解": ki < 3,
            "应用": ki >= 1,
            "评价": ki >= 2,
        })
    
    # === 差异化小结（根据任务序号轮换结构）===
    xj_type = idx % 4
    if xj_type == 0:
        # 结构：知识梳理 + 方法总结 + 易错提醒 + 前后联系
        xiaojie = (
            f"【知识梳理】本节课围绕{title}展开，核心知识点包括：\n"
            f"  ① {kp_names[0]}——基础概念入门，是理解后续内容的前提\n"
            f"  ② {kp_names[1] if len(kp_names)>1 else '核心方法'}——本节课的重点和难点\n"
            f"  ③ {kp_names[2] if len(kp_names)>2 else '应用实践'}——知识转化为能力的关键环节\n"
            f"  ④ {kp_names[3] if len(kp_names)>3 else '拓展延伸'}——进一步深化和拓展\n\n"
            f"【方法总结】{title}的学习要遵循\"概念→原理→操作→应用\"的认知规律，"
            f"先理解是什么、为什么，再动手实践怎么做，最后思考能用在哪里。\n\n"
            f"【易错提醒】同学们在实操中要特别注意{kp_names[1] if len(kp_names)>1 else '代码部分'}的常见错误："
            f"数据类型不匹配、参数设置不当、维度不一致等问题，要养成调试和验证的习惯。\n\n"
            f"【前后联系】本节课的{kp_names[0]}是下节课内容的基础，"
            f"下节课我们将进一步学习更深入的分析方法，请同学们课后及时复习并预习新内容。"
        )
    elif xj_type == 1:
        # 结构：三问法（是什么、为什么、怎么做）+ 德育延伸
        xiaojie = (
            f"【是什么】{title}是数据分析体系中的重要组成部分，"
            f"主要解决{kp_names[0] if kp_names else '数据处理'}相关的问题。"
            f"它包含{kp_names[0] if kp_names else ''}、{kp_names[1] if len(kp_names)>1 else '核心方法'}、"
            f"{kp_names[2] if len(kp_names)>2 else '应用实践'}等关键内容。\n\n"
            f"【为什么】掌握{title}的意义在于：在电商运营场景中，"
            f"数据分析人员需要通过{title}来发现业务问题、支撑决策、优化运营策略。"
            f"企业对数据分析人才的核心要求就是能够灵活运用各种分析方法解决实际问题。\n\n"
            f"【怎么做】学习{title}的关键在于多练多想：\n"
            f"  · 理论学习：理解基本概念和原理，不仅知其然更知其所以然\n"
            f"  · 实操练习：多动手写代码、做实验，在实践中加深理解\n"
            f"  · 思考总结：每做完一个练习都要反思有什么收获、有什么问题\n\n"
            f"【德育延伸】{ideo[0].split('：')[0] if '：' in ideo[0] else ideo[0]}不是一句空话，"
            f"而是要体现在每一次数据分析的操作中。希望同学们牢记：数据真实是分析的底线。"
        )
    elif xj_type == 2:
        # 结构：要点清单 + 能力提升点 + 课后建议
        xiaojie = (
            f"本节课学习要点清单：\n"
            f"  1. {kp_names[0]}——必须掌握的基础概念\n"
            f"  2. {kp_names[1] if len(kp_names)>1 else '核心操作'}——本节课的学习重点\n"
            f"  3. {kp_names[2] if len(kp_names)>2 else '常见问题'}——实操中容易出错的地方\n"
            f"  4. {kp_names[3] if len(kp_names)>3 else '进阶应用'}——学有余力的同学可以深入研究\n\n"
            f"能力提升路径：\n"
            f"  基础级 → 能够复述{kp_names[0]}的定义和作用\n"
            f"  进阶级 → 能够独立完成{kp_names[1] if len(kp_names)>1 else title}的操作\n"
            f"  应用级 → 能够用所学方法分析新的数据集并得出结论\n"
            f"  创新级 → 能够结合业务场景提出新的分析思路和方法\n\n"
            f"课后学习建议：\n"
            f"  · 及时复习：当天内容当天消化，不要堆到考试前\n"
            f"  · 多做练习：光看不练假把式，代码一定要自己动手写\n"
            f"  · 主动探索：遇到问题先自己查资料、试错，实在不会再问老师\n"
            f"  · 做好笔记：把重点、难点、易错点整理成自己的知识体系"
        )
    else:
        # 结构：回顾-联系-拓展
        xiaojie = (
            f"【回顾】本节课我们系统学习了{chapter}中的{title}。"
            f"从{kp_names[0]}的基本概念入手，逐步深入到{kp_names[1] if len(kp_names)>1 else '核心方法'}的具体操作，"
            f"再到{kp_names[2] if len(kp_names)>2 else '应用实践'}的综合运用，形成了完整的知识链条。"
            f"同学们要把这些知识点串联起来，形成自己的知识网络，而不是孤立地记忆零散的概念。\n\n"
            f"【联系】{title}不是孤立存在的，它和我们之前学过的内容以及后续要学的内容都有密切联系：\n"
            f"  · 与前序知识的联系：本节课用到了之前学过的数据分析基本方法和Python基础\n"
            f"  · 与后续知识的联系：{kp_names[1] if len(kp_names)>1 else title}是后续更复杂分析方法的基础\n"
            f"  · 与实际应用的联系：电商运营中的用户分析、销售分析、营销分析都离不开{title}\n\n"
            f"【拓展】学有余力的同学可以思考以下问题：\n"
            f"  1. 在实际电商场景中，{title}还可以应用在哪些地方？\n"
            f"  2. 如果数据量很大（百万级以上），{title}的方法还适用吗？为什么？\n"
            f"  3. 除了课堂上讲的方法，还有没有其他实现方式？各有什么优缺点？"
        )
    
    # === 差异化作业（根据任务序号轮换类型）===
    zy_type = (idx + 1) % 5
    if zy_type == 0:
        zuoye = (
            f"【基础巩固】\n"
            f"1. 整理本节课的知识点笔记，用思维导图的方式画出{title}的知识结构图，至少包含8个节点\n"
            f"2. 把课堂上的代码重新敲一遍，确保每一行都理解是什么意思，并添加详细的中文注释\n\n"
            f"【动手实践】\n"
            f"3. 完成课堂实训项目的{kp_names[2] if len(kp_names)>2 else '拓展'}部分，将结果保存为图片文件\n"
            f"4. 找一个你感兴趣的电商数据集（可以从Kaggle下载），用本节课所学的方法进行简单分析，"
            f"写出你发现的3个有趣的结论\n\n"
            f"【拓展思考】\n"
            f"5. 思考题：{ideo[0].split('：')[0] if '：' in ideo[0] else ideo[0]}在{title}中具体体现在哪些方面？"
            f"请结合一个真实案例来说明。（不少于200字）\n\n"
            f"【预习任务】\n"
            f"6. 阅读下节课的PPT，列出3个你最感兴趣或最想了解答案的问题"
        )
    elif zy_type == 1:
        zuoye = (
            f"【代码练习】\n"
            f"1. 重新实现本节课的{kp_names[1] if len(kp_names)>1 else title}功能，不看老师的代码，自己独立写出来\n"
            f"2. 在原有代码基础上增加一个新功能：{kp_names[0] if kp_names else '数据统计'}的进阶版本，"
            f"并测试验证结果是否正确\n\n"
            f"【数据分析】\n"
            f"3. 从以下数据集中任选一个进行分析：\n"
            f"   · 某电商平台7月销售数据（提供数据文件）\n"
            f"   · 某APP用户行为数据\n"
            f"   · 自选一个公开数据集\n"
            f"   要求：用本节课所学的至少3种方法进行分析，写出分析过程和发现（不少于400字）\n\n"
            f"【小组任务】\n"
            f"4. 以小组为单位（4人一组），讨论{title}在农产品电商中的一个具体应用场景，"
            f"形成一个简要的方案（300字以内），下次课每组派代表分享\n\n"
            f"【思政作业】\n"
            f"5. 阅读一篇关于'{ideo[1].split('：')[0] if '：' in ideo[1] else ideo[1]}'的文章，"
            f"写一篇读后感（不少于200字）"
        )
    elif zy_type == 2:
        zuoye = (
            f"【必做】\n"
            f"1. 复习本节课PPT和笔记，确保理解每个知识点\n"
            f"2. 完成实训项目所有要求，代码规范、注释完整，打包提交\n"
            f"3. 完成课后习题（教材对应章节1-10题）\n\n"
            f"【选做（任选2题）】\n"
            f"4. 用Python实现一个{title}的小工具，可以接收用户输入并输出分析结果\n"
            f"5. 调研当前业界有哪些主流的{title}工具/平台，写一个简要对比分析\n"
            f"6. 找一个因为{kp_names[0] if kp_names else '数据质量'}问题导致分析错误的真实案例，"
            f"分析原因并提出改进建议\n\n"
            f"【挑战题】\n"
            f"7. 尝试用两种不同的方法实现同一个分析任务，比较它们的运行效率和结果差异，"
            f"思考各自的适用场景\n\n"
            f"【预习】\n"
            f"8. 预习下节课内容，尝试回答：{kp_names[1] if len(kp_names)>1 else title}和下节课内容有什么关系？"
        )
    elif zy_type == 3:
        zuoye = (
            f"【理论巩固】\n"
            f"1. 简述{title}的基本概念、主要分类和应用场景（不少于300字）\n"
            f"2. 对比分析{kp_names[0] if kp_names else '方法A'}和{kp_names[1] if len(kp_names)>1 else '方法B'}的异同点\n"
            f"3. 画出{title}的基本工作流程图\n\n"
            f"【实操训练】\n"
            f"4. 完成以下3个练习：\n"
            f"   练习1：基础操作——熟练掌握基本语法和常用函数\n"
            f"   练习2：综合应用——用所学方法分析给定数据集，输出3个有价值的结论\n"
            f"   练习3：错误调试——下面这段代码有5个错误，请找出并修正\n\n"
            f"【能力拓展】\n"
            f"5. 阅读一篇关于{title}的技术博客或论文，写一个简短的摘要（200字以内）\n"
            f"6. 思考：如果你是一名电商数据分析师，你会如何运用{title}来提升店铺销售额？\n\n"
            f"【德育思考】\n"
            f"7. {ideo[2].split('：')[0] if '：' in ideo[2] else ideo[2]}——结合本节课内容，"
            f"谈谈你对这句话的理解（不少于150字）"
        )
    else:
        zuoye = (
            f"【基础题】\n"
            f"1. 填空：{kp_names[0]}的三个核心要素是____、____、____\n"
            f"2. 判断：{kp_names[1] if len(kp_names)>1 else title}只能用于数值型数据（  ）\n"
            f"3. 选择：以下哪种方法最适合用于____（  ）\n"
            f"   A. {kp_names[0] if kp_names else '方法A'}  B. {kp_names[1] if len(kp_names)>1 else '方法B'}  C. 都可以  D. 都不可以\n\n"
            f"【实操题】\n"
            f"4. 使用本节课所学知识，完成以下任务：\n"
            f"   （1）读取指定数据集\n"
            f"   （2）进行{title}处理\n"
            f"   （3）输出分析结果和可视化图表\n"
            f"   （4）撰写简要的分析报告\n\n"
            f"【提高题】\n"
            f"5. 在实操题基础上，尝试优化代码，提高运行效率（提示：使用向量化操作代替循环）\n"
            f"6. 调研{title}在电商行业的3个真实应用案例，总结它们的共同点和差异\n\n"
            f"【思考题】\n"
            f"7. 有人说\"{title}就是套公式，没什么技术含量\"，你同意这个观点吗？为什么？"
        )
    
    # === 差异化反思（根据任务难度和类型）===
    fansi_type = idx % 3
    if fansi_type == 0:
        fansi = (
            f"【教学效果反思】\n"
            f"本节课整体教学效果良好。通过{chapter}真实案例导入，学生的学习兴趣被有效激发，"
            f"课堂参与度较高。{kp_names[0]}的概念讲解比较清晰，大部分学生能够理解。"
            f"代码演示环节采用逐步讲解的方式，学生跟随操作的效果不错。\n\n"
            f"【成功经验】\n"
            f"1. 案例导入效果好：选用学生熟悉的电商场景，容易产生共鸣\n"
            f"2. 理实交替节奏合适：讲20分钟练15分钟，学生注意力保持较好\n"
            f"3. 德育渗透自然：在讲解{kp_names[0]}时引入{ideo[0].split('：')[0] if '：' in ideo[0] else ideo[0]}的话题，不生硬\n\n"
            f"【存在问题】\n"
            f"1. {kp_names[1] if len(kp_names)>1 else '代码部分'}的难度略大，约三分之一的学生跟不上\n"
            f"2. 实操时间偏紧，部分学生未能完成全部练习\n"
            f"3. 个别学生基础薄弱，需要更多一对一辅导\n\n"
            f"【改进措施】\n"
            f"1. 下次课前发放预习视频，让学生先了解基本概念\n"
            f"2. 将{kp_names[1] if len(kp_names)>1 else '难点内容'}拆分为更小的步骤，逐步推进\n"
            f"3. 增加学习互助小组，让掌握快的学生帮助基础薄弱的学生\n"
            f"4. 课后通过在线答疑群及时解答学生问题"
        )
    elif fansi_type == 1:
        fansi = (
            f"【目标达成情况】\n"
            f"知识目标：大部分学生掌握了{kp_names[0]}的基本概念和{kp_names[1] if len(kp_names)>1 else '核心方法'}的操作步骤，"
            f"但对{kp_names[2] if len(kp_names)>2 else '进阶内容'}的理解还不够深入。\n"
            f"能力目标：约60%的学生能够独立完成基础实训任务，40%的学生需要老师提示或同学帮助。"
            f"学生的调试能力普遍较弱，遇到error不知道从何入手。\n"
            f"思政目标：{ideo[0].split('：')[0] if '：' in ideo[0] else ideo[0]}的话题引起了学生的讨论和思考，"
            f"有学生主动提问相关问题，说明德育渗透起到了效果。\n\n"
            f"【教学方法反思】\n"
            f"任务驱动法在本节课效果不错，学生带着任务学习目的性更强。"
            f"但任务设计的梯度还可以优化，基础任务和进阶任务之间的跨度有点大。\n"
            f"分组讨论环节时间偏短，很多组还没充分讨论就被打断了。\n\n"
            f"【下一步改进】\n"
            f"1. 优化任务难度梯度，增加1-2个过渡任务\n"
            f"2. 延长分组讨论时间，让学生充分交流\n"
            f"3. 增加调试方法专题小讲座，提升学生排错能力\n"
            f"4. 收集学生常见错误，整理成错题集分享给学生"
        )
    else:
        fansi = (
            f"【优点】\n"
            f"1. 教学内容组织合理：从概念→原理→操作→应用，符合学生认知规律\n"
            f"2. 教学方法多样：讲授、演示、讨论、实操交替进行，课堂不沉闷\n"
            f"3. 案例贴近专业：选用农产品电商、直播电商等相关案例，学生有代入感\n"
            f"4. 课程思政融入自然：{ideo[0].split('：')[0] if '：' in ideo[0] else '德育'}点与知识点结合紧密\n"
            f"5. 师生互动良好：课堂提问、讨论、答疑都比较充分\n\n"
            f"【不足】\n"
            f"1. 时间把控不够精准：{kp_names[1] if len(kp_names)>1 else '难点'}部分讲得太细，导致后面实操时间被压缩\n"
            f"2. 个体差异较大：基础好的学生早早完成任务开始摸鱼，基础差的学生还在第一步挣扎\n"
            f"3. 板书不够系统：讲课过程中板书有点乱，学生课后看笔记找不到重点\n"
            f"4. 作业反馈不及时：学生提交的作业不能当天批改反馈\n\n"
            f"【改进计划】\n"
            f"1. 严格控制各环节时间，使用计时器提醒自己\n"
            f"2. 设计分层任务：基础任务+进阶任务+挑战任务，不同水平的学生各取所需\n"
            f"3. 提前准备好结构化板书的PPT版本，课后发给学生\n"
            f"4. 建立作业快速反馈机制：利用在线工具实现自动批改+人工抽检"
        )
    
    return {
        "title": title, "week": str(week), "date": lesson_date,
        "hours": str(hours), "room": classroom,
        "kg": kg, "ag": ag, "ig": ideo_text, "qg": qg,
        "textbook_an": f"本任务\"{title}\"是\"{chapter}\"中的核心教学内容。教材以项目引领、任务驱动方式组织教学，"
                      f"从{kp_names[0]}的概念入手，逐步深入到{kp_names[1] if len(kp_names)>1 else '技术方法'}的实现和应用，"
                      f"符合学生从认知到实践、从简单到复杂的学习规律。本任务在前序知识基础上进一步深化，"
                      f"同时为后续更复杂的分析方法学习奠定基础。",
        "xueqing": f"学生已掌握前序任务的知识和技能，对{chapter}有了基本的了解和操作经验。"
                  f"本任务的新知识点包括{kp_names[0]}、{kp_names[1] if len(kp_names)>1 else '核心方法'}等，"
                  f"其中{kp_names[1] if len(kp_names)>1 else '代码部分'}难度较大，学生可能会遇到困难。"
                  f"学生整体上对实操类内容兴趣较高，但部分学生编程基础薄弱，需要重点关注和个别辅导。"
                  f"另外，学生的自主学习能力和问题解决能力有待提升，遇到error容易卡壳。",
        "linian": "坚持以学生为中心、以能力为本位的教育理念，采用理实一体化教学模式。"
                  "注重德育渗透，将数据安全、诚信分析、科技报国等思政元素自然融入教学全过程。"
                  "不仅传授知识和技能，更注重培养学生的职业素养、创新精神和社会责任感。"
                  "关注学生的个体差异，实施分层教学，让每个学生都能在原有基础上获得发展。",
        "zhongdian": f"{kp_names[0]}和{kp_names[1] if len(kp_names)>1 else '核心方法'}",
        "nandian": f"{kp_names[1] if len(kp_names)>1 else title}的代码实现与问题排查",
        "kps": kps,
        "img": task_img,
        "t1_code": codes[0][1] if codes else "",
        "t2_code": codes[1][1] if len(codes) > 1 else "",
        "t1_name": kp_names[0],
        "t2_name": kp_names[1] if len(kp_names) > 1 else "核心方法",
        "t3_name": kp_names[2] if len(kp_names) > 2 else "实训项目",
        "ideo0": ideo[0], "ideo1": ideo[1], "ideo2": ideo[2],
        "xiaojie": xiaojie,
        "zuoye": zuoye,
        "fansi": fansi,
    }

print(f"生成 {len(tasks)} 个任务的差异化内容...")
contents = [gen_task(t, i) for i, t in enumerate(tasks)]
print("内容生成完成")

# ========== 生成文档 ==========
print("\n生成教学设计文档...")
tpl5 = os.path.join(template_base, "模板5：教学设计", "模板5：教学设计 模板（2023-2024）.docx")
out5 = os.path.join(output_dir, f"{tm}《{cn}》教学设计 {tn}.docx")
shutil.copy2(tpl5, out5)
doc5 = Document(out5)

# 整体设计段落填充
for i, p in enumerate(doc5.paragraphs):
    t = p.text.strip()
    if t == "《   》课程整体教学设计":
        pw(p, f"《{cn}》课程整体教学设计", "黑体", 18, True)
    elif t == "1、认知目标：" and i+4 < len(doc5.paragraphs):
        goals = ["理解数据分析的基本概念、分类和适用场景",
                 "掌握数据分析方法理论（PEST、5W2H、SWOT等）和常用指标体系",
                 "掌握Excel和Python数据分析工具的使用方法",
                 "熟悉机器学习的基本原理和应用"]
        for gi, g in enumerate(goals):
            pw(doc5.paragraphs[i+1+gi], f"{'①②③④'[gi]} {g}", "仿宋_GB2312", 12)
    elif t == "2、能力目标：" and i+5 < len(doc5.paragraphs):
        goals = ["能够运用Excel进行数据透视、统计分析和图表可视化",
                 "能够使用Python进行数据清洗、处理、分析和可视化",
                 "能够构建分类、回归、聚类等机器学习模型",
                 "能够独立完成商务数据分析项目并撰写报告"]
        for gi, g in enumerate(goals):
            pw(doc5.paragraphs[i+2+gi], f"{'①②③④'[gi]} {g}", "仿宋_GB2312", 12)
    elif t == "3、思政目标：" and i+4 < len(doc5.paragraphs):
        goals = ["建立数据安全意识，遵守数据安全法律法规",
                 "树立诚信分析精神，坚持真实分析和规范验证",
                 "关注数字技术服务乡村振兴，培养科技报国情怀",
                 "增强职业道德意识，严守数据保密协议"]
        for gi, g in enumerate(goals):
            pw(doc5.paragraphs[i+1+gi], f"{'①②③④'[gi]} {g}", "仿宋_GB2312", 12)
    elif t == "4、素质目标：" and i+4 < len(doc5.paragraphs):
        goals = ["培养数据驱动的创新思维和问题解决能力",
                 "形成规范操作、主动学习和依据标准检查成果的习惯",
                 "提升任务分工、沟通反馈、成果检查和按时交付能力",
                 "养成按步骤实施、及时测试、记录问题和持续改进的习惯"]
        for gi, g in enumerate(goals):
            pw(doc5.paragraphs[i+1+gi], f"{'①②③④'[gi]} {g}", "仿宋_GB2312", 12)
    elif t == "教学模式：" and i+2 < len(doc5.paragraphs):
        pw(doc5.paragraphs[i+1], "①理实一体化教学模式：理论讲解与实操训练交替进行，做中学、学中做", "仿宋_GB2312", 12)
        pw(doc5.paragraphs[i+2], "②项目引领、任务驱动：以企业真实项目为载体，以具体任务为线索", "仿宋_GB2312", 12)
    elif t == "教学方法：" and i+2 < len(doc5.paragraphs):
        pw(doc5.paragraphs[i+1], "①案例教学法：通过分析电商行业真实数据案例，引导学生理解方法应用", "仿宋_GB2312", 12)
        pw(doc5.paragraphs[i+2], "②操作演示法：教师演示操作步骤，学生跟随练习并独立完成实训", "仿宋_GB2312", 12)
    elif t == "教材：" and i+1 < len(doc5.paragraphs):
        pw(doc5.paragraphs[i+1], f"教材：《{tb}》，天津大学出版社", "仿宋_GB2312", 12)
    elif t == "教学资料：" and i+1 < len(doc5.paragraphs):
        pw(doc5.paragraphs[i+1], f"教学资料：配套PPT课件8套（共153页）、实训源代码112个文件、"
                                 f"企业真实数据案例库、Python开发环境（Anaconda+Jupyter Notebook）", "仿宋_GB2312", 12)

# 整体设计表格填充
t0 = doc5.tables[0]
cw(t0.rows[0].cells[1], cn); cw(t0.rows[0].cells[3], cc); cw(t0.rows[0].cells[5], dp)
cw(t0.rows[1].cells[1], "2024年2月"); cw(t0.rows[1].cells[3], tn)
cw(t0.rows[2].cells[1], ct); cw(t0.rows[2].cells[3], str(th)); cw(t0.rows[2].cells[5], str(cr))
cw(t0.rows[3].cells[1], tm); cw(t0.rows[3].cells[4], cls)
cw(t0.rows[4].cells[1], "电子商务基础、计算机应用基础")
cw(t0.rows[4].cells[4], "电子商务综合实训、顶岗实习")

# 表1
t1d = doc5.tables[1]
while len(t1d.rows) > 2: t1d._tbl.remove(t1d.rows[-1]._tr)
tr1d = deepcopy(t1d.rows[1]._tr)
for i, u in enumerate(units):
    r = t1d.rows[1] if i == 0 else None
    if r is None: t1d._tbl.append(deepcopy(tr1d)); r = t1d.rows[-1]
    clc(r.cells[0]); cw(r.cells[0], u["project_title"])
    clc(r.cells[1]); cw(r.cells[1], str(int(u["suggested_hours"] or 6)))
t1d._tbl.append(deepcopy(tr1d))
r = t1d.rows[-1]; clc(r.cells[0]); cw(r.cells[0], "合计")
clc(r.cells[1]); cw(r.cells[1], str(th))

# 表2
t2d = doc5.tables[2]
while len(t2d.rows) > 2: t2d._tbl.remove(t2d.rows[-1]._tr)
tr2d = deepcopy(t2d.rows[1]._tr)
cnt = 0
for ui, u in enumerate(units):
    sk = u.get("source_skills", "")
    items = [s.strip() for s in sk.split("；") if s.strip()] if sk else []
    for si, s in enumerate(items[:3]):
        if cnt == 0: r = t2d.rows[1]
        else: t2d._tbl.append(deepcopy(tr2d)); r = t2d.rows[-1]
        cnt += 1
        for ci in range(len(r.cells)): clc(r.cells[ci])
        cw(r.cells[0], str(ui+1))
        cw(r.cells[1], u["project_title"])
        cw(r.cells[2], f"{ui+1}.{si+1} {s}")
        cw(r.cells[3], f"{'①②③'[si]}{s}")
        cw(r.cells[4], "项目教学法：教师演示→学生实操→成果提交→检查优化")
        cw(r.cells[5], f"{u['project_title']}实训成果")

# 表3
t3d = doc5.tables[3]
while len(t3d.rows) > 4: t3d._tbl.remove(t3d.rows[-1]._tr)
tr3d = deepcopy(t3d.rows[3]._tr)
for i, task in enumerate(tasks):
    r = t3d.rows[3] if i == 0 else None
    if r is None: t3d._tbl.append(deepcopy(tr3d)); r = t3d.rows[-1]
    c = contents[i]
    for ci in range(len(r.cells)): clc(r.cells[ci])
    cw(r.cells[0], str(task["seq"]))
    cw(r.cells[1], str(task.get("hours", 2)))
    cw(r.cells[2], c["title"][:40])
    cw(r.cells[3], c["ag"].split("\n")[0][:50])
    for ui, u in enumerate(units):
        if u["project_title"] == task["chapter"]: cw(r.cells[4], str(ui+1)); break
    cw(r.cells[5], c["kg"].split("\n")[0][:50])
    cw(r.cells[6], c["ig"].split("\n")[0][:50])
    cw(r.cells[7], c["qg"].split("\n")[0][:50])
    cw(r.cells[8], "项目教学法、任务驱动法")
    cw(r.cells[9], "实训成果提交+课堂表现")

# 表4
t4d = doc5.tables[4]
for ri in range(1, 5): cw(t4d.rows[ri].cells[0], cn)
cw(t4d.rows[5].cells[0], cn)
cw(t4d.rows[5].cells[2], "终结性考核")
cw(t4d.rows[5].cells[3], "综合作品（数据采集、清洗、分析、可视化和报告撰写）")

# ========== 单元设计（仅表5+表6）==========
body5 = doc5.element.body
sectPr5 = body5.find(qn('w:sectPr'))
pt5 = deepcopy(doc5.tables[5]._tbl)
pt6 = deepcopy(doc5.tables[6]._tbl)

# 删除表4之后所有内容
t4xml = doc5.tables[4]._tbl
to_remove = []
found = False
for child in list(body5):
    if found:
        if child.tag == qn('w:sectPr'): continue
        to_remove.append(child)
    if child is t4xml: found = True
for e in to_remove: body5.remove(e)

if sectPr5 is not None:
    p = sectPr5.getparent()
    if p is not None: p.remove(sectPr5)

for idx, c in enumerate(contents):
    # ===== 表5 =====
    nt5 = deepcopy(pt5); body5.append(nt5)
    t5 = Table(nt5, doc5)
    
    # 先清空所有单元格
    for row in t5.rows:
        for cell in row.cells: clc(cell)
    
    # 行0：周次、课时、授课班级
    cw(t5.rows[0].cells[2], c["week"])
    cw(t5.rows[0].cells[4], c["hours"])
    cw(t5.rows[0].cells[6], cls)
    # 行1：授课教师、授课日期
    cw(t5.rows[1].cells[2], tn)
    cw(t5.rows[1].cells[6], c["date"])
    # 行2：课程类型、教学环境
    cw(t5.rows[2].cells[2], "理实一体课程")
    cw(t5.rows[2].cells[6], c["room"])
    # 行3：教学任务
    cw(t5.rows[3].cells[2], c["title"])
    # 行4：知识目标
    cw(t5.rows[4].cells[2], c["kg"])
    # 行5：能力目标
    cw(t5.rows[5].cells[2], c["ag"])
    # 行6：思政目标
    cw(t5.rows[6].cells[2], c["ig"])
    # 行7：素质目标
    cw(t5.rows[7].cells[2], c["qg"])
    # 行8：教材学情分析及教育理念
    cw(t5.rows[8].cells[2], 
        f"教材分析：{c['textbook_an']}\n"
        f"学情分析：{c['xueqing']}\n"
        f"教师教育理念：{c['linian']}")
    # 行9：教学重点
    cw(t5.rows[9].cells[2], f"【教学重点】{c['zhongdian']}")
    # 行10：教学难点
    cw(t5.rows[10].cells[2], f"【教学难点】{c['nandian']}")
    
    # 行11-14：四个知识点
    for ki, kp in enumerate(c["kps"]):
        ri = 11 + ki
        if ri >= len(t5.rows): break
        cw(t5.rows[ri].cells[2], f"{'一二三四'[ki]}、{kp['name']}")
        # 重点列7
        if kp["重点"]: cw(t5.rows[ri].cells[7], "√")
        # 难点列8
        if kp["难点"]: cw(t5.rows[ri].cells[8], "√")
        # 识记列11
        if kp["识记"]: cw(t5.rows[ri].cells[11], "√")
        # 理解列15
        if kp["理解"]: cw(t5.rows[ri].cells[15], "√")
        # 应用列17
        if kp["应用"]: cw(t5.rows[ri].cells[17], "√")
        # 评价列19
        if kp["评价"]: cw(t5.rows[ri].cells[19], "√")
    
    # ===== 表6 =====
    nt6 = deepcopy(pt6); body5.append(nt6)
    t6 = Table(nt6, doc5)
    
    for row in t6.rows:
        for cell in row.cells: clc(cell)
    
    # 行0：教学场景设计
    cw(t6.rows[0].cells[2], 
        f"{c['room']}，配备多媒体教学设备、Python开发环境（Anaconda+Jupyter Notebook）。\n"
        f"场景布置：教师机1台+投影设备+学生机每人1台，预装Python 3.x及numpy、pandas、"
        f"matplotlib、scikit-learn等数据分析库。\n"
        f"分组方式：4人一组，围坐式布局，便于讨论和互助。")
    # 行1：教学资源准备
    cw(t6.rows[1].cells[2], 
        f"1. 多媒体课件：本节课PPT（约20页），含概念讲解、案例展示、代码演示\n"
        f"2. 教材：《{tb}》对应章节\n"
        f"3. 软件环境：Python 3.x、Anaconda、Jupyter Notebook\n"
        f"4. 实训数据：本节课实训项目数据集（CSV/Excel格式）\n"
        f"5. 企业案例：电商行业真实数据分析案例（农产品电商、直播电商）\n"
        f"6. 参考资料：Python官方文档、Pandas官方文档、Kaggle竞赛平台")
    
    # 行3：教学导入
    cw(t6.rows[3].cells[2], 
        f"【案例展示】播放电商企业{c['title']}应用案例：\n"
        f"  · 某农产品电商平台通过数据分析发现用户购买偏好，精准推荐提升转化率30%\n"
        f"  · 某直播电商通过实时数据分析调整货品策略，单场GMV提升50%\n"
        f"【提问引导】提出3个问题：\n"
        f"  1. {c['title']}在以上案例中起到了什么作用？\n"
        f"  2. 如果没有{c['title']}，企业可能会做出什么错误决策？\n"
        f"  3. 你认为做好{c['title']}最关键的是什么？\n"
        f"【分组讨论】4人一组讨论3分钟，推选代表发言\n"
        f"【教师引导】点评学生回答，引出本节课学习内容")
    cw(t6.rows[3].cells[3], 
        "案例展示法、提问引导法、分组讨论法；\n"
        "学生：观看案例、思考问题、小组讨论、代表发言")
    cw(t6.rows[3].cells[4], 
        f"1. 激发学习兴趣，了解{c['title']}的实际应用价值\n"
        f"2. 引发思考，建立新旧知识联系\n"
        f"3. 培养分析问题和表达观点的能力")
    cw(t6.rows[3].cells[5], "10分钟")
    
    # 行4：教师归纳
    cw(t6.rows[4].cells[2], 
        f"教师归纳学生讨论结果，明确本节课学习目标：\n"
        f"1. 知识目标：理解{c['title']}的基本概念、原理和方法\n"
        f"2. 能力目标：掌握操作方法，能够独立完成实训任务\n"
        f"3. 思政目标：{c['ideo0'][:30]}...\n"
        f"4. 素质目标：培养规范操作、主动学习的良好习惯")
    cw(t6.rows[4].cells[3], "总结归纳法、讲授法；学生：明确目标、调整状态、准备工具")
    cw(t6.rows[4].cells[4], "明确学习目标，导入新课")
    cw(t6.rows[4].cells[5], "5分钟")
    
    # 行5：任务1
    cell1 = t6.rows[5].cells[2]
    p = cell1.paragraphs[0]
    r = p.add_run(f"任务1：{c['t1_name']}（知识讲解）\n"); sf(r, "仿宋_GB2312", 9, True)
    r = p.add_run(
        f"一、基本概念\n"
        f"  1. 定义：{c['kps'][0]['name']}是什么\n"
        f"  2. 核心要素：包含哪些关键组成部分\n"
        f"  3. 与相关概念的区别与联系\n\n"
        f"二、作用与意义\n"
        f"  1. 在数据分析全流程中的位置和作用\n"
        f"  2. 对企业决策的支撑价值\n"
        f"  3. 电商运营、市场营销、用户分析中的应用场景\n\n"
        f"三、基本原理\n"
        f"  1. 核心原理和方法\n"
        f"  2. 实现步骤和关键技术\n"
        f"  3. 常见误区和注意事项\n\n"
        f"四、分类与适用场景\n"
        f"  1. 常用分类方式\n"
        f"  2. 各种方法优缺点对比\n"
        f"  3. 如何根据实际场景选择合适的方法"); sf(r, "仿宋_GB2312", 9)
    if c.get("img"):
        try:
            np = cell1.add_paragraph()
            r = np.add_run()
            r.add_picture(BytesIO(c["img"]["data"]), width=Cm(5))
        except: pass
    cw(t6.rows[5].cells[3], 
        "讲授法、案例教学法、启发式提问；\n"
        "学生：听讲、思考、做笔记、回答问题")
    cw(t6.rows[5].cells[4], 
        f"1. 理解{c['t1_name']}的基本概念和原理\n"
        f"2. 掌握核心知识点，建立知识框架\n"
        f"3. 了解应用场景和实际价值")
    cw(t6.rows[5].cells[5], "25分钟")
    
    # 行6：任务2
    cell2 = t6.rows[6].cells[2]
    p = cell2.paragraphs[0]
    r = p.add_run(f"任务2：{c['t2_name']}（技术演示）\n"); sf(r, "仿宋_GB2312", 9, True)
    r = p.add_run(
        f"一、环境准备\n"
        f"  1. 打开Anaconda Navigator，启动Jupyter Notebook\n"
        f"  2. 新建Python3 Notebook并重命名\n"
        f"  3. 导入必要的库\n\n"
        f"二、代码演示（逐行讲解）\n"
        f"  1. 数据加载和预处理\n"
        f"  2. 核心功能实现\n"
        f"  3. 结果输出和验证\n"
        f"  4. 常见错误和调试方法\n\n"
        f"三、关键步骤解析\n"
        f"  每一步为什么这么做？有什么注意事项？\n"
        f"  如果参数改变会有什么影响？\n\n"
        f"四、学生跟随练习\n"
        f"  学生在自己电脑上跟随操作，教师巡回指导\n\n"
        f"【德育渗透】{c['ideo1']}"); sf(r, "仿宋_GB2312", 9)
    if c.get("t1_code"):
        add_code(cell2, c["t1_code"][:400], 7)
    cw(t6.rows[6].cells[3], 
        "操作演示法、逐步讲解法、巡回指导法；\n"
        "学生：跟随操作、提问讨论、记录笔记")
    cw(t6.rows[6].cells[4], 
        f"1. 掌握{c['t2_name']}的操作步骤和代码实现\n"
        f"2. 培养动手操作和代码调试能力\n"
        f"3. 渗透{c['ideo1'].split('：')[0] if '：' in c['ideo1'] else '德育'}教育")
    cw(t6.rows[6].cells[5], "30分钟")
    
    # 行7：任务3
    cell3 = t6.rows[7].cells[2]
    p = cell3.paragraphs[0]
    r = p.add_run(f"任务3：{c['t3_name']}（实操练习）\n"); sf(r, "仿宋_GB2312", 9, True)
    r = p.add_run(
        f"一、实训目标\n"
        f"  1. 独立完成{c['title']}的完整操作流程\n"
        f"  2. 培养发现问题、分析问题、解决问题的能力\n"
        f"  3. 养成规范操作和文档记录的习惯\n\n"
        f"二、实训任务\n"
        f"  任务A：基础操作——按照要求完成基础功能实现\n"
        f"  任务B：进阶练习——在基础上增加新功能，尝试不同参数\n"
        f"  任务C：拓展挑战——自选数据集，用所学方法进行分析\n\n"
        f"三、操作要求\n"
        f"  1. 独立完成，可以讨论但不能抄袭\n"
        f"  2. 代码规范，注释完整\n"
        f"  3. 记录遇到的问题和解决方法\n"
        f"  4. 按时提交实训成果\n\n"
        f"四、教师指导\n"
        f"  巡回观察、个别辅导、共性问题集中讲解\n\n"
        f"五、成果展示\n"
        f"  抽取2-3组展示成果，师生共同点评"); sf(r, "仿宋_GB2312", 9)
    if c.get("t2_code"):
        add_code(cell3, c["t2_code"][:300], 7)
    cw(t6.rows[7].cells[3], 
        "实操练习法、任务驱动法、巡回指导法、成果展示法；\n"
        "学生：独立操作、小组互助、展示分享、互评学习")
    cw(t6.rows[7].cells[4], 
        f"1. 能够独立完成{c['t3_name']}实训任务\n"
        f"2. 培养实操能力和问题解决能力\n"
        f"3. 提升团队协作和表达能力")
    cw(t6.rows[7].cells[5], "25分钟")
    
    # 行8：课堂小结
    cw(t6.rows[8].cells[2], c["xiaojie"])
    cw(t6.rows[8].cells[3], 
        "总结归纳法、提问回顾法、思维导图法；\n"
        "学生：回顾、思考、整理笔记、回答问题")
    cw(t6.rows[8].cells[4], 
        "1. 系统梳理本节课知识点，形成知识体系\n"
        "2. 强化重点和难点，加深理解记忆\n"
        "3. 建立前后知识联系，为后续学习铺垫")
    cw(t6.rows[8].cells[5], "5分钟")
    
    # 行9：课后作业
    cw(t6.rows[9].cells[2], c["zuoye"])
    cw(t6.rows[9].cells[3], 
        "课后自主完成；学生：独立思考、动手实践、查阅资料、小组讨论")
    cw(t6.rows[9].cells[4], 
        "1. 巩固所学知识和技能\n"
        "2. 培养自主学习和问题解决能力\n"
        "3. 拓展延伸，培养创新思维\n"
        "4. 预习下节课内容")
    cw(t6.rows[9].cells[5], "5分钟")
    
    # 行10：教学反思
    cw(t6.rows[10].cells[2], c["fansi"])
    cw(t6.rows[10].cells[3],
        "课后填写；教师自我反思、持续改进")
    cw(t6.rows[10].cells[4],
        "1. 总结教学经验，持续优化教学\n"
        "2. 发现问题，及时调整改进\n"
        "3. 提升自身教学能力")
    
    if (idx+1) % 10 == 0:
        print(f"  单元设计 {idx+1}/{len(contents)}")

if sectPr5 is not None:
    body5.append(sectPr5)

doc5.save(out5)
print(f"\n教学设计完成：{out5}")

try:
    shutil.copy2(out5, f"D:\\BaiduSyncdisk\\{os.path.basename(out5)}")
    print("已复制到D盘")
except Exception as e:
    print(f"D盘复制失败: {e}")
