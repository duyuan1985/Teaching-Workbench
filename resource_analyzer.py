import ast
import json
import re
from collections import Counter
from pathlib import Path

from docx import Document
from PIL import Image
from pptx import Presentation

import store
from resource_indexer import legacy_ppt_text


TEXT_EXTENSIONS = {".py", ".html", ".htm", ".css", ".js", ".json", ".xml", ".csv", ".txt", ".md"}


# ============================================================
# 文档类型识别（基于文件名关键词）
# ============================================================

# 权重映射：不同类型的文档在检索时的权重
DOCUMENT_WEIGHTS = {
    "教学大纲": 1.5,    # 课程标准、教学要求，最重要
    "课程标准": 1.5,
    "教案": 1.3,        # 教学过程设计，参考价值高
    "教学设计": 1.3,
    "实训指导书": 1.3,   # 实训操作指导
    "实验指导": 1.3,
    "实训大纲": 1.3,
    "实验手册": 1.2,
    "实训方案": 1.2,
    "教材": 1.0,         # 通用教材，基础权重
    "课本": 1.0,
    "习题答案": 0.6,     # 答案参考，权重较低
    "参考答案": 0.6,
    "考核方案": 0.8,     # 考核方式参考
    "案例": 0.9,         # 案例素材
    "需求文档": 0.7,     # 需求说明
    "交付": 0.6,         # 交付文档
    "设计说明": 0.8,     # 设计说明
}

# 文档类型识别关键词（按优先级排序，先匹配的优先）
DOCUMENT_TYPE_PATTERNS = [
    ("教学大纲", ["教学大纲", "课程标准", "课程教学大纲"]),
    ("教案", ["教案", "教学设计"]),
    ("实训指导书", ["实训指导书", "实验指导", "实训指导", "实验手册", "实训手册"]),
    ("实训大纲", ["实训大纲", "实验大纲"]),
    ("实训方案", ["实训方案", "实训设计"]),
    ("习题答案", ["习题答案", "参考答案", "答案"]),
    ("考核方案", ["考核方案", "考核标准", "考试方案"]),
    ("案例", ["案例", "项目案例"]),
    ("需求文档", ["需求文档", "需求规格"]),
    ("教材", ["教材", "课本", "教程"]),
    ("设计说明", ["设计说明", "设计说明书"]),
    ("交付文档", ["交付", "结项"]),
]


def classify_document_type(file_path):
    """根据文件名判断文档类型，返回 (类型名, 权重)"""
    name = Path(file_path).stem.lower()
    for doc_type, keywords in DOCUMENT_TYPE_PATTERNS:
        for kw in keywords:
            if kw.lower() in name:
                return doc_type, DOCUMENT_WEIGHTS.get(doc_type, 1.0)
    return "其他文档", 1.0


# ============================================================
# PPT 噪音过滤
# ============================================================

def _detect_ppt_noise_texts(presentation):
    """检测PPT中重复出现的页眉页脚文字"""
    all_slide_texts = []
    for slide in presentation.slides:
        texts = set()
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.add(shape.text.strip())
        all_slide_texts.append(texts)

    total = len(all_slide_texts)
    if total < 3:
        return set()  # 页数太少，不做过滤

    # 统计每页都出现的文字
    text_counts = Counter()
    for texts in all_slide_texts:
        for t in texts:
            text_counts[t] += 1

    threshold = max(3, total * 0.5)  # 出现于50%以上页面
    noise = set()
    for text, count in text_counts.items():
        if count >= threshold:
            # 太短的数字（页码）也过滤
            if re.fullmatch(r'\d{1,3}', text):
                noise.add(text)
                continue
            # 包含公司网址、版权声明等关键词
            noise_keywords = ['http://', 'https://', 'www.', '版权', '出品', '集团', '科技']
            if any(kw in text for kw in noise_keywords):
                noise.add(text)
                continue
            # 纯数字加单位（如"第X页"）
            if re.fullmatch(r'第\s*\d+\s*页', text):
                noise.add(text)
                continue
            # 长文本不太可能是页眉，跳过
            if len(text) > 30:
                continue
            noise.add(text)

    return noise


PPT_NOISE_KEYWORDS = ("LOGO", "logo", "Logo")

def _is_noise_text(text, noise_set):
    """判断文本是否为噪音"""
    text = text.strip()
    if not text:
        return True
    if text in noise_set:
        return True
    # 纯数字（可能是页码）
    if re.fullmatch(r'\d{1,3}', text):
        return True
    # "第X页 / 共Y页" 格式
    if re.match(r'第\s*\d+\s*页', text) and len(text) < 20:
        return True
    # PPT模板占位符（LOGO等）
    if text in PPT_NOISE_KEYWORDS:
        return True
    return False


def _clean(text):
    return re.sub(r"\s+", " ", str(text or "")).strip()


def _read_text(path):
    for encoding in ("utf-8", "utf-8-sig", "gb18030"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, OSError):
            continue
    return ""


def _pptx(path):
    presentation = Presentation(path)
    noise_texts = _detect_ppt_noise_texts(presentation)
    slides = []
    facts = []
    for index, slide in enumerate(presentation.slides, 1):
        blocks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and _clean(shape.text):
                text = _clean(shape.text)
                if not _is_noise_text(text, noise_texts):
                    blocks.append(text)
        text = "；".join(blocks)
        slides.append({"slide": index, "text": text})
        if text:
            facts.append(("ppt_slide", f"第{index}页", text, f"slide:{index}", 1.0))
    return {"kind": "pptx", "slides": slides, "noise_removed": len(noise_texts)}, facts


def _legacy_ppt(path):
    text = legacy_ppt_text(path)
    slides = []
    facts = []
    for block in text.splitlines():
        match = re.match(r"第(\d+)页：(.*)", block)
        if not match:
            continue
        index, value = int(match.group(1)), _clean(match.group(2))
        slides.append({"slide": index, "text": value})
        if value:
            facts.append(("ppt_slide", f"第{index}页", value, f"slide:{index}", 1.0))
    return {"kind": "ppt", "slides": slides}, facts


def _is_toc_paragraph(text):
    """判断段落是否属于目录页"""
    # 目录标题
    if re.fullmatch(r'\s*目\s*录\s*', text):
        return True
    # 典型目录行：标题 + 省略号/点 + 页码
    if re.search(r'[·\.]+\s*\d+\s*$', text) and len(text) > 10:
        return True
    return False


def _is_heading_paragraph(paragraph):
    """判断段落是否为标题（基于样式或加粗特征）"""
    style = paragraph.style.name if paragraph.style else ''
    # 样式判断
    if re.match(r'Heading\s*\d+', style, re.I):
        return True
    if '标题' in style:
        return True
    # 加粗且短文本判断（伪标题）
    text = paragraph.text.strip()
    if not text or len(text) > 40:
        return False
    if paragraph.runs:
        bold_runs = [r for r in paragraph.runs if r.text.strip() and r.font.bold]
        text_runs = [r for r in paragraph.runs if r.text.strip()]
        if text_runs and len(bold_runs) == len(text_runs):
            # 全部加粗且短文本
            # 进一步判断：中文标题特征（数字开头+顿号/点+内容）
            if re.match(r'^[一二三四五六七八九十\d]+[、.．]\s*\S', text):
                return True
            if re.match(r'^第[一二三四五六七八九十\d]+[章节部分讲]\s*\S', text):
                return True
    return False


def _heading_level(paragraph):
    """获取标题级别，返回 1-6，非标题返回 0"""
    style = paragraph.style.name if paragraph.style else ''
    m = re.match(r'Heading\s*(\d+)', style, re.I)
    if m:
        return int(m.group(1))
    if style == 'Title':
        return 0  # 文档总标题
    text = paragraph.text.strip()
    if not text:
        return 0
    # 伪标题级别判断
    if re.match(r'^第[一二三四五六七八九十\d]+[章节]\s*', text):
        return 1
    if re.match(r'^[一二三四五六七八九十]+[、.．]\s*\S', text):
        return 2
    if re.match(r'^\d+[、.．]\s*\S', text) and len(text) < 30:
        return 3
    if re.match(r'^[（(][一二三四五六七八九十\d]+[)）]\s*\S', text):
        return 4
    return 0


def _docx(path):
    document = Document(path)

    # 判断文档类型和权重
    doc_type, doc_weight = classify_document_type(path)

    paragraphs = []
    sections = []  # 章节结构
    facts = []
    current_section = None
    in_toc = False
    toc_ended = False

    for para_index, paragraph in enumerate(document.paragraphs, 1):
        text = _clean(paragraph.text)
        if not text:
            continue

        # 目录检测与跳过
        if not toc_ended:
            if _is_toc_paragraph(text):
                in_toc = True
                continue
            if in_toc:
                # 目录通常在文档开头，遇到第一个一级标题或30段后结束
                if _heading_level(paragraph) == 1 or para_index > 50:
                    toc_ended = True
                    in_toc = False
                else:
                    continue

        # 标题识别
        is_heading = _is_heading_paragraph(paragraph)
        level = _heading_level(paragraph) if is_heading else 0

        if is_heading and level > 0:
            # 更新当前章节
            if level == 1:
                current_section = text
                sections.append({"level": level, "title": text, "para_index": para_index})
            elif current_section:
                sections.append({"level": level, "title": text, "parent": current_section, "para_index": para_index})

            # 标题也作为事实存入（用于检索章节名称）
            section_path = current_section if level > 1 and current_section else ""
            facts.append((
                "document_heading",
                f"{'　' * (level - 1)}{text}",
                text,
                f"paragraph:{para_index}",
                doc_weight * 1.2,  # 标题权重略高
            ))

        paragraphs.append(text)
        # 正文段落
        if not is_heading:
            context = f"[{current_section}] {text}" if current_section else text
            facts.append((
                "document_paragraph",
                f"段落{para_index}",
                context,
                f"paragraph:{para_index}",
                doc_weight,
            ))

    # 表格处理
    tables = []
    for table_index, table in enumerate(document.tables):
        rows = [[_clean(cell.text) for cell in row.cells] for row in table.rows]
        tables.append(rows)
        for row_index, row in enumerate(rows):
            value = " | ".join(item for item in row if item)
            if value:
                facts.append((
                    "document_table_row",
                    f"表{table_index + 1}行{row_index + 1}",
                    value,
                    f"table:{table_index}/row:{row_index}",
                    doc_weight * 0.9,
                ))

    return {
        "kind": "docx",
        "doc_type": doc_type,
        "doc_weight": doc_weight,
        "paragraphs": paragraphs,
        "sections": sections,
        "tables": tables,
        "section_count": len(sections),
    }, facts


def _pdf(path):
    try:
        from pypdf import PdfReader
    except ImportError:
        return {"kind": "pdf", "pages": [], "error": "缺少pypdf"}, []
    pages = []
    facts = []
    for index, page in enumerate(PdfReader(path).pages, 1):
        text = _clean(page.extract_text())
        pages.append({"page": index, "text": text})
        if text:
            facts.append(("pdf_page", f"第{index}页", text, f"page:{index}", 0.95))
    return {"kind": "pdf", "pages": pages}, facts


def _python(path, text):
    imports, classes, functions = [], [], []
    try:
        tree = ast.parse(text)
        for node in ast.walk(tree):
            if isinstance(node, ast.Import):
                imports.extend(alias.name for alias in node.names)
            elif isinstance(node, ast.ImportFrom):
                imports.append(node.module or "")
            elif isinstance(node, ast.ClassDef):
                classes.append(node.name)
            elif isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
                functions.append(node.name)
    except SyntaxError:
        pass
    structured = {"kind": "source", "language": "python", "imports": sorted(set(imports)), "classes": classes, "functions": functions, "text": text[:20000]}
    facts = []
    for key, values in (("imports", imports), ("classes", classes), ("functions", functions)):
        if values:
            facts.append(("source_structure", key, "、".join(dict.fromkeys(values)), key, 1.0))
    snippet = "\n".join(line.rstrip() for line in text.splitlines() if line.strip())[:500]
    if snippet:
        facts.append(("source_excerpt", path.name, snippet, "source:excerpt", 0.95))
    return structured, facts


def _source(path):
    text = _read_text(path)
    if path.suffix.lower() == ".py":
        return _python(path, text)
    tags = sorted(set(re.findall(r"<([a-zA-Z][\w-]*)\b", text))) if path.suffix.lower() in {".html", ".htm"} else []
    selectors = re.findall(r"(?:^|\})\s*([^@][^{]+)\s*\{", text)[:80] if path.suffix.lower() == ".css" else []
    structured = {"kind": "source", "language": path.suffix.lower().lstrip("."), "html_tags": tags, "css_selectors": [_clean(x) for x in selectors], "text": text[:20000]}
    facts = []
    if tags:
        facts.append(("source_structure", "html_tags", "、".join(tags), "html", 1.0))
    if selectors:
        facts.append(("source_structure", "css_selectors", "、".join(_clean(x) for x in selectors), "css", 0.95))
    snippet = "\n".join(line.rstrip() for line in text.splitlines() if line.strip())[:500]
    if snippet:
        facts.append(("source_excerpt", path.name, snippet, "source:excerpt", 0.95))
    return structured, facts


def _image(path):
    try:
        with Image.open(path) as image:
            width, height = image.size
            return {"kind": "image", "width": width, "height": height, "mode": image.mode}, [
                ("image_metadata", "dimensions", f"{width}x{height}", "image", 1.0)
            ]
    except Exception as error:
        return {"kind": "image", "readable": False, "error": str(error)}, []


def analyze_resource(path):
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return _pptx(path)
    if suffix == ".ppt":
        return _legacy_ppt(path)
    if suffix == ".docx":
        return _docx(path)
    if suffix == ".pdf":
        return _pdf(path)
    if suffix in TEXT_EXTENSIONS:
        return _source(path)
    if suffix in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}:
        return _image(path)
    return {"kind": "binary", "suffix": suffix}, []


def project_hint(path, units):
    normalized = re.sub(r"\W+", "", str(path).lower())
    for unit in units:
        title = re.sub(r"\W+", "", unit["project_title"].lower())
        if title and title in normalized:
            return unit["project_title"]
    match = re.search(r"source[-_ ]*core[-_ ]*0?(\d+)", str(path), re.I)
    if not match:
        match = re.search(r"core[-_ ]*0*(\d+)", str(path), re.I)
    if match:
        index = int(match.group(1))
        if 1 <= index <= len(units):
            return units[index - 1]["project_title"]
    project_match = re.search(r"项目\s*0*(\d+)", str(path), re.I)
    if project_match:
        index = int(project_match.group(1))
        if 1 <= index <= len(units):
            return units[index - 1]["project_title"]
    return ""


def analyze_offering_resources(offering_id):
    units = store.rows("SELECT seq,project_title FROM curriculum_units WHERE offering_id=? ORDER BY seq", (offering_id,))
    items = store.rows("SELECT * FROM resource_items WHERE offering_id=? ORDER BY id", (offering_id,))
    analyzed = failed = fact_count = 0
    with store.connect() as db:
        db.execute("DELETE FROM resource_facts WHERE offering_id=?", (offering_id,))
        for item in items:
            path = Path(item["file_path"])
            hint = project_hint(path, units)
            try:
                structured, facts = analyze_resource(path)
                status = "已解析"
                analyzed += 1
            except Exception as error:
                structured, facts = {"kind": "error", "error": str(error)}, []
                status = "解析失败"
                failed += 1
            db.execute(
                "UPDATE resource_items SET structured_json=?,extraction_status=?,project_hint=? WHERE id=?",
                (json.dumps(structured, ensure_ascii=False), status, hint, item["id"]),
            )
            db.executemany(
                "INSERT INTO resource_facts (offering_id,resource_item_id,project_hint,fact_type,fact_key,fact_value,locator,confidence) VALUES (?,?,?,?,?,?,?,?)",
                [(offering_id, item["id"], hint, *fact) for fact in facts],
            )
            fact_count += len(facts)
        db.commit()
    return {"resources": len(items), "analyzed": analyzed, "failed": failed, "facts": fact_count}
