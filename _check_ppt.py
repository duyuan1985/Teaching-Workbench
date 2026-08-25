"""检查PPT文件和源码文件结构"""
import os
from pptx import Presentation

ppt_base = r"原始资料\教材\商务数据分析\大数据分析方法项目实战\03 课程PPT"
code_base = r"原始资料\教材\商务数据分析\大数据分析方法项目实战\04 实训源代码"

print("=== PPT文件 ===")
if os.path.exists(ppt_base):
    for f in sorted(os.listdir(ppt_base)):
        if f.endswith(".pptx"):
            fp = os.path.join(ppt_base, f)
            try:
                prs = Presentation(fp)
                slide_count = len(prs.slides)
                # 提取前3页的文字
                texts = []
                for si, slide in enumerate(prs.slides):
                    if si >= 3:
                        break
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t and len(t) < 200:
                                    texts.append(t)
                print(f"\n{f} ({slide_count}页):")
                for t in texts[:10]:
                    print(f"  {t[:80]}")
            except Exception as e:
                print(f"  {f}: 读取失败 {e}")
else:
    print(f"路径不存在: {ppt_base}")

print(f"\n=== 源码文件 ===")
if os.path.exists(code_base):
    for root, dirs, files in os.walk(code_base):
        for f in sorted(files):
            if f.endswith(".py"):
                rel = os.path.relpath(os.path.join(root, f), code_base)
                print(f"  {rel}")
else:
    print(f"路径不存在: {code_base}")

# 检查tasks表中的chapter分布
import store
tasks = store.rows("SELECT * FROM tasks WHERE offering_id=20 ORDER BY seq")
print(f"\n=== 任务章节分布 ===")
chapters = {}
for t in tasks:
    ch = t["chapter"]
    if ch not in chapters:
        chapters[ch] = []
    chapters[ch].append(t["title"])
for ch, titles in chapters.items():
    print(f"\n{ch} ({len(titles)}个任务):")
    for t in titles:
        print(f"  - {t}")
