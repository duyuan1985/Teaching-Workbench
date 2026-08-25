"""
教学设计生成器 - 增强版
- 封面填充（段落+文本框）
- 表0-4填充（课程信息/内容设计/能力训练/进度表/考核方案）
- 表5+6克隆30对（每任务一对）
- AI内容接入（authored_sections）
"""
import os, shutil, random, json
from copy import deepcopy
from io import BytesIO
from docx import Document
from docx.shared import Pt, Cm
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from generate import (
    safe_set_text, clear_indent, set_indent_chars, fix_spacing_auto,
    set_alignment, set_page_break_before, set_space_before, set_line_spacing,
    write_cell, write_cell_lines, _set_font_name, _set_font_size,
    merge_v, merge_h, CN_NUMS,
    gen_summary, gen_homework, gen_reflection, get_ideo, get_opening_semester,
    find_unit_for_task, find_sess_for_task
)

random.seed(42)


def _load_ai_content(offering_id):
    """从authored_sections加载AI生成的unit_design内容"""
    import store
    rows = store.rows(
        "SELECT * FROM authored_sections WHERE offering_id=? AND section_key='unit_design' ORDER BY CAST(repeat_key AS INTEGER)",
        [offering_id]
    )
    ai_contents = {}
    for r in rows:
        try:
            data = json.loads(r["content_json"]) if r["content_json"] else {}
            rk = int(r.get("repeat_key", 0) or 0)
            ai_contents[rk] = data
        except (json.JSONDecodeError, ValueError):
            continue
    return ai_contents


def _fill_cover(doc, offering):
    """填充封面段落和文本框"""
    cn = offering["course_name"]
    tn = offering.get("teacher_name") or "杜媛"
    cls = offering.get("teaching_class", "")
    term = offering.get("term", "")
    mj = offering.get("major", "")
    textbook = offering.get("textbook_version", "")

    # 解析学期: "2023-2024-2" -> "2023——2024学年  第二学期"
    parts = term.split("-")
    if len(parts) == 3:
        y1, y2, sem = parts
        sem_text = f"{y1}——{y2}学年  {'第一' if sem == '1' else '第二'}学期"
    else:
        sem_text = term

    # 填充文本框 (w:txbxContent)
    body = doc.element.body
    txbx_list = body.findall('.//' + qn('w:txbxContent'))
    for txbx in txbx_list:
        for p in txbx.findall(qn('w:p')):
            full_text = ""
            for r in p.findall(qn('w:r')):
                for t in r.findall(qn('w:t')):
                    full_text += (t.text or '')

            new_text = None
            if '学年' in full_text and '学期' in full_text and '课程名称' not in full_text:
                new_text = sem_text
            elif '课程名称' in full_text:
                new_text = f"课程名称：{cn}"
            elif '班' in full_text and '级' in full_text:
                new_text = f"班    级：{cls}"
            elif '教' in full_text and '材' in full_text:
                new_text = f"教    材：{textbook}" if textbook else None
            elif '授课教师' in full_text or '教师' in full_text:
                new_text = f"授课教师：{tn}"

            if new_text is not None:
                runs = p.findall(qn('w:r'))
                if runs:
                    for t in runs[0].findall(qn('w:t')):
                        runs[0].remove(t)
                    for r in runs[1:]:
                        p.remove(r)
                    t_elem = OxmlElement('w:t')
                    t_elem.text = new_text
                    runs[0].append(t_elem)

    # 填充封面段落
    cover_map = {
        "山西林业职业技术学院": None,  # 校名保持不变
        "《": f"《{cn}》",  # 课程名
        "课程整体教学设计": None,  # 保持不变
        "授课教师": f"授课教师：{tn}",
        "授课班级": f"授课班级：{cls}",
        "授课学期": f"授课学期：{sem_text}",
    }
    for p in doc.paragraphs:
        text = p.text.strip()
        if text.startswith("《") and text.endswith("》"):
            for r in p.runs:
                r.text = ""
            if p.runs:
                p.runs[0].text = f"《{cn}》"
        elif "授课教师" in text and "：" in text:
            for r in p.runs:
                r.text = ""
            if p.runs:
                p.runs[0].text = f"授课教师：{tn}"
        elif "授课班级" in text and "：" in text:
            for r in p.runs:
                r.text = ""
            if p.runs:
                p.runs[0].text = f"授课班级：{cls}"
        elif "授课学期" in text and "：" in text:
            for r in p.runs:
                r.text = ""
            if p.runs:
                p.runs[0].text = f"授课学期：{sem_text}"
        elif text == "《   》课程整体教学设计" or text == "《   》课程整体教学设计":
            for r in p.runs:
                r.text = ""
            if p.runs:
                p.runs[0].text = f"《{cn}》课程整体教学设计"


def _fill_design_body(doc, offering, tasks, units, sessions):
    """填充教学设计正文的九个章节，避免模板标题后留空。"""
    cn = offering["course_name"]
    major = offering.get("major", "") or "本专业"
    ct = offering.get("course_type", "") or "专业核心课"
    th = int(offering.get("total_hours", 60) or 60)
    first_task = tasks[0]["title"] if tasks else "课程导入任务"
    chapter_names = []
    for task in tasks:
        if task["chapter"] not in chapter_names:
            chapter_names.append(task["chapter"])
    chapters = "、".join(chapter_names)
    resources = "、".join(["课程PPT", "教材", "实训数据集", "示例代码", "Anaconda/Jupyter", "学习通在线资源"])
    contents = {
        "二、课程定位": f"《{cn}》是{major}开设的{ct}，总学时{th}学时。课程面向电商数据处理、分析与可视化等岗位能力，围绕{chapters}组织教学，前接电子商务基础、计算机应用基础和Python程序设计，后续支撑平台运营、综合实训和顶岗实习。",
        "三、教学设计的理念": f"课程坚持以学生为中心、以真实任务为引领、以职业能力为本位，采用做中学、学中做的理实一体化理念。教学中融入数据安全、隐私保护、诚信分析和工匠精神，利用案例、演示、实操和成果评价促进学生知识、能力和素质协同发展。",
        "（四）课程教学模式和教学方法设计": "采用项目化、任务驱动和线上线下混合式教学。课前通过学习通推送微课和任务单，课中实施案例导入、教师演示、学生分组实操、巡回指导和成果展示，课后完成作业、反思和拓展任务。主要方法包括案例教学法、讲授法、操作演示法、任务驱动法和小组协作法。",
        "（五）课程教学进度表设计": f"课程按{len(chapter_names)}个教学单元组织，共{len(tasks)}个任务、{th}学时。教学进度依据任务表、排课日期和校历安排，理论讲解与实践训练交替进行，补课任务按数据库排课记录纳入相应周次。",
        "五、考核方案设计": "课程采用过程性考核与终结性考核相结合的方式：签到10%、课堂表现10%、作业20%，过程性考核合计40%；综合作品或课程项目成果占60%，重点评价数据处理规范、分析方法应用、结果表达和职业素养。",
        "六、课程教学实施条件": "教学场所配备多媒体设备、稳定网络和实训电脑，安装Excel、Anaconda、Jupyter Notebook及课程所需Python库。教师准备课程PPT、案例数据、操作演示和任务单，学生按小组使用学习通提交成果并开展互评。",
        "七、教学资源": f"主要资源包括：{resources}。资源内容覆盖{chapters}的知识讲解、操作示范、源代码、练习数据和评价量表，并结合行业案例和数据安全规范持续更新。",
        "八、需要说明的其他问题": "课程实施中根据学生基础差异进行分层指导，对代码调试和数据清洗等难点提供示例与补充练习；涉及真实业务数据时进行脱敏处理，严格遵守数据安全和个人信息保护要求。",
        "九、第一节课设计梗概": f"第一节课以{first_task}为载体，先展示电商数据分析案例，提出真实工作问题，再讲解课程目标、学习方式和评价标准。学生完成环境检查、案例讨论和首个小任务，形成对课程内容、学习要求和成果形式的整体认识。",
    }
    paragraphs = doc.paragraphs
    for i, p in enumerate(paragraphs):
        heading = p.text.strip()
        matched = next((key for key in contents if key in heading), None)
        if not matched:
            continue
        if i + 1 < len(paragraphs):
            nxt = paragraphs[i + 1]
            if not nxt.text.strip() or len(nxt.text.strip()) < 20:
                safe_set_text(nxt, contents[matched], font="仿宋_GB2312", size=10.5, bold=False)
                clear_indent(nxt)
                set_line_spacing(nxt, 1.5)
    # 汇总课程任务和章节目标，填写课程目标设计中的圈号条目。
    chapter_summary = "、".join(chapter_names)
    objective_sections = {
        "1、认知目标：": [
            f"理解{cn}的基本概念、工作流程和典型应用场景",
            f"掌握{chapter_summary}等课程核心知识与操作原理",
            "掌握数据采集、清洗、分析、可视化和结果表达的规范方法",
            "了解大数据、人工智能等新技术在业务数据分析中的应用",
        ],
        "2、能力目标：": [
            "能够依据任务要求选择合适的方法和工具完成数据处理与分析",
            "能够使用课程配套软件、示例代码和数据集完成实训任务",
            "能够检查运行结果、定位常见问题并持续优化分析成果",
            "能够整理过程材料、规范表达结论并完成项目成果汇报",
        ],
        "3、思政目标：": [
            "树立科技报国理想，认识数字技术服务产业发展和乡村振兴的价值",
            "坚持实事求是和诚信分析，不篡改、不隐瞒数据与分析结果",
            "遵守个人信息保护、数据安全和网络安全要求，按必要范围处理业务数据",
            "践行敬业、责任与法治精神，恪守数据分析岗位职业道德",
        ],
        "4、素质目标：": [
            "形成规范操作、及时测试、记录问题和持续改进的学习习惯",
            "提升独立分析、主动学习、表达技术思路和解决问题的能力",
            "提升任务分工、沟通反馈、成果检查和按时交付的团队协作能力",
            "培养创新意识、质量意识和精益求精的工匠精神",
        ],
    }
    markers = "①②③④⑤⑥"
    for heading, goals in objective_sections.items():
        heading_index = next((i for i, paragraph in enumerate(paragraphs) if paragraph.text.strip() == heading), None)
        if heading_index is None:
            continue
        goal_index = 0
        for paragraph in paragraphs[heading_index + 1:]:
            current = paragraph.text.strip()
            if current in objective_sections and current != heading:
                break
            if current.startswith("（二）"):
                break
            if current and current[0] in markers and goal_index < len(goals):
                safe_set_text(paragraph, f"{markers[goal_index]} {goals[goal_index]}", font="仿宋_GB2312", size=10.5, bold=False)
                clear_indent(paragraph)
                set_line_spacing(paragraph, 1.5)
                goal_index += 1
    for p in paragraphs:
        if "《   》课程单元教学设计" in p.text:
            safe_set_text(p, f"《{cn}》课程单元教学设计", font="方正小标宋简体", size=16, bold=True)
    # 删除封面文本框中由模板带来的下划线
    for u in doc.element.body.findall(".//" + qn("w:u")):
        u.set(qn("w:val"), "none")


def _fill_table0(t, offering, units, sessions=None):
    """填充表0：课程信息表。"""
    from datetime import date, timedelta
    cn = offering["course_name"]
    tn = offering.get("teacher_name") or "杜媛"
    term = offering.get("term", "")
    th = int(offering.get("total_hours", 60) or 60)
    cr = offering.get("credits", 3) or 3
    cls = offering.get("teaching_class", "")
    course_type = offering.get("course_type") or offering.get("course_nature") or "专业核心课"
    first_date = None
    for item in sessions or []:
        raw = item.get("lesson_date")
        if raw:
            try:
                current = date.fromisoformat(str(raw)[:10])
                first_date = current if first_date is None or current < first_date else first_date
            except ValueError:
                pass
    developed = (first_date - timedelta(days=7)).strftime("%Y年%m月").replace("年0", "年") if first_date else term.rsplit("-", 1)[0] + "年"
    parts = term.split("-")
    term_text = f"{parts[0]}-{parts[1]}学年{'第一' if parts[2] == '1' else '第二'}学期" if len(parts) == 3 else term
    prerequisites = offering.get("prerequisite_courses") or "电子商务基础、计算机应用基础、Python程序设计"
    followups = offering.get("followup_courses") or "新媒体平台运营与推广、电子商务综合实训"
    values = {(0,1):cn,(0,3):offering.get("course_code",""),(0,5):offering.get("department","经济贸易系"),(1,1):developed,(1,3):tn,(2,1):course_type,(2,3):str(th),(2,5):str(cr).rstrip("0").rstrip(".") if isinstance(cr,float) else str(cr),(3,1):get_opening_semester(offering) or term_text,(3,4):cls,(4,1):prerequisites,(4,4):followups}
    for (row, col), value in values.items():
        if row < len(t.rows) and col < len(t.rows[row].cells):
            write_cell(t.rows[row].cells[col], value, "仿宋_GB2312", 10.5, False)


def _fill_table1(t, offering, tasks, units):
    """填充表1：课程内容设计（章节名+学时+合计）。"""
    chapter_map = {}
    chapter_order = []
    for task in tasks:
        chapter = task["chapter"]
        if chapter not in chapter_map:
            chapter_map[chapter] = []
            chapter_order.append(chapter)
        chapter_map[chapter].append(task)
    while len(t.rows) > 1:
        t.rows[-1]._element.getparent().remove(t.rows[-1]._element)
    total_hours = 0
    for i, chapter in enumerate(chapter_order):
        row = t.add_row()
        hours = sum(int(task.get("hours", 2) or 2) for task in chapter_map[chapter])
        total_hours += hours
        write_cell(row.cells[0], f"第{CN_NUMS[min(i, 9)]}章 {chapter}", "仿宋_GB2312", 9, False)
        write_cell(row.cells[1], str(hours), "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
    total = t.add_row()
    write_cell(total.cells[0], "合计", "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
    write_cell(total.cells[1], str(total_hours), "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)


def _fill_table2(t, offering, tasks, units):
    """填充表2：职业能力训练项目设计（每任务一行）"""
    while len(t.rows) > 1:
        t.rows[-1]._element.getparent().remove(t.rows[-1]._element)

    for idx, task in enumerate(tasks):
        ch = task["chapter"]
        title = task["title"].split("：", 1)[1] if "：" in task["title"] else task["title"]
        unit = find_unit_for_task(task, units)
        skills = unit.get("source_skills", "") if unit else ""
        sks = [s.strip() for s in skills.split("；") if s.strip()]

        row = t.add_row()
        cells = row.cells
        n = len(cells)

        def safe_cw(ci, text, font="仿宋_GB2312", size=9, bold=False, align=None):
            if ci < n:
                write_cell(cells[ci], text, font, size, bold, align)

        safe_cw(0, str(idx + 1), "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        safe_cw(1, f"{ch}：{title}", "仿宋_GB2312", 9, False)
        safe_cw(2, task.get("ability_goal", "")[:80] or f"掌握{title}的操作技能", "仿宋_GB2312", 9, False)
        safe_cw(3, "\n".join(sks[:4]) if sks else title, "仿宋_GB2312", 9, False)
        safe_cw(4, "任务驱动法、案例教学法、小组协作法", "仿宋_GB2312", 9, False)
        safe_cw(5, f"完成{title}实训", "仿宋_GB2312", 9, False)


def _fill_table3(t, offering, tasks, sessions):
    """填充表3：课程教学进度表"""
    while len(t.rows) > 3:
        t.rows[-1]._element.getparent().remove(t.rows[-1]._element)

    for idx, task in enumerate(tasks):
        row = t.add_row()
        cells = row.cells
        n = len(cells)

        def safe_cw(ci, text, font="仿宋_GB2312", size=9, bold=False, align=None):
            if ci < n:
                write_cell(cells[ci], text, font, size, bold, align)

        chapter = task["chapter"]
        title = task["title"].split("：", 1)[1] if "：" in task["title"] else task["title"]
        week = task.get("week_no", "")
        date = task.get("lesson_date", "")

        safe_cw(0, str(idx + 1), "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        safe_cw(1, "2", "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        safe_cw(2, f"第{week}周", "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        safe_cw(3, f"{chapter}：{title}", "仿宋_GB2312", 9, False)

        goals = task.get("knowledge_goal", "")
        safe_cw(4, goals[:100] if goals else f"掌握{title}的知识和技能", "仿宋_GB2312", 9, False)

        if n > 5:
            safe_cw(5, "理实一体化", "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        if n > 6:
            safe_cw(6, "801教室" if not date else f"{date}", "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        if n > 7:
            safe_cw(7, "任务驱动", "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        if n > 8:
            hw = task.get("ability_goal", "")[:60] if task.get("ability_goal") else f"完成{title}实训"
            safe_cw(8, hw, "仿宋_GB2312", 9, False)
        if n > 9:
            safe_cw(9, "杜媛", "仿宋_GB2312", 9, False, WD_ALIGN_PARAGRAPH.CENTER)


def _fill_table4(t, offering):
    """填充表4：考核方案表（过程性40%+终结性60%，细分行）"""
    ncols = len(t.columns) if t.rows else 0

    if ncols >= 5:
        # 旧模板：项目/学习内容/考核方式/考核标准/权重
        while len(t.rows) > 1:
            t.rows[-1]._element.getparent().remove(t.rows[-1]._element)

        # 过程性考核 - 签到
        r1 = t.add_row()
        write_cell(r1.cells[0], "过程性考核\n（40%）", "仿宋", 9, True, WD_ALIGN_PARAGRAPH.CENTER)
        write_cell(r1.cells[1], "签到", "仿宋", 9, False)
        write_cell(r1.cells[2], "考勤记录", "仿宋", 9, False)
        write_cell(r1.cells[3], "10%", "仿宋", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        write_cell(r1.cells[4], "40%", "仿宋", 9, True, WD_ALIGN_PARAGRAPH.CENTER)

        # 过程性考核 - 课堂表现
        r2 = t.add_row()
        write_cell(r2.cells[1], "课堂表现", "仿宋", 9, False)
        write_cell(r2.cells[2], "课堂观察", "仿宋", 9, False)
        write_cell(r2.cells[3], "10%", "仿宋", 9, False, WD_ALIGN_PARAGRAPH.CENTER)

        # 过程性考核 - 作业
        r3 = t.add_row()
        write_cell(r3.cells[1], "作业", "仿宋", 9, False)
        write_cell(r3.cells[2], "作业批改", "仿宋", 9, False)
        write_cell(r3.cells[3], "20%", "仿宋", 9, False, WD_ALIGN_PARAGRAPH.CENTER)

        # 终结性考核
        r4 = t.add_row()
        write_cell(r4.cells[0], "终结性考核\n（60%）", "仿宋", 9, True, WD_ALIGN_PARAGRAPH.CENTER)
        write_cell(r4.cells[1], "综合作品", "仿宋", 9, False)
        write_cell(r4.cells[2], "作品评分量表", "仿宋", 9, False)
        write_cell(r4.cells[3], "60%", "仿宋", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
        write_cell(r4.cells[4], "60%", "仿宋", 9, True, WD_ALIGN_PARAGRAPH.CENTER)

        # 合并：过程性考核 项目列(0)、权重列(4)
        merge_v(t, 1, 3, 0)
        merge_v(t, 1, 3, 4)
    else:
        # 新模板：评价类型/评价内容/评价主体/评价方式
        while len(t.rows) > 1:
            t.rows[-1]._element.getparent().remove(t.rows[-1]._element)

        r1 = t.add_row()
        if len(r1.cells) >= 4:
            write_cell(r1.cells[0], "过程性考核\n（40%）", "仿宋", 9, True, WD_ALIGN_PARAGRAPH.CENTER)
            write_cell(r1.cells[1], "签到(10%)+课堂表现(10%)+作业(20%)", "仿宋", 9, False)
            write_cell(r1.cells[2], "教师评价\n学生互评", "仿宋", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
            write_cell(r1.cells[3], "签到记录+课堂观察+作业批改", "仿宋", 9, False)

        r2 = t.add_row()
        if len(r2.cells) >= 4:
            write_cell(r2.cells[0], "终结性考核\n（60%）", "仿宋", 9, True, WD_ALIGN_PARAGRAPH.CENTER)
            write_cell(r2.cells[1], "综合作品：完整项目报告", "仿宋", 9, False)
            write_cell(r2.cells[2], "教师评价", "仿宋", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
            write_cell(r2.cells[3], "作品评分量表", "仿宋", 9, False)

        if len(t.rows) >= 3:
            merge_v(t, 1, 2, 2)
            merge_v(t, 1, 2, 3)


def _clone_task_tables(doc, n_tasks):
    """克隆表5+6对，为每个任务创建一对表格"""
    if len(doc.tables) < 7:
        return
    if n_tasks <= 1:
        return

    # 获取模板中的表5和表6的XML元素
    t5_elem = doc.tables[5]._element
    t6_elem = doc.tables[6]._element

    # 获取表5前面的"教学设计·基本信息"段落
    t5_prev = t5_elem.getprevious()
    t6_prev = t6_elem.getprevious()  # "教学设计·教学组织"段落

    # 最后一个元素（表6）
    last_elem = t6_elem

    for i in range(1, n_tasks):
        # 克隆段落和表格
        new_p5 = deepcopy(t5_prev) if t5_prev is not None else None
        new_t5 = deepcopy(t5_elem)
        new_p6 = deepcopy(t6_prev) if t6_prev is not None else None
        new_t6 = deepcopy(t6_elem)

        # 在最后一个元素后面依次插入
        if new_p5 is not None:
            last_elem.addnext(new_p5)
            last_elem = new_p5
        last_elem.addnext(new_t5)
        last_elem = new_t5
        if new_p6 is not None:
            last_elem.addnext(new_p6)
            last_elem = new_p6
        last_elem.addnext(new_t6)
        last_elem = new_t6

def generate_design(offering, tasks, sessions, units, output_dir=None):
    """生成教学设计"""
    from generate import get_template_path, get_output_path

    template_path = get_template_path("design", offering)
    if not template_path or not os.path.exists(template_path):
        raise FileNotFoundError(f"教学设计模板不存在: {template_path}")

    fp = get_output_path("design", offering, output_dir)
    shutil.copy2(template_path, fp)
    doc = Document(fp)

    cn = offering["course_name"]
    tn = offering.get("teacher_name") or "杜媛"
    th = int(offering.get("total_hours", 60))
    cr = int(offering.get("credits", 3))
    mj = offering.get("major", "")
    cls = offering.get("teaching_class", "")

    # 读取PPT和源码资料
    ppt_data, code_files = _load_resources(offering)

    # 按学习情境分组
    chapter_map = {}
    chapter_order = []
    for task in tasks:
        ch = task["chapter"]
        if ch not in chapter_map:
            chapter_map[ch] = []
            chapter_order.append(ch)
        chapter_map[ch].append(task)

    # 单元PPT映射
    unit_ppt = {}
    for i, ch in enumerate(chapter_order, 1):
        unit_ppt[i] = f"CORE-{i:02d}"

    # ideo映射
    ideo_map = {
        "初识": ["数据安全意识：遵守《个人信息保护法》《数据安全法》",
                 "诚信分析精神：坚持真实分析和规范验证",
                 "科技报国情怀：服务乡村振兴战略"],
        "Excel": ["严谨细致：公式和数据透视不能出错",
                  "规范操作：遵循数据处理规范",
                  "工匠精神：精益求精"],
        "Numpy": ["科学精神：严格遵循数学原理",
                  "严谨细致：数组维度不能错",
                  "规范操作：遵循编码规范"],
        "Pandas": ["数据诚信：如实记录处理过程",
                  "隐私保护：数据脱敏处理",
                  "工匠精神：严格把关数据质量"],
        "SciPy": ["科学求真：结论要有数据支撑",
                  "逻辑思维：从现象发现本质",
                  "创新探索：尝试不同方法"],
        "Sklearn": ["算法伦理：确保模型公平公正",
                   "数据偏见：培养批判性思维",
                   "科技向善：让科技服务社会"],
        "Seaborn": ["数据美学：以清晰直观方式呈现数据",
                   "严谨表达：图表标注完整准确",
                   "创新思维：探索可视化新方法"],
        "综合": ["团队协作：分工合作互相帮助",
                "责任担当：按时高质量交付",
                "持续学习：树立终身学习理念"],
    }

    def get_ideo_local(ch):
        for k, v in ideo_map.items():
            if k in ch: return v
        return ideo_map["综合"]

    # ============================================================
    # 填充封面（段落+文本框）
    # ============================================================
    _fill_cover(doc, offering)
    _fill_design_body(doc, offering, tasks, units, sessions)

    # ============================================================
    # 填充表0-4（课程信息/内容设计/能力训练/进度表/考核方案）
    # ============================================================
    if len(doc.tables) >= 5:
        _fill_table0(doc.tables[0], offering, units, sessions)
        _fill_table1(doc.tables[1], offering, tasks, units)
        _fill_table2(doc.tables[2], offering, tasks, units)
        _fill_table3(doc.tables[3], offering, tasks, sessions)
        _fill_table4(doc.tables[4], offering)

    # ============================================================
    # 加载AI内容
    # ============================================================
    ai_contents = _load_ai_content(offering["id"])

    # ============================================================
    # 克隆表5+6对（为每个任务创建一对表格）
    # ============================================================
    _clone_task_tables(doc, len(tasks))

    # ============================================================
    # 填充每对表5+6（基本信息+教学组织）
    # 克隆后: 表5,7,9...=基本信息表; 表6,8,10...=教学组织表
    # ============================================================
    for idx, task in enumerate(tasks):
        t5_idx = 5 + idx * 2
        t6_idx = 6 + idx * 2

        if t5_idx >= len(doc.tables) or t6_idx >= len(doc.tables):
            break

        t5 = doc.tables[t5_idx]
        t6 = doc.tables[t6_idx]

        unit = find_unit_for_task(task, units)
        sess = find_sess_for_task(task, sessions)
        chapter = task["chapter"]
        title = task["title"].split("：", 1)[1] if "：" in task["title"] else task["title"]
        week = task.get("week_no", "")
        hours = task.get("hours", 2)
        classroom = sess.get("classroom", "801教室") if sess else "801教室"
        lesson_date = sess.get("lesson_date", "") if sess else ""

        ideo = get_ideo_local(chapter)

        # 获取知识点
        skills = unit.get("source_skills", "") if unit else ""
        sks = [s.strip() for s in skills.split("；") if s.strip()]

        # PPT要点
        ui = None
        for i2, u in enumerate(units):
            if u["project_title"] == chapter:
                ui = i2 + 1
                break
        pkey = unit_ppt.get(ui, "CORE-01")
        pdat = ppt_data.get(pkey, {"slides": [], "images": []})

        hl = []
        for sl in pdat.get("slides", [])[:20]:
            for tt in sl["texts"][:4]:
                if 3 < len(tt) < 120:
                    hl.append(tt)
                if len(hl) >= 25:
                    break
            if len(hl) >= 25:
                break
        if not hl:
            hl = sks if sks else [title]

        # 知识点
        if sks and len(sks) >= 4:
            kp_names = sks[:4]
        elif sks:
            kp_names = sks + hl[:4 - len(sks)]
        else:
            kp_names = hl[:4]
        kp_names = [k[:50] for k in kp_names[:4]]

        # 生成差异化内容
        summary = gen_summary(task, idx, kp_names, title, chapter, ideo)
        homework = gen_homework(task, idx, kp_names, title, chapter, ideo)
        reflection = gen_reflection(task, idx, kp_names, title, chapter, ideo)

        # AI内容
        ai_data = ai_contents.get(idx + 1, {})

        # 填充表5（基本信息）- 整个表格
        _fill_task_basic_info(t5, offering, task, sess, unit, ideo, idx,
                              cn, mj, cls, tn, ai_data)

        # 填充表6（教学组织）- 整个表格
        _fill_task_teaching_org(t6, task, offering, unit, sess,
                                kp_names, ideo, summary, homework, reflection,
                                title, chapter, pdat, code_files, ui, ai_data)

    # 清除模板残留
    _clean_residues(doc)
    _set_rows_auto(doc)

    doc.save(fp)
    return fp


def _load_resources(offering):
    """加载PPT和源码资源"""
    ppt_data = {}
    code_files = {}

    # 尝试多个可能的路径
    cn = offering["course_name"]
    possible_ppt_bases = [
        os.path.join("原始资料", "教材", cn, "大数据分析方法项目实战", "03 课程PPT"),
        os.path.join("原始资料", "教材", cn),
    ]
    possible_code_bases = [
        os.path.join("原始资料", "教材", cn, "大数据分析方法项目实战", "04 实训源代码"),
        os.path.join("原始资料", "教材", cn),
    ]

    ppt_base = None
    for p in possible_ppt_bases:
        if os.path.exists(p):
            ppt_base = p
            break

    code_base = None
    for p in possible_code_bases:
        if os.path.exists(p):
            code_base = p
            break

    # 读取PPT
    if ppt_base:
        for f in sorted(os.listdir(ppt_base)):
            if f.endswith(".pptx"):
                key = f.replace(".pptx", "")
                try:
                    prs = Presentation(os.path.join(ppt_base, f))
                    slides = []
                    images = []
                    for si, slide in enumerate(prs.slides):
                        texts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    t = para.text.strip()
                                    if t and len(t) < 200:
                                        texts.append(t)
                            if shape.shape_type == 13:
                                try:
                                    images.append({"slide": si + 1, "data": shape.image.blob})
                                except:
                                    pass
                        if texts:
                            slides.append({"index": si + 1, "texts": texts})
                    ppt_data[key] = {"slides": slides, "images": images}
                except:
                    pass

    # 读取源码
    if code_base:
        for root, dirs, files in os.walk(code_base):
            for f in files:
                if f.endswith(".py"):
                    fp = os.path.join(root, f)
                    rel = os.path.relpath(fp, code_base)
                    try:
                        with open(fp, "r", encoding="utf-8", errors="ignore") as fh:
                            code_files[rel] = fh.read(2000)
                    except:
                        pass

    return ppt_data, code_files


def _fill_task_basic_info(t, offering, task, sess, unit, ideo, idx, cn, mj, cls, tn, ai_data):
    """填充表5，使用模板标签定位行，兼容新旧模板。"""
    chapter = task["chapter"]
    title = task["title"].split("：", 1)[1] if "：" in task["title"] else task["title"]
    week = task.get("week_no", "") or idx + 1
    date = task.get("lesson_date", "") or (sess.get("lesson_date", "") if sess else "")
    classroom = sess.get("classroom", "801教室") if sess else "801教室"
    hours = task.get("hours", 2) or 2
    ai_info = ai_data.get("基本信息", {})
    ai_goals = ai_data.get("教学目标", {})
    kg = "\n".join(ai_goals.get("知识目标", [])) or task.get("knowledge_goal", "") or f"理解{title}的基本概念"
    ag = "\n".join(ai_goals.get("能力目标", [])) or task.get("ability_goal", "") or f"能够运用{title}完成对应任务"
    ig = "\n".join(ai_goals.get("思政目标", [])) or task.get("ideological_goal", "") or ideo[0]
    qg = "\n".join(ai_goals.get("素质目标", [])) or task.get("quality_goal", "") or "形成规范操作和严谨分析的习惯"
    skills = unit.get("source_skills", "") if unit else ""
    sks = [item.strip() for item in skills.split("；") if item.strip()]
    analysis = ai_info.get("教材学情", "") or f"教材内容围绕{chapter}组织知识与实践，强调概念、方法和操作的衔接。学生已具备相关基础知识，但对{title}的综合应用仍需通过案例演示、分步训练和成果评价加以巩固。"

    def set_c(row, col, text, font="仿宋", size=9, bold=False, align=None):
        if 0 <= row < len(t.rows) and 0 <= col < len(t.rows[row].cells):
            write_cell(t.rows[row].cells[col], text, font, size, bold, align)

    def row_for(*labels, default=None):
        for ri, row in enumerate(t.rows):
            label = "|".join(dict.fromkeys(cell.text.replace("\n", "").strip() for cell in row.cells[:2]))
            if any(key in label for key in labels):
                return ri
        return default

    set_c(0, 2, str(week), "宋体", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
    set_c(0, 4, str(hours), "宋体", 9, False, WD_ALIGN_PARAGRAPH.CENTER)
    set_c(0, 6, cls, "宋体", 9)
    set_c(1, 2, tn)
    set_c(1, 6, date)
    set_c(2, 2, offering.get("course_type") or "理实一体化课程")
    set_c(2, 6, classroom)
    set_c(3, 2, f"{chapter}\n子情景：{title.replace('、', '；')}")

    for row, text in [
        (row_for("知识目标", default=4), kg),
        (row_for("能力目标", default=5), ag),
        (row_for("思政目标", default=6), ig),
        (row_for("素质目标", default=7), qg),
        (row_for("教材学情", default=8), analysis),
    ]:
        set_c(row, 2, text)

    requirement_row = row_for("知识＆技能", "知识&技能", default=9)
    item_start = min(requirement_row + 2, len(t.rows))
    for offset, row in enumerate(range(item_start, min(item_start + 4, len(t.rows)))):
        text = sks[offset] if offset < len(sks) else ""
        set_c(row, 2, text)
        if text:
            set_c(row, 7, "√", align=WD_ALIGN_PARAGRAPH.CENTER)


def _fill_task_teaching_org(t, task, offering, unit, sess,
                             kp_names, ideo, summary, homework, reflection,
                             title, chapter, pdat, code_files, ui, ai_data):
    """填充表6，保留模板表头并写入正确的内容锚点。"""
    ai_org = ai_data.get("教学组织", {})

    def set_c(row, col, text, font="仿宋", size=9, bold=False, align=None):
        if 0 <= row < len(t.rows) and 0 <= col < len(t.rows[row].cells):
            write_cell(t.rows[row].cells[col], text, font, size, bold, align)

    classroom = sess.get("classroom", "801教室") if sess else "801教室"
    set_c(0, 2, f"{classroom}，配备多媒体教学设备、课程所需软件环境及实训数据集")
    set_c(1, 2, f"1. 多媒体课件：{chapter}PPT\n2. 微课视频：{title}\n3. 实训数据集和任务单\n4. 示例代码或操作素材\n5. 在线学习资源")

    intro = ai_org.get("教学导入", "") or f"展示与{title}相关的业务案例，提出问题，引导学生明确本次任务。"
    set_c(3, 2, intro)
    set_c(3, 3, "案例展示法、提问引导法、分组讨论法")
    set_c(3, 4, "激发学习兴趣，建立任务认知")
    set_c(3, 5, "5分钟", align=WD_ALIGN_PARAGRAPH.CENTER)
    set_c(4, 2, "教师归纳讨论结果，说明学习目标、成果要求和评价标准；学生明确任务并做好准备。")
    set_c(4, 3, "总结归纳法、讲授法")
    set_c(4, 4, "明确学习目标和完成标准")
    set_c(4, 5, "5分钟", align=WD_ALIGN_PARAGRAPH.CENTER)

    times = ("25分钟", "30分钟", "20分钟")
    for offset in range(3):
        row = 5 + offset
        key = f"任务{offset + 1}"
        point = kp_names[offset] if offset < len(kp_names) else title
        content = ai_org.get(key, "") or f"任务{offset + 1}：{point}\n教师讲解核心知识并演示操作，学生完成分步练习、检查结果并交流改进。"
        set_c(row, 2, content)
        set_c(row, 3, "任务驱动法、操作演示法、巡回指导法")
        set_c(row, 4, f"掌握{point}并能完成相应操作")
        set_c(row, 5, times[offset], align=WD_ALIGN_PARAGRAPH.CENTER)

    set_c(8, 2, ai_org.get("课堂小结", "") or summary)
    set_c(8, 3, "归纳总结、互动答疑")
    set_c(8, 4, "梳理知识与技能，查漏补缺")
    set_c(8, 5, "5分钟", align=WD_ALIGN_PARAGRAPH.CENTER)
    set_c(9, 2, ai_org.get("课后作业", "") or homework)
    set_c(10, 2, ai_org.get("教学反思", "") or reflection)


def _set_rows_auto(doc):
    """取消模板固定行高，避免教学组织内容截断。"""
    for table in doc.tables:
        for row in table.rows:
            trPr = row._tr.get_or_add_trPr()
            height = trPr.find(qn("w:trHeight"))
            if height is not None:
                height.set(qn("w:hRule"), "auto")
            else:
                height = OxmlElement("w:trHeight")
                height.set(qn("w:hRule"), "auto")
                trPr.append(height)


def _clean_residues(doc):
    """清除模板残留"""
    import re as re_mod
    # 正则模式（需要转义的）
    re_patterns = [r'×+', r'XX+', r'XXX+']
    # 字面匹配模式
    literal_patterns = [
        "（注：", "(注：", "（建议", "【模板", "（模板",
        "体例之一", "巩固知识", "课堂总结",
        "XXXX教学法", "（对涉及的教学法",
    ]

    def is_residue(text):
        for pat in re_patterns:
            if re_mod.search(pat, text):
                return True
        for pat in literal_patterns:
            if pat in text:
                return True
        return False

    for p in doc.paragraphs:
        if is_residue(p.text):
            for r in p.runs:
                r.text = ""

    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                if is_residue(cell.text):
                    for p in cell.paragraphs:
                        for r in p.runs:
                            r.text = ""
